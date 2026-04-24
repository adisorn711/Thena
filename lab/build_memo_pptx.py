from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# === Color Palette ===
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─── Helpers ──────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.75, RGBColor(0x0A, 0x14, 0x20))
    rect(slide, 0, 0.75, 13.33, 0.05, C_ACCENT)
    txt(slide, title, 0.5, 0.1, 11, 0.6, size=24, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.55, 11, 0.3, size=11, color=C_ACCENT)

def footer(slide, n, note="Data Science Team | เมษายน 2025"):
    rect(slide, 0, 7.15, 13.33, 0.35, RGBColor(0x0A, 0x14, 0x20))
    txt(slide, note, 0.5, 7.18, 11, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

# ─── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
circ = s.shapes.add_shape(9, Inches(8.5), Inches(-0.5), Inches(6), Inches(6))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x52)
circ.line.fill.background()

rect(s, 0, 2.8, 0.1, 2.0, C_ACCENT)
txt(s, "Customer Intelligence Platform", 0.6, 1.6, 11, 1.0,
    size=44, bold=True, color=C_WHITE)
txt(s, "แผนงานสำหรับการนำเสนอ", 0.6, 2.8, 9, 0.6,
    size=24, color=C_ACCENT)
txt(s, "ทีม Data Science  |  เมษายน 2025", 0.6, 3.6, 9, 0.4,
    size=14, color=C_LIGHT)
footer(s, 1)

# ─── Slide 2: Big Picture — 5 Layers ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ภาพใหญ่ — เราจะไปถึงไหน",
       "Hyundai Card UNIVERSE Platform — 5 Layers of Intelligence")

txt(s, "ธุรกิจบัตรเครดิตที่แข็งแกร่งที่สุดไม่ได้แข่งด้วย Cashback แต่แข่งด้วยความเข้าใจลูกค้า",
    0.5, 1.05, 12.3, 0.45, size=14, color=C_LIGHT)

layers = [
    ("1", "Customer Tags",          "รู้ว่าลูกค้าเป็นใคร",           "Foundation ทั้งหมด",              C_ACCENT,  True),
    ("2", "Personalized Trigger",    "เสนอ offer ที่ใช่ ให้คนที่ใช่ ในเวลาที่เหมาะ", "ลด CAC, เพิ่ม Conversion", RGBColor(0x4D,0xD0,0xFF), False),
    ("3", "Personalized Products",  "Card ที่ benefit ตรงแต่ละคน",    "ลด Churn, เพิ่ม Spend per Card",  RGBColor(0x7B,0xC8,0xFF), False),
    ("4", "Merchant Intelligence",  "ขายความแม่นยำให้ Merchant",       "Revenue Stream ใหม่จาก Data",     RGBColor(0xA8,0xD8,0xEA), False),
    ("5", "Lifestyle Platform",     "ส่วนหนึ่งของชีวิตลูกค้า",         "Brand Moat คู่แข่ง copy ไม่ได้",  RGBColor(0xB0,0xC4,0xDE), False),
]

for i, (num, name, goal, result, col, active) in enumerate(layers):
    y = 1.6 + i * 1.0
    bg_col = RGBColor(0x00, 0x30, 0x45) if active else C_CARD
    rect(s, 0.5, y, 12.3, 0.85, bg_col)
    if active:
        rect(s, 0.5, y, 12.3, 0.85, RGBColor(0x00, 0x30, 0x45))
        rect(s, 0.5, y, 0.1, 0.85, C_ACCENT)
    rect(s, 0.5, y, 0.9, 0.85, col)
    txt(s, num, 0.5, y + 0.15, 0.9, 0.55, size=24, bold=True,
        color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name, 1.55, y + 0.08, 3.2, 0.4, size=14, bold=True, color=col)
    txt(s, goal, 4.9, y + 0.1, 3.9, 0.65, size=12, color=C_WHITE)
    txt(s, "→  " + result, 8.9, y + 0.1, 4.0, 0.65, size=12, color=C_LIGHT)
    if active:
        txt(s, "◀ เราอยู่ตรงนี้", 1.55, y + 0.5, 3.0, 0.3,
            size=10, color=C_ACCENT)

footer(s, 2)

# ─── Slide 3 (NEW): สิ่งที่เราได้ — Customer Tag Profile ──────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1: สิ่งที่เราได้ — Customer Tag Profile",
       "Tag Engine รัน daily → Profile 1 ชุดต่อลูกค้า 1 คน บอกว่า \"เขาเป็นใคร\"")

# Left: key info
txt(s, "Output คืออะไร?", 0.5, 1.0, 5.5, 0.4, size=14, bold=True, color=C_WHITE)
key_points = [
    ("323 Tags", "ครอบ 13 มิติพฤติกรรม"),
    ("Score 0–1", "ต่อ Tag ต่อลูกค้า — ยิ่งสูงยิ่ง active"),
    ("12 เดือน", "lookback window — ไม่ใช่แค่ปัจจุบัน"),
    ("Daily", "refresh — profile เปลี่ยนตาม behavior จริง"),
]
for j, (kw, desc) in enumerate(key_points):
    y = 1.5 + j * 0.9
    rect(s, 0.5, y, 5.5, 0.78, C_CARD)
    rect(s, 0.5, y, 0.08, 0.78, C_ACCENT)
    txt(s, kw, 0.75, y + 0.08, 1.6, 0.3, size=13, bold=True, color=C_ACCENT)
    txt(s, desc, 0.75, y + 0.4, 4.8, 0.3, size=11, color=C_LIGHT)

txt(s, "→  ลูกค้าคนเดียวกันมีได้หลาย Tags พร้อมกัน",
    0.5, 5.25, 5.5, 0.45, size=12, color=C_WHITE)

# Right: sample profile table
rect(s, 6.5, 1.0, 6.4, 0.48, RGBColor(0x00, 0x30, 0x45))
txt(s, "ตัวอย่าง: ลูกค้า A", 6.65, 1.07, 6.0, 0.35, size=13, bold=True, color=C_ACCENT)

profile = [
    ("flash_sale_hunter",       "ช้อปหนักช่วง sale event",    "0.82", C_ACCENT),
    ("japanese_food_lover",     "ชอบอาหารญี่ปุ่น",            "0.74", C_YELLOW),
    ("wealth_builder",          "กำลังสร้างความมั่งคั่ง",      "0.61", C_GREEN),
    ("weekend_warrior",         "ใช้จ่ายหนักช่วง weekend",    "0.55", C_ORANGE),
    ("convenience_store_addict","7-Eleven บ่อย",              "0.49", C_LIGHT),
]
for j, (tag, meaning, score, col) in enumerate(profile):
    y = 1.58 + j * 1.05
    row_bg = RGBColor(0x12, 0x24, 0x38) if j % 2 == 0 else C_CARD
    rect(s, 6.5, y, 6.4, 0.95, row_bg)
    txt(s, tag,     6.65, y + 0.08, 3.0, 0.35, size=11, bold=True, color=col)
    txt(s, meaning, 6.65, y + 0.52, 3.0, 0.35, size=10, color=C_LIGHT)
    txt(s, score,   9.75, y + 0.18, 1.0, 0.5,  size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    bar_fill = min(float(score), 1.0) * 2.6
    rect(s, 10.85, y + 0.35, 2.6,  0.18, RGBColor(0x0A, 0x18, 0x28))
    if bar_fill > 0:
        rect(s, 10.85, y + 0.35, bar_fill, 0.18, col)

footer(s, 3)

# ─── Slide 4 (NEW): จาก Tag สู่ Action — Automation Flow ──────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "จาก Tag สู่ Action — Automation Flow",
       "Tags ไม่ใช่แค่ข้อมูล — พวกมันคือ trigger ของทุก campaign และ offer")

steps = [
    ("①", "Tag\nEngine",      "รัน daily\nรู้ว่าลูกค้า\nเป็นใคร",     C_ACCENT),
    ("②", "Trigger\nEngine",  "ประเมินทุกวัน\nเลือกว่าใคร\nควรได้อะไร", RGBColor(0x4D,0xD0,0xFF)),
    ("③", "Campaign\nManager","เลือก offer\nปรับ message\nตาม tag",      C_YELLOW),
    ("④", "Delivery\nChannel","Push / SMS\nEmail\nIn-app",               C_GREEN),
    ("⑤", "Response\nTracker","วัดผล\nfeedback กลับ\nปรับ model",        C_ORANGE),
]

box_w = 2.2
gap   = 0.28
start_x = 0.45

for i, (num, name, desc, col) in enumerate(steps):
    x = start_x + i * (box_w + gap)

    # Box background
    rect(s, x, 1.3, box_w, 4.8, C_CARD)
    rect(s, x, 1.3, box_w, 0.55, col)

    # Number badge
    txt(s, num,  x + 0.05, 1.33, box_w - 0.1, 0.48,
        size=22, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # Name
    txt(s, name, x + 0.1, 2.0, box_w - 0.2, 0.85,
        size=15, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Description
    txt(s, desc, x + 0.1, 2.95, box_w - 0.2, 1.5,
        size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Arrow (not on last item)
    if i < len(steps) - 1:
        ax = x + box_w + 0.03
        txt(s, "→", ax, 3.2, gap + 0.1, 0.5,
            size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Bottom note
rect(s, 0.45, 6.3, 12.4, 0.72, RGBColor(0x0A, 0x20, 0x30))
rect(s, 0.45, 6.3, 0.08, 0.72, C_ACCENT)
txt(s, "ทีม DS ออกแบบ rule — ทุก offer มาจาก Tag ของลูกค้า ไม่ว่าจะ automate หรือ manual ก็ตาม",
    0.65, 6.42, 12.0, 0.45, size=12, color=C_WHITE)

footer(s, 4)

# ─── Slide 5 (NEW): ตัวอย่าง — 3 Use Cases ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ตัวอย่าง — Tag ได้แล้วเอาไปใช้ยังไง",
       "3 use cases จาก tag เดียว → trigger → campaign → outcome")

cases = [
    {
        "tag":     "flash_sale_hunter",
        "col":     C_ACCENT,
        "icon":    "🛒",
        "trigger": "5 วันก่อน 11.11 / 12.12",
        "action":  "Push: \"ผ่อน 0% 3 เดือน\"\npre-approval limit พิเศษ",
        "channel": "Push Notification",
        "outcome": "capture spend ก่อนลูกค้าเลือกบัตรอื่น",
    },
    {
        "tag":     "japanese_food_lover",
        "col":     C_YELLOW,
        "icon":    "🍣",
        "trigger": "ร้านญี่ปุ่นใหม่เปิดใน 5km",
        "action":  "In-app: \"แคชแบ็ก 15%\"\nที่ร้านใหม่นี้เฉพาะคุณ",
        "channel": "In-app Banner",
        "outcome": "drive trial + merchant partnership revenue",
    },
    {
        "tag":     "wealth_builder",
        "col":     C_GREEN,
        "icon":    "📈",
        "trigger": "ช่วง Q4 / ต้นปีใหม่",
        "action":  "Email: แนะนำ\ninvestment product / premium card",
        "channel": "Email",
        "outcome": "upsell + เพิ่ม LTV ระยะยาว",
    },
]

card_w = 3.9
for i, c in enumerate(cases):
    x = 0.45 + i * (card_w + 0.3)
    col = c["col"]

    # Card background
    rect(s, x, 1.05, card_w, 5.95, C_CARD)
    rect(s, x, 1.05, card_w, 0.65, col)

    # Tag name + icon
    txt(s, c["icon"] + "  " + c["tag"], x + 0.15, 1.1, card_w - 0.3, 0.5,
        size=12, bold=True, color=C_BG, align=PP_ALIGN.LEFT)

    rows = [
        ("Trigger",  c["trigger"], RGBColor(0x00,0x20,0x35)),
        ("Action",   c["action"],  RGBColor(0x0A,0x28,0x20)),
        ("Channel",  c["channel"], RGBColor(0x20,0x20,0x0A)),
        ("Outcome",  c["outcome"], RGBColor(0x20,0x0A,0x0A)),
    ]
    row_cols = [RGBColor(0x00,0x20,0x35), RGBColor(0x0A,0x25,0x18),
                RGBColor(0x20,0x18,0x00), RGBColor(0x20,0x08,0x08)]
    label_cols = [C_ACCENT, C_GREEN, C_YELLOW, C_ORANGE]

    for j, (label, content, _) in enumerate(rows):
        y = 1.85 + j * 1.25
        rect(s, x + 0.15, y, card_w - 0.3, 1.1, row_cols[j])
        txt(s, label, x + 0.25, y + 0.05, card_w - 0.5, 0.3,
            size=9, bold=True, color=label_cols[j])
        txt(s, content, x + 0.25, y + 0.38, card_w - 0.5, 0.65,
            size=11, color=C_WHITE)

footer(s, 5)

# ─── Slide 6 (NEW): 5 Building Blocks — สิ่งที่ต้องมีก่อน Automate ────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "สิ่งที่ต้องมีเพื่อให้ Automate ได้จริง",
       "5 building blocks ที่ต้องพร้อมพร้อมกัน — นี่คือ roadmap ที่เราจะสร้าง")

blocks = [
    ("①", "Tag Engine",                  "รัน daily pipeline สร้าง customer_tag_profile",
     "Path A: 254 tags | Path B: 69 tags | refresh ทุกคืน",          C_ACCENT),
    ("②", "Trigger Rule Engine",         "ประเมิน condition ทุกวัน → สร้าง eligible customer list",
     "rule: IF tag = X AND score ≥ Y AND condition Z → FIRE",         RGBColor(0x4D,0xD0,0xFF)),
    ("③", "CRM / Campaign Tool",         "รับ list + tag จาก Trigger Engine → personalize + schedule",
     "ต้องรองรับ tag-based targeting และ dynamic message template",    C_YELLOW),
    ("④", "Frequency Cap & Consent",     "ป้องกันส่ง message ถี่เกินไป + เคารพ opt-out ของลูกค้า",
     "max 2 campaigns/week/customer | honor opt-out flag",            C_GREEN),
    ("⑤", "Response Tracking Pipeline", "วัดว่า campaign ได้ผลไหม → feedback กลับปรับ model",
     "click / convert / ignore → reinforce หรือ adjust threshold",    C_ORANGE),
]

for i, (num, name, desc, detail, col) in enumerate(blocks):
    y = 1.1 + i * 1.15
    rect(s, 0.45, y, 12.4, 1.0, C_CARD)
    rect(s, 0.45, y, 0.75, 1.0, col)
    txt(s, num, 0.45, y + 0.2, 0.75, 0.55,
        size=22, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name,   1.35, y + 0.08, 3.5, 0.38, size=13, bold=True, color=col)
    txt(s, desc,   1.35, y + 0.52, 4.8, 0.38, size=11, color=C_WHITE)
    txt(s, detail, 6.3,  y + 0.08, 6.4, 0.8,  size=10, color=C_LIGHT)

footer(s, 6)

# ─── Slide 7: Layer 1 — Problem & Proposal ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1: Customer Tags — ที่เราอยู่ตอนนี้")

# Problem box
rect(s, 0.5, 1.05, 12.3, 1.5, RGBColor(0x2A, 0x10, 0x10))
rect(s, 0.5, 1.05, 0.1, 1.5, C_ORANGE)
txt(s, "ปัญหาที่เราเจออยู่", 0.75, 1.1, 11, 0.4, size=13, bold=True, color=C_ORANGE)
txt(s, "ปัจจุบันเราแบ่งลูกค้าเป็นกลุ่มใหญ่ๆ แล้วส่ง Campaign เดียวกันทั้งกลุ่ม\nทั้งที่ลูกค้าที่ใช้จ่ายเท่ากันอาจมีพฤติกรรมต่างกันอย่างสิ้นเชิง\nผลคือ Campaign ตรงสำหรับบางคน แต่พลาดสำหรับคนส่วนใหญ่",
    0.75, 1.55, 11.8, 0.85, size=13, color=C_LIGHT)

# Proposal box
rect(s, 0.5, 2.75, 12.3, 3.5, C_CARD)
rect(s, 0.5, 2.75, 0.1, 3.5, C_ACCENT)
txt(s, "สิ่งที่เราจะทำ", 0.75, 2.82, 11, 0.4, size=13, bold=True, color=C_ACCENT)
txt(s, "เราจะสร้างระบบ Customer Tagging ที่เปลี่ยนจาก \"ลูกค้าอยู่กลุ่มไหน\" มาเป็น \"ลูกค้าเป็นคนแบบไหน\"\nระบบนี้อ่านพฤติกรรมจาก Transaction ที่มีอยู่แล้ว ไม่ต้องลงทุนซื้อข้อมูลใหม่ และ Refresh ทุกวัน",
    0.75, 3.28, 11.8, 0.7, size=13, color=C_WHITE)

# Two stats
for i, (num, label) in enumerate([("254 Tags", "ทำได้ทันที"), ("69 Tags", "ต้องใช้ Enrichment")]):
    x = 1.5 + i * 5.5
    col = C_ACCENT if i == 0 else C_ORANGE
    rect(s, x, 4.2, 4.5, 1.7, RGBColor(0x0D, 0x2A, 0x3A) if i == 0 else RGBColor(0x2A, 0x1A, 0x0D))
    txt(s, num, x + 0.2, 4.35, 4.0, 0.8, size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, label, x + 0.2, 5.1, 4.0, 0.5, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

txt(s, "รวม 323 Tags ครอบ 13 มิติพฤติกรรม", 0.5, 6.1, 12.3, 0.4,
    size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

footer(s, 7)

# ─── Slide 8: 117 Tags — No Enrichment ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "254 Tags — ทำได้ทันที",
       "ใช้ข้อมูล Transaction ที่มีอยู่แล้ว ไม่รอ ไม่ต้องลงทุนเพิ่ม")

categories = [
    ("🍽️", "Food & Dining",     "17", ["frequent_diner", "cafe_hopper", "delivery_dependent", "late_night_diner", "budget_eater"]),
    ("✈️", "Travel",            "13", ["frequent_traveler", "luxury_traveler", "business_traveler", "budget_traveler"]),
    ("🛍️", "Shopping",          "19", ["marketplace_shopper", "luxury_shopper", "convenience_store_addict", "online_first_shopper"]),
    ("🎭", "Entertainment",     "13", ["movie_lover", "gym_member", "streaming_subscriber", "gamer"]),
    ("🚗", "Transport",         "16", ["car_owner", "grab_dependent", "ev_driver", "daily_commuter"]),
    ("💊", "Health",             "15", ["health_conscious", "pharmacy_regular", "hospital_visitor"]),
    ("💳", "Financial",         "25", ["high_spender", "installment_user", "reward_maximizer", "consistent_spender"]),
    ("👨‍👩‍👧", "Life Stage",      "15", ["parent", "pet_owner", "homeowner", "new_parent"]),
    ("⏰", "Time Patterns",     "15", ["night_owl_spender", "weekend_warrior", "morning_person"]),
    ("📱", "Digital",            "15", ["online_native", "crypto_curious", "fintech_user"]),
]

cols = 5
for i, (icon, cat, count, tags) in enumerate(categories):
    col = i % cols
    row = i // cols
    x = 0.4 + col * 2.58
    y = 1.1 + row * 2.95
    rect(s, x, y, 2.45, 2.7, C_CARD)
    rect(s, x, y, 2.45, 0.55, RGBColor(0x0A, 0x30, 0x45))
    txt(s, icon + " " + cat, x + 0.1, y + 0.05, 2.1, 0.45, size=11, bold=True, color=C_ACCENT)
    txt(s, count + " tags", x + 0.1, y + 0.6, 2.2, 0.35, size=11, color=C_YELLOW, bold=True)
    tag_text = "\n".join(["• " + t for t in tags[:3]])
    txt(s, tag_text, x + 0.1, y + 0.95, 2.3, 1.6, size=9, color=C_LIGHT)

footer(s, 8)

# ─── Slide 9: Merchant Enrichment Explained ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "69 Tags — ต้องใช้ Merchant Enrichment",
       "เพิ่มความเข้าใจร้านค้าในเชิงลึก เพื่อ Tags ที่ละเอียดขึ้นอีกชั้น")

txt(s, "Merchant Enrichment คืออะไร?", 0.5, 1.05, 12, 0.4,
    size=16, bold=True, color=C_WHITE)
txt(s, "ปัจจุบัน Transaction บอกได้ว่าลูกค้าไปร้านไหน และหมวดหมู่กว้างๆ เป็นอะไร\nEnrichment คือการเพิ่ม Label ให้ร้านแต่ละแห่งในเชิงลึก เพื่อให้เราแยกแยะได้มากขึ้น",
    0.5, 1.5, 12.3, 0.65, size=13, color=C_LIGHT)

rect(s, 0.5, 2.35, 5.9, 3.3, RGBColor(0x1A, 0x10, 0x10))
rect(s, 0.5, 2.35, 5.9, 0.45, RGBColor(0x3A, 0x10, 0x10))
txt(s, "❌  ก่อน Enrichment", 0.7, 2.4, 5.5, 0.35, size=13, bold=True, color=C_ORANGE)
before = [
    ("ร้าน Fuji", "→ รู้แค่ \"ร้านอาหาร\""),
    ("Marriott", "→ รู้แค่ \"โรงแรม\""),
    ("ผล:", "แยกได้แค่ว่า \"ชอบกินข้าวนอกบ้าน\""),
]
for j, (a, b) in enumerate(before):
    y = 2.95 + j * 0.8
    txt(s, a, 0.7, y, 1.6, 0.45, size=12, bold=True, color=C_WHITE)
    txt(s, b, 2.4, y, 3.8, 0.45, size=12, color=C_LIGHT)

rect(s, 6.9, 2.35, 5.9, 3.3, RGBColor(0x0A, 0x1E, 0x10))
rect(s, 6.9, 2.35, 5.9, 0.45, RGBColor(0x0A, 0x2E, 0x15))
txt(s, "✅  หลัง Enrichment", 7.1, 2.4, 5.5, 0.35, size=13, bold=True, color=C_GREEN)
after = [
    ("ร้าน Fuji", "→ \"อาหารญี่ปุ่น / ในห้าง / กลาง\""),
    ("Marriott", "→ \"City Hotel / Luxury / ธุรกิจ\""),
    ("ผล:", "แยกได้ว่า \"ชอบอาหารญี่ปุ่นโดยเฉพาะ\""),
]
for j, (a, b) in enumerate(after):
    y = 2.95 + j * 0.8
    txt(s, a, 7.1, y, 1.7, 0.45, size=12, bold=True, color=C_WHITE)
    txt(s, b, 8.9, y, 3.7, 0.45, size=12, color=C_LIGHT)

txt(s, "→", 6.3, 3.6, 0.6, 0.5, size=24, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

rect(s, 0.5, 5.85, 12.3, 0.8, C_CARD)
txt(s, "Tags ที่ unlock ได้หลัง Enrichment:",
    0.7, 5.9, 4.0, 0.35, size=12, bold=True, color=C_ACCENT)
enrichment_tags = "japanese_food_lover  |  korean_food_lover  |  resort_lover  |  mall_shopper  |  golf_player  |  yoga_practitioner  |  young_professional  |  ..."
txt(s, enrichment_tags, 0.7, 6.25, 12.0, 0.35, size=11, color=C_LIGHT)

footer(s, 9)

# ─── Slide 10: Tag Examples — Food, Travel, Shopping, Entertainment ─────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ตัวอย่าง Tags — ชีวิตประจำวัน",
       "แต่ละ Tag มาพร้อม Score 0–1 และ Refresh ทุกวันตามพฤติกรรมจริง")

groups = [
    ("🍽️ Food & Dining", C_ACCENT, [
        ("frequent_diner",       "กินข้าวนอกบ้าน ≥ 10 ครั้ง/เดือน"),
        ("cafe_hopper",          "ร้านกาแฟ ≥ 6 ครั้ง/เดือน หลาย merchant"),
        ("delivery_dependent",   "≥ 50% ของ dining ผ่าน GrabFood/LINE MAN"),
        ("late_night_diner",     "กินหลัง 22:00 ≥ 4 ครั้ง/เดือน"),
        ("fine_dining_enthusiast","avg ticket ร้านอาหาร > P75"),
    ]),
    ("✈️ Travel", RGBColor(0x7B,0xC8,0xFF), [
        ("frequent_traveler",    "airline + hotel ≥ 4 ครั้ง/ปี"),
        ("luxury_traveler",      "full-service airline + hotel avg > P75"),
        ("business_traveler",    "บินวันจันทร์-ศุกร์ + hotel weekday"),
        ("budget_traveler",      "LCC airline + hotel avg ต่ำ"),
        ("international_traveler","foreign currency txn ≥ 2 ครั้ง/ปี"),
    ]),
    ("🛍️ Shopping", C_YELLOW, [
        ("marketplace_shopper",  "Shopee/Lazada ≥ 4 ครั้ง/เดือน"),
        ("luxury_shopper",       "avg ticket > P90 หรือ luxury brand"),
        ("convenience_store_addict","7-Eleven/FamilyMart ≥ 15 ครั้ง/เดือน"),
        ("online_first_shopper", "≥ 50% ของ txn ทั้งหมดเป็น online"),
        ("grocery_regular",      "Tops/BigC/Lotus ≥ 3 ครั้ง/เดือน"),
    ]),
    ("🎭 Entertainment", C_GREEN, [
        ("movie_lover",          "โรงหนัง ≥ 2 ครั้ง/เดือน"),
        ("streaming_subscriber", "Netflix/Disney+ recurring ≥ 2 บริการ"),
        ("gym_member",           "ค่า fitness รายเดือนสม่ำเสมอ ≥ 6 เดือน"),
        ("gamer",                "Steam/PS Store ≥ 12 ครั้ง/ปี"),
        ("concert_goer",         "ticketing event ≥ 3 ครั้ง/ปี"),
    ]),
]

for i, (title, col, tags) in enumerate(groups):
    x = 0.4 + (i % 2) * 6.45
    y = 1.1 + (i // 2) * 2.95
    rect(s, x, y, 6.2, 2.75, C_CARD)
    rect(s, x, y, 6.2, 0.5, RGBColor(0x0A, 0x25, 0x35))
    txt(s, title, x + 0.15, y + 0.07, 5.8, 0.38, size=13, bold=True, color=col)
    for j, (tag, desc) in enumerate(tags):
        ty = y + 0.6 + j * 0.42
        txt(s, tag, x + 0.2, ty, 2.2, 0.38, size=10, bold=True, color=col)
        txt(s, desc, x + 2.5, ty, 3.5, 0.38, size=10, color=C_LIGHT)

footer(s, 10)

# ─── Slide 11: Tag Examples — Transport, Health, Financial, Life Stage ──────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ตัวอย่าง Tags — ไลฟ์สไตล์และการเงิน",
       "ครอบคลุม 13 มิติ — ทุก Tag มาจากสิ่งที่ลูกค้าทำอยู่แล้ว ไม่ใช่การเดา")

groups = [
    ("🚗 Transport", C_ORANGE, [
        ("car_owner",            "ครบ 3 อย่าง: fuel + parking + auto service"),
        ("grab_dependent",       "Grab ≥ 15 ครั้ง/เดือน"),
        ("ev_driver",            "มี txn ที่ EV Charging Station"),
        ("daily_commuter",       "fuel/transit weekday ≥ 20 วัน/เดือน"),
        ("fuel_brand_loyal",     "≥ 70% fuel txn ที่ปั๊มแบรนด์เดียวกัน"),
    ]),
    ("💊 Health & Wellness", RGBColor(0xFF,0x8C,0xD4), [
        ("health_conscious",     "gym + pharmacy + health food รวมกัน"),
        ("pharmacy_regular",     "ร้านขายยา ≥ 3 ครั้ง/เดือน"),
        ("hospital_visitor",     "โรงพยาบาล ≥ 2 ครั้ง/ไตรมาส"),
        ("supplement_buyer",     "Watsons/Boots ≥ 1 ครั้ง/เดือน"),
        ("gym_member",           "fitness center recurring ≥ 6 เดือน"),
    ]),
    ("💳 Financial Behavior", C_ACCENT, [
        ("high_spender",         "total spend > P75 ของ portfolio"),
        ("installment_user",     "มี installment ≥ 1 ครั้ง/ไตรมาส"),
        ("reward_maximizer",     "≥ 50% spend ที่ merchant points สูง"),
        ("consistent_spender",   "monthly spend variation ต่ำ สม่ำเสมอ"),
        ("paycheck_spender",     "spend spike ต้นเดือนทุกเดือน"),
    ]),
    ("👨‍👩‍👧 Life Stage", C_GREEN, [
        ("parent",               "baby/kids products + family restaurant บ่อย"),
        ("pet_owner",            "pet shop/vet ≥ 1 ครั้ง/เดือน"),
        ("homeowner",            "HomePro/SCG + utility ≥ 1 ครั้ง/ไตรมาส"),
        ("new_parent",           "spike ที่ baby products ใน 12 เดือนล่าสุด"),
        ("teen_spender",         "gaming + fastfood + convenience dominant"),
    ]),
]

for i, (title, col, tags) in enumerate(groups):
    x = 0.4 + (i % 2) * 6.45
    y = 1.1 + (i // 2) * 2.95
    rect(s, x, y, 6.2, 2.75, C_CARD)
    rect(s, x, y, 6.2, 0.5, RGBColor(0x0A, 0x25, 0x35))
    txt(s, title, x + 0.15, y + 0.07, 5.8, 0.38, size=13, bold=True, color=col)
    for j, (tag, desc) in enumerate(tags):
        ty = y + 0.6 + j * 0.42
        txt(s, tag, x + 0.2, ty, 2.2, 0.38, size=10, bold=True, color=col)
        txt(s, desc, x + 2.5, ty, 3.5, 0.38, size=10, color=C_LIGHT)

footer(s, 11)

# ─── Slide 12: Plan & Ask ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "แผนการทำงานและสิ่งที่ขออนุมัติ")

txt(s, "เราเสนอให้เริ่ม Phase 1 ด้วย 10 Tags ที่ชัดเจนและวัดผลได้ก่อน ทดสอบกับ Campaign จริง เทียบผลกับ approach เดิม และ scale เมื่อเห็นตัวเลขที่ดีขึ้น",
    0.5, 1.05, 12.3, 0.55, size=13, color=C_LIGHT)

phases = [
    ("Phase 1\nทันที",         C_ACCENT,                        ["10 Quick Win Tags", "Campaign A/B Test", "วัดผลเทียบ approach เดิม"]),
    ("Phase 2\nหลัง Pilot",    RGBColor(0x4D,0xD0,0xFF),        ["ขยาย 254 Tags ครบ", "Personalized Trigger", "Tags เป็น ML Features"]),
    ("Phase 3\nFull Intel",    RGBColor(0x7B,0xC8,0xFF),        ["69 Tags Enrichment", "Merchant Partnership", "Revenue Stream ใหม่"]),
    ("Phase 4\nLong-term",     RGBColor(0xB0,0xC4,0xDE),        ["Personalized Products", "Lifestyle Platform", "Brand Moat"]),
]

for i, (phase, col, bullets) in enumerate(phases):
    x = 0.4 + i * 3.15
    rect(s, x, 1.8, 3.0, 3.8, C_CARD)
    rect(s, x, 1.8, 3.0, 0.7, col)
    txt(s, phase, x + 0.15, 1.85, 2.7, 0.6, size=13, bold=True,
        color=C_BG, align=PP_ALIGN.LEFT)
    for j, b in enumerate(bullets):
        by = 2.65 + j * 0.85
        rect(s, x + 0.2, by, 2.6, 0.72, RGBColor(0x0D, 0x2A, 0x3A))
        txt(s, "→  " + b, x + 0.35, by + 0.1, 2.3, 0.55, size=11, color=C_WHITE)

txt(s, "สิ่งที่ขอการตัดสินใจ 3 เรื่อง", 0.5, 5.82, 12, 0.4,
    size=14, bold=True, color=C_WHITE)
asks = [
    ("01", "Approve Phase 1", "เริ่ม implement 10 Tags แรก"),
    ("02", "กำหนด KPI", "วัดความสำเร็จก่อน scale — Conversion, Cost per Acquisition"),
    ("03", "ระบุ Campaign", "Campaign แรกที่จะใช้ทดสอบ"),
]
for i, (num, title, desc) in enumerate(asks):
    x = 0.5 + i * 4.3
    rect(s, x, 6.3, 4.1, 0.85, C_CARD)
    rect(s, x, 6.3, 0.6, 0.85, C_ACCENT)
    txt(s, num, x, 6.38, 0.6, 0.65, size=18, bold=True,
        color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.75, 6.33, 3.2, 0.35, size=12, bold=True, color=C_ACCENT)
    txt(s, desc, x + 0.75, 6.68, 3.2, 0.4, size=10, color=C_LIGHT)

footer(s, 12)

# ─── Slide 13: Closing ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
circ = s.shapes.add_shape(9, Inches(-1.5), Inches(3.5), Inches(6), Inches(6))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x52)
circ.line.fill.background()

lines = [
    ("เรามีข้อมูลลูกค้าอยู่แล้ว", 30, C_WHITE, True),
    ("คำถามคือ  เราอ่านมันได้ดีแค่ไหน", 20, C_LIGHT, False),
    ("", 8, C_WHITE, False),
    ("Customer Tagging คือการเปลี่ยนจาก", 16, C_LIGHT, False),
    ("\"รู้ว่าลูกค้าใช้จ่ายเท่าไหร่\"", 22, C_WHITE, True),
    ("เป็น", 16, C_LIGHT, False),
    ("\"รู้ว่าลูกค้าเป็นใคร\"", 34, C_ACCENT, True),
]
y = 1.5
for text, size, color, bold in lines:
    h = size / 30.0 * 0.65
    txt(s, text, 1.5, y, 10.5, h + 0.1, size=size, color=color,
        bold=bold, align=PP_ALIGN.CENTER)
    y += h + 0.18

txt(s, "พร้อมนำเสนอรายละเอียดเพิ่มเติมเมื่อสะดวก",
    1.5, 6.5, 10.5, 0.4, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

footer(s, 13)

# ─── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_MemoV2_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  (13 main slides — run append_tags_pptx.py to re-add appendix)")
