"""
insert_agenda_keymsg.py
สร้าง Agenda + Key Message slide แทรกที่ position 3
(หลัง Agenda เก่า, ก่อน Big Picture)
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE   = "/Users/adisornj/Desktop/Thena/lab"
TARGET = f"{BASE}/CustomerTagging_FullDeck.pptx"

C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_DARK   = RGBColor(0x0A, 0x14, 0x20)

def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def footer(slide, n, note="Data Science Team | เมษายน 2025"):
    rect(slide, 0, 7.15, 13.33, 0.35, C_DARK)
    txt(slide, note, 0.5, 7.18, 11, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

# ── Build slide ────────────────────────────────────────────────────────────────
tmp  = Presentation()
tmp.slide_width  = Inches(13.33)
tmp.slide_height = Inches(7.5)
s = tmp.slides.add_slide(tmp.slide_layouts[6])

bg(s)

# Header
rect(s, 0, 0, 13.33, 0.75, C_DARK)
rect(s, 0, 0.75, 13.33, 0.05, C_ACCENT)
txt(s, "Agenda — Key Message", 0.5, 0.1, 10, 0.6,
    size=24, bold=True, color=C_WHITE)
txt(s, "สิ่งสำคัญที่เราอยากให้จำได้จากแต่ละ section",
    0.5, 0.55, 10, 0.28, size=11, color=C_ACCENT)

footer(s, 3)

# ── 5 Section cards ────────────────────────────────────────────────────────────
sections = [
    (
        "01", "Big Picture — 5 Layers of Intelligence",
        "ธุรกิจที่ชนะระยะยาวไม่ได้แข่งด้วย cashback — แต่แข่งด้วยความเข้าใจลูกค้าที่ลึกกว่า",
        C_ACCENT
    ),
    (
        "02", "Layer 1: Customer Tags",
        "ทุก transaction ที่ลูกค้ารูด คือข้อมูลที่บอกว่าเขาเป็นใคร — 323 tags คือภาษากลางที่ทั้งองค์กรใช้พูดถึงลูกค้าคนเดียวกัน",
        C_GREEN
    ),
    (
        "03", "Layer 2: Personalized Trigger",
        "Campaign Manager ตั้ง rule ครั้งเดียว — ระบบทำงานแทนทุกวัน",
        C_YELLOW
    ),
    (
        "04", "Roadmap & Next Steps",
        "Layer 1-2 คือรากฐาน — ยิ่งเราเก็บ data นานขึ้น intelligence ยิ่งลึกขึ้น และ competitor ตามยิ่งยาก",
        C_LIGHT
    ),
    (
        "05", "Appendix: Tag Library",
        "323 tags — reference สำหรับทีมที่อยากดูรายละเอียด",
        RGBColor(0x70, 0x70, 0x70)
    ),
]

CARD_L   = 0.4
CARD_W   = 12.33
CARD_H   = 1.02
GAP      = 0.07
START_T  = 0.9

for i, (num, title, keymsg, accent) in enumerate(sections):
    top = START_T + i * (CARD_H + GAP)

    # card bg
    rect(s, CARD_L, top, CARD_W, CARD_H, C_CARD)
    # left accent bar
    rect(s, CARD_L, top, 0.07, CARD_H, accent)

    # number
    txt(s, num, CARD_L + 0.18, top + 0.1, 0.55, 0.38,
        size=20, bold=True, color=accent)

    # title
    txt(s, title, CARD_L + 0.82, top + 0.08, 4.2, 0.38,
        size=13, bold=True, color=C_WHITE)

    # divider
    rect(s, CARD_L + 0.82, top + 0.5, CARD_W - 1.1, 0.02,
         RGBColor(0x2A, 0x40, 0x58))

    # key message
    txt(s, f"Key Message: {keymsg}",
        CARD_L + 0.82, top + 0.56, CARD_W - 1.1, 0.42,
        size=11, color=C_LIGHT)

# ── Merge into FullDeck at position 3 (index 2) ────────────────────────────────
prs = Presentation(TARGET)

blank     = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank)
sp_tree   = new_slide.shapes._spTree
for el in list(sp_tree):
    sp_tree.remove(el)
for shape in s.shapes:
    sp_tree.append(copy.deepcopy(shape.element))

# Move from last → index 2 (position 3)
xml_slides = prs.slides._sldIdLst
entries    = list(xml_slides)
last       = entries[-1]
xml_slides.remove(last)
xml_slides.insert(2, last)

prs.save(TARGET)
print(f"✅  Agenda + Key Message slide inserted at position 3")
print(f"   Total slides: {len(prs.slides)}")
