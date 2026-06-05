"""
build_task_planning_pptx.py
Customer Tag Pipeline — Task Planning PPTX
Slide 1: Cover
Slide 2: Gantt Overview (16 weeks)
Slide 3: Phase Details 1–5
Slide 4: Phase Details 6–10
Slide 5: Milestones
Slide 6: Risk & Buffer
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_TaskPlanning.pptx"

# ── Palette ────────────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_HDR    = RGBColor(0x06, 0x0F, 0x18)
C_DARK   = RGBColor(0x0A, 0x14, 0x20)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_LGRAY  = RGBColor(0x6A, 0x8A, 0x9A)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
C_PURPLE = RGBColor(0x8A, 0x7A, 0xFF)
C_STEEL  = RGBColor(0x7A, 0xB0, 0xCC)
C_TEAL   = RGBColor(0x00, 0xC5, 0x8A)
C_RED    = RGBColor(0xFF, 0x6B, 0x6B)
C_BLUE   = RGBColor(0x4A, 0xA5, 0xFF)

# ── Phase definitions ─────────────────────────────────────────────────────────
#  (label, week_start, week_end, color, duration_text, tasks)
PHASES = [
    ("Foundation",           1,  1,  C_STEEL,  "1 สัปดาห์",
     ["Database schema design", "Pipeline framework setup", "Dev/Test environment"]),

    ("Scoring Engine",       2,  2,  C_ACCENT, "1 สัปดาห์",
     ["Core formula: 0.4×freq + 0.4×recency + 0.2×breadth",
      "Tier / Trending / Recency classifiers", "Unit tests"]),

    ("Path A (221 tags)",    3,  6,  C_GREEN,  "4 สัปดาห์",
     ["Batch 1 (W3–4): Food, Travel, Shopping, Entertainment (~95 tags)",
      "Batch 2 (W5–6): Transport, Health, Financial, Life Stage, Digital (~80 tags)",
      "Batch 3 (W5–6): Education, RE, Investment, Payment, Card, Momentum (~46 tags)"]),

    ("Sanity Validation",    6,  6,  C_YELLOW, "3 วัน",
     ["Coverage check (tag > 1% population)", "Score distribution histogram",
      "Null rate check (< 5%)"]),

    ("Quality + Stability",  7,  8,  C_ORANGE, "2 สัปดาห์",
     ["Tag churn rate (2 consecutive windows)", "PSI — distribution shift",
      "Threshold sensitivity analysis"]),

    ("Business Signal",      8,  8,  C_PURPLE, "3 วัน",
     ["Behavioral diff: tagged vs non-tagged", "Lift analysis per tag category"]),

    ("Profile Pipeline",     7,  7,  C_TEAL,   "1 สัปดาห์",
     ["Daily refresh pipeline", "Performance test (full portfolio)",
      "Monitoring + alerting"]),

    ("Path B (67 tags)",     9,  10, C_RED,    "1.5 สัปดาห์",
     ["⚠️  ขึ้นกับ Merchant Enrichment team",
      "Map enrichment → customer tag condition",
      "Implement + sanity check 67 tags"]),

    ("Campaign Integration", 10, 11, C_BLUE,   "2 สัปดาห์",
     ["Connect tag profile → campaign tool", "Query template library (10 templates)",
      "UAT + train campaign managers"]),

    ("Pilot Campaign",       12, 15, C_YELLOW, "4 สัปดาห์",
     ["เลือก 10 tags สำหรับ pilot", "ออกแบบ A/B test + control group",
      "Run campaign + วัดผล + สรุป"]),
]

MILESTONES = [
    (2,  "M1", "Engine Ready",      "Scoring engine + pipeline ทำงานได้"),
    (6,  "M2", "Path A Done",       "221 tags ใน pipeline, sanity ผ่าน"),
    (8,  "M3", "Validated",         "Quality, stability, biz signal ผ่าน"),
    (11, "M4", "Integrated",        "Campaign manager ใช้งานได้จริง"),
    (15, "M5", "Pilot Result",      "มีผล A/B test พร้อม present"),
]

RISKS = [
    ("Validation พบ tag ที่ไม่ผ่าน", "สูง", "+1–2 สัปดาห์",
     "Build buffer ใน Phase 5", C_ORANGE),
    ("Enrichment team ช้า (Path B)", "กลาง", "Path B ล่าช้า",
     "Path A ออก production ก่อน — enrichment ตามทีหลัง", C_YELLOW),
    ("Data quality issues", "กลาง", "+1 สัปดาห์",
     "Sanity check Phase 4 จับก่อนถึง Validation", C_YELLOW),
    ("Campaign integration ซับซ้อน", "ต่ำ", "+1 สัปดาห์",
     "UAT เริ่มเร็ว ก่อน M4", C_GREEN),
]

# ── Gantt layout constants ────────────────────────────────────────────────────
N_WEEKS  = 16
CHART_L  = 2.72   # inches — left edge of Gantt area
CHART_W  = 10.3   # inches — total Gantt width
WEEK_W   = CHART_W / N_WEEKS  # 0.644" per week

# Month labels: May W1–4, Jun W5–8, Jul W9–12, Aug W13–16
MONTHS = [("พ.ค.", 0, 4), ("มิ.ย.", 4, 8), ("ก.ค.", 8, 12), ("ส.ค.", 12, 16)]

# ── pptx helpers ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(sl, l, t, w, h, color, alpha=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h, size=11, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, font="Helvetica Neue", wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb

def multirun(sl, items, l, t, w, h, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for text, size, bold, color in items:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Helvetica Neue"

def header(sl, title, subtitle=""):
    rect(sl, 0, 0, 13.33, 7.5, C_BG)
    rect(sl, 0, 0, 13.33, 0.62, C_HDR)
    rect(sl, 0, 0, 0.08,  0.62, C_ACCENT)
    txt(sl, title, 0.2, 0.08, 9, 0.48, size=18, bold=True)
    if subtitle:
        txt(sl, subtitle, 9.5, 0.14, 3.6, 0.38, size=11,
            color=C_ACCENT, align=PP_ALIGN.RIGHT)

def footer(sl, n):
    rect(sl, 0, 7.14, 13.33, 0.36, C_DARK)
    txt(sl, "Customer Tag Pipeline — Task Planning  |  Data Science Team  |  2026",
        0.3, 7.17, 10, 0.28, size=8, color=C_LGRAY)
    txt(sl, str(n), 12.8, 7.17, 0.4, 0.28, size=8, color=C_LGRAY, align=PP_ALIGN.RIGHT)

def bullet_block(sl, items, x, y, w, gap=0.36, size=11, dot_color=C_ACCENT):
    for i, item in enumerate(items):
        yt = y + i * gap
        txt(sl, "•", x, yt, 0.25, gap, size=size, bold=True, color=dot_color)
        txt(sl, item, x + 0.22, yt, w - 0.22, gap, size=size, color=C_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, C_BG)
circ = s.shapes.add_shape(9, Inches(8.0), Inches(-1.0), Inches(7.5), Inches(7.5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x18, 0x28)
circ.line.fill.background()
rect(s, 0, 2.6, 0.12, 2.3, C_ACCENT)

txt(s, "Customer Tag Pipeline", 0.5, 1.1, 11, 1.0, size=44, bold=True)
txt(s, "Task Planning & Timeline", 0.5, 2.25, 10, 0.7, size=24, color=C_ACCENT)
txt(s, "ทีม Data Science  |  พ.ค.–ส.ค. 2026  |  16 สัปดาห์", 0.5, 3.1, 9, 0.5, size=14, color=C_LIGHT)

# Summary stat cards
cards = [
    ("10", "Phases", C_ACCENT),
    ("16", "Weeks", C_GREEN),
    ("288", "Tags", C_YELLOW),
    ("5", "Milestones", C_ORANGE),
]
for i, (num, label, col) in enumerate(cards):
    x = 0.5 + i * 2.2
    rect(s, x, 4.3, 2.0, 1.6, RGBColor(0x0A, 0x18, 0x28))
    rect(s, x, 4.3, 2.0, 0.07, col)
    txt(s, num,   x+0.15, 4.5,  1.7, 0.75, size=38, bold=True, color=col)
    txt(s, label, x+0.15, 5.32, 1.7, 0.42, size=13, color=C_LIGHT)

footer(s, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Gantt Overview
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Timeline Overview — 16 สัปดาห์", "พ.ค.–ส.ค. 2026")

GANTT_T = 0.68   # top of gantt area
ROW_H   = 0.52   # height per phase row
ROW_GAP = 0.04
LABEL_W = 2.45   # width of phase label column

# Month header boxes
month_y = GANTT_T
for month, w_start, w_end in MONTHS:
    bx = CHART_L + w_start * WEEK_W
    bw = (w_end - w_start) * WEEK_W
    rect(s, bx, month_y, bw - 0.04, 0.30, RGBColor(0x0A, 0x14, 0x20))
    txt(s, month, bx, month_y, bw - 0.04, 0.30,
        size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Week tick numbers
tick_y = month_y + 0.32
for w in range(1, N_WEEKS + 1):
    wx = CHART_L + (w - 1) * WEEK_W
    txt(s, str(w), wx, tick_y, WEEK_W, 0.18,
        size=7.5, color=C_LGRAY, align=PP_ALIGN.CENTER)

# Vertical week grid lines
grid_top = tick_y + 0.19
grid_bot = grid_top + len(PHASES) * (ROW_H + ROW_GAP) + 0.22
for w in range(N_WEEKS + 1):
    wx = CHART_L + w * WEEK_W
    line_color = RGBColor(0x1A, 0x2E, 0x44) if w % 4 != 0 else RGBColor(0x2A, 0x44, 0x60)
    lw = 0.008 if w % 4 != 0 else 0.014
    rect(s, wx, grid_top, lw, grid_bot - grid_top, line_color)

# Phase rows
for i, (name, ws, we, color, dur, _) in enumerate(PHASES):
    row_y = grid_top + i * (ROW_H + ROW_GAP)

    # Row background
    rect(s, 0.15, row_y, CHART_L - 0.2 + CHART_W, ROW_H,
         RGBColor(0x0A, 0x18, 0x28))

    # Phase label (left)
    rect(s, 0.15, row_y, 0.06, ROW_H, color)
    ph_num = str(i + 1)
    txt(s, ph_num, 0.25, row_y + 0.06, 0.3, ROW_H - 0.1,
        size=9, bold=True, color=color)
    txt(s, name, 0.57, row_y + 0.07, LABEL_W - 0.45, ROW_H - 0.1,
        size=9.5, bold=False, color=C_WHITE)

    # Duration pill
    txt(s, dur, 0.57, row_y + ROW_H - 0.22, LABEL_W - 0.45, 0.2,
        size=7.5, color=C_LGRAY)

    # Gantt bar
    bx = CHART_L + (ws - 1) * WEEK_W + 0.03
    bw = max((we - ws + 1) * WEEK_W - 0.06, WEEK_W * 0.6)
    bar_h = ROW_H - 0.10
    bar_y = row_y + 0.05

    # Bar background (subtle)
    rect(s, CHART_L, row_y + 0.02, CHART_W, ROW_H - 0.04,
         RGBColor(0x08, 0x14, 0x20))

    # Main bar
    rect(s, bx, bar_y, bw, bar_h, color)

    # Bar label (week range)
    wk_label = f"W{ws}" if ws == we else f"W{ws}–{we}"
    txt(s, wk_label, bx + 0.06, bar_y + 0.04, bw - 0.08, bar_h - 0.06,
        size=8, bold=True, color=C_BG, align=PP_ALIGN.LEFT)

    # Dependency flag for Path B
    if "⚠️" in (PHASES[i][5][0] if PHASES[i][5] else ""):
        txt(s, "⚠️ ext.", bx + bw - 0.7, bar_y + 0.04, 0.65, bar_h - 0.06,
            size=7, color=RGBColor(0xFF, 0xD7, 0x00), align=PP_ALIGN.RIGHT)

# Milestone markers (diamond shape at bottom)
ms_y = grid_bot + 0.01
rect(s, CHART_L, ms_y, CHART_W, 0.28, RGBColor(0x06, 0x10, 0x1C))
txt(s, "MILESTONE", CHART_L, ms_y + 0.02, 1.0, 0.22,
    size=7, bold=True, color=C_LGRAY)

m_colors = [C_ACCENT, C_GREEN, C_ORANGE, C_PURPLE, C_YELLOW]
for (wk, code, short, _), col in zip(MILESTONES, m_colors):
    mx = CHART_L + wk * WEEK_W - 0.02
    rect(s, mx, ms_y, 0.04, 0.28, col)
    txt(s, code, mx - 0.28, ms_y + 0.03, 0.6, 0.22,
        size=8, bold=True, color=col, align=PP_ALIGN.CENTER)

footer(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Phase Details 1–5
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Phase Details — Foundation → Validation", "Phase 1–6")

CARD_TOP = 0.75
CARD_H   = 1.52
CARD_W   = 3.92
CARD_GAP = 0.18

phase_subset = PHASES[:6]  # phases 1-6
for i, (name, ws, we, color, dur, tasks) in enumerate(phase_subset):
    col = i % 3
    row = i // 3
    cx = 0.2 + col * (CARD_W + CARD_GAP)
    cy = CARD_TOP + row * (CARD_H + 0.2)

    # Card background
    rect(s, cx, cy, CARD_W, CARD_H, RGBColor(0x0A, 0x18, 0x28))
    rect(s, cx, cy, CARD_W, 0.07, color)

    # Phase number + name
    wk_label = f"W{ws}" if ws == we else f"W{ws}–{we}"
    txt(s, f"Phase {i+1}  ·  {wk_label}  ·  {dur}",
        cx+0.12, cy+0.10, CARD_W-0.2, 0.22, size=8, color=color)
    txt(s, name, cx+0.12, cy+0.30, CARD_W-0.2, 0.36,
        size=13, bold=True, color=C_WHITE)

    # Divider
    rect(s, cx+0.12, cy+0.68, CARD_W-0.24, 0.03, RGBColor(0x1A, 0x2E, 0x44))

    # Tasks
    for j, task in enumerate(tasks[:3]):
        ty = cy + 0.76 + j * 0.24
        txt(s, "›", cx+0.12, ty, 0.18, 0.22, size=9, bold=True, color=color)
        txt(s, task, cx+0.28, ty, CARD_W-0.38, 0.22, size=8.5, color=C_LIGHT)

footer(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Phase Details 6–10
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Phase Details — Profile → Pilot", "Phase 7–10")

phase_subset2 = PHASES[6:]  # phases 7-10
CARD_W2 = 5.95
CARD_H2 = 2.5
CARD_GAP2 = 0.25

for i, (name, ws, we, color, dur, tasks) in enumerate(phase_subset2):
    col = i % 2
    row = i // 2
    cx = 0.2 + col * (CARD_W2 + CARD_GAP2)
    cy = 0.75 + row * (CARD_H2 + 0.2)

    rect(s, cx, cy, CARD_W2, CARD_H2, RGBColor(0x0A, 0x18, 0x28))
    rect(s, cx, cy, CARD_W2, 0.08, color)

    ph_n = i + 7
    wk_label = f"W{ws}" if ws == we else f"W{ws}–{we}"
    txt(s, f"Phase {ph_n}  ·  {wk_label}  ·  {dur}",
        cx+0.15, cy+0.12, CARD_W2-0.25, 0.24, size=9, color=color)
    txt(s, name, cx+0.15, cy+0.35, CARD_W2-0.25, 0.44,
        size=16, bold=True, color=C_WHITE)

    # External dep badge for Path B
    if ph_n == 8:
        rect(s, cx + CARD_W2 - 1.5, cy + 0.14, 1.35, 0.22,
             RGBColor(0x2A, 0x18, 0x00))
        txt(s, "⚠️  External Dep.", cx + CARD_W2 - 1.48, cy + 0.15,
            1.3, 0.20, size=7.5, color=C_YELLOW)

    rect(s, cx+0.15, cy+0.84, CARD_W2-0.3, 0.03, RGBColor(0x1A, 0x2E, 0x44))

    for j, task in enumerate(tasks):
        ty = cy + 0.92 + j * 0.44
        txt(s, "›", cx+0.15, ty, 0.2, 0.38, size=10, bold=True, color=color)
        txt(s, task, cx+0.33, ty, CARD_W2-0.45, 0.38, size=10.5, color=C_LIGHT)

footer(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Milestones
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Milestones — 5 จุดหมาย สำคัญ")

m_colors2 = [C_ACCENT, C_GREEN, C_ORANGE, C_PURPLE, C_YELLOW]
M_W = 2.3
M_H = 4.8

# Timeline connector bar
rect(s, 0.45, 3.85, 12.4, 0.06, RGBColor(0x1A, 0x2E, 0x44))

for i, ((wk, code, short, desc), col) in enumerate(zip(MILESTONES, m_colors2)):
    mx = 0.45 + i * (M_W + 0.1)

    # Milestone card
    rect(s, mx, 0.75, M_W, M_H, RGBColor(0x0A, 0x18, 0x28))
    rect(s, mx, 0.75, M_W, 0.07, col)

    # Code badge
    rect(s, mx + 0.12, 0.86, 0.62, 0.38, col)
    txt(s, code, mx+0.12, 0.86, 0.62, 0.38, size=14, bold=True,
        color=C_BG, align=PP_ALIGN.CENTER)

    # Week label
    txt(s, f"สิ้นสัปดาห์ที่ {wk}", mx+0.82, 0.90, M_W-0.9, 0.30,
        size=9, color=col)

    # Title
    txt(s, short, mx+0.12, 1.32, M_W-0.22, 0.50, size=14, bold=True)

    # Divider
    rect(s, mx+0.12, 1.86, M_W-0.24, 0.03, RGBColor(0x1A, 0x2E, 0x44))

    # Description
    txt(s, desc, mx+0.12, 1.95, M_W-0.22, 0.8, size=10.5, color=C_LIGHT)

    # Circle on timeline
    circ_x = mx + M_W/2 - 0.14
    circ_y = 3.71
    c = s.shapes.add_shape(9, Inches(circ_x), Inches(circ_y),
                           Inches(0.28), Inches(0.28))
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.line.fill.background()

    # Week number on timeline
    txt(s, f"W{wk}", mx + M_W/2 - 0.25, 3.99, 0.5, 0.22,
        size=9, bold=True, color=col, align=PP_ALIGN.CENTER)

    # What it means
    txt(s, "ทำไมสำคัญ", mx+0.12, 2.84, M_W-0.22, 0.22,
        size=8, bold=True, color=C_LGRAY)
    why_texts = [
        "Pipeline พร้อม\nสำหรับ tag logic",
        "เริ่ม validate\nได้ทันที",
        "Tags พร้อมใช้\nใน production",
        "Campaign team\nใช้งานได้จริง",
        "มีผลพิสูจน์\nให้ stakeholder",
    ]
    txt(s, why_texts[i], mx+0.12, 3.06, M_W-0.22, 0.55,
        size=9.5, color=C_LIGHT)

footer(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Risk & Buffer
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Risk & Buffer Planning")

# Risk table
RISK_COLS = ["Risk", "โอกาส", "Impact", "Mitigation"]
RISK_CX = [0.2, 5.2, 6.5, 8.0]
RISK_CW = [4.8, 1.1, 1.3, 5.0]
RISK_TH = 0.36
RISK_RH = 0.90

# Table header
rect(s, 0.2, 0.75, 12.9, RISK_TH, C_HDR)
for j, (label, cx, cw) in enumerate(zip(RISK_COLS, RISK_CX, RISK_CW)):
    txt(s, label, cx+0.08, 0.78, cw, RISK_TH-0.04,
        size=9, bold=True, color=C_LGRAY)

# Risk rows
for i, (risk, prob, impact, mitigation, col) in enumerate(RISKS):
    row_y = 0.75 + RISK_TH + i * (RISK_RH + 0.06)
    bg = RGBColor(0x0A, 0x18, 0x28) if i % 2 == 0 else RGBColor(0x0D, 0x1E, 0x30)
    rect(s, 0.2, row_y, 12.9, RISK_RH, bg)
    rect(s, 0.2, row_y, 0.07, RISK_RH, col)

    texts = [risk, prob, impact, mitigation]
    for j, (t, cx, cw) in enumerate(zip(texts, RISK_CX, RISK_CW)):
        c = col if j in [1, 2] else C_LIGHT
        b = True if j == 0 else False
        txt(s, t, cx+0.1, row_y+0.12, cw-0.12, RISK_RH-0.2,
            size=10, bold=b, color=c)

# Overall buffer note
ry = 0.75 + RISK_TH + len(RISKS) * (RISK_RH + 0.06) + 0.15
rect(s, 0.2, ry, 12.9, 1.4, RGBColor(0x00, 0x1C, 0x30))
rect(s, 0.2, ry, 0.08, 1.4, C_ACCENT)

txt(s, "Overall Buffer Recommendation", 0.38, ry+0.1, 7, 0.35,
    size=12, bold=True, color=C_ACCENT)
buffer_pts = [
    "เพิ่ม 2 สัปดาห์ buffer หลัง Phase 5 (Validation) สำหรับ fix tag ที่ไม่ผ่าน",
    "Path B ไม่ block Phase A — ออก production ได้ก่อน แล้ว Path B ตามเมื่อ enrichment พร้อม",
    "Pilot Campaign (Phase 10) สามารถเริ่มด้วย 5 tags ก่อน ไม่ต้องรอครบ 10 tags",
]
for k, pt in enumerate(buffer_pts):
    txt(s, "•", 0.38, ry+0.50+k*0.28, 0.2, 0.26, size=10, bold=True, color=C_ACCENT)
    txt(s, pt, 0.56, ry+0.50+k*0.28, 12.3, 0.26, size=10.5, color=C_LIGHT)

footer(s, 6)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  CustomerTagging_TaskPlanning.pptx — 6 slides")
print(f"   Slide 1: Cover")
print(f"   Slide 2: Gantt Overview (16 weeks)")
print(f"   Slide 3: Phase Details 1–6")
print(f"   Slide 4: Phase Details 7–10")
print(f"   Slide 5: Milestones (5 จุด)")
print(f"   Slide 6: Risk & Buffer")
