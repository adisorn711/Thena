from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x0F, 0x3D, 0x6B)
GOLD = RGBColor(0xE2, 0xB9, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)
MID_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
TEAL = RGBColor(0x1A, 0x85, 0x7E)
SOFT_BLUE = RGBColor(0x2E, 0x86, 0xC1)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_text_multiline(slide, lines, left, top, width, height,
                        font_size=12, bold=False, color=WHITE,
                        align=PP_ALIGN.LEFT, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────
# SLIDE 1: COVER
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

# Full background
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=NAVY)

# Gold accent bar left
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

# Title area box
add_rect(slide, 0.5, 1.8, 9, 1.1, fill_color=ACCENT)

add_text(slide, "D-Tag Framework", 0.6, 1.85, 8.5, 1.0,
         font_size=40, bold=True, color=WHITE)

add_text(slide, "Customer Intelligence Architecture for Super Customization",
         0.6, 3.1, 9, 0.5, font_size=16, color=GOLD, italic=True)

add_text(slide, "ระบบ Tag ลูกค้าอัจฉริยะ — ได้รับแรงบันดาลใจจาก Hyundai Card",
         0.6, 3.7, 9, 0.5, font_size=14, color=MID_GRAY)

# Bottom info bar
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_color=RGBColor(0x0A, 0x0A, 0x1A))
add_text(slide, "Data Science Team  |  2026", 0.6, 6.85, 6, 0.4,
         font_size=11, color=MID_GRAY)


# ─────────────────────────────────────────────
# SLIDE 2: PROBLEM — WHY D-TAG
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "ทำไมต้อง D-Tag?", 0.3, 0.15, 10, 0.8,
         font_size=26, bold=True, color=WHITE)

# Old model box
add_rect(slide, 0.5, 1.4, 5.5, 5.2, fill_color=WHITE,
         line_color=MID_GRAY, line_width=1)
add_rect(slide, 0.5, 1.4, 5.5, 0.55, fill_color=RGBColor(0xCC, 0x33, 0x33))
add_text(slide, "Traditional Segmentation", 0.7, 1.45, 5, 0.45,
         font_size=14, bold=True, color=WHITE)

old_lines = [
    "• ลูกค้า 1 คน = 1 กลุ่มเท่านั้น",
    "• แบ่งจาก demographics เป็นหลัก",
    "• Static — ไม่อัปเดตตาม behavior",
    "• One-size-fits-all ภายในกลุ่ม",
    "• ไม่สามารถจับ nuance ได้",
]
add_text_multiline(slide, old_lines, 0.7, 2.1, 5, 4.2,
                   font_size=13, color=DARK_GRAY)

# Arrow
add_text(slide, "→", 6.2, 3.5, 0.7, 0.6, font_size=30, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)

# New model box
add_rect(slide, 7.0, 1.4, 5.8, 5.2, fill_color=WHITE,
         line_color=TEAL, line_width=2)
add_rect(slide, 7.0, 1.4, 5.8, 0.55, fill_color=TEAL)
add_text(slide, "D-Tag Approach", 7.2, 1.45, 5.3, 0.45,
         font_size=14, bold=True, color=WHITE)

new_lines = [
    "• ลูกค้า 1 คน = หลาย Tag พร้อมกัน",
    "• ใช้ transaction behavior จริง",
    "• Dynamic — refresh รายเดือน",
    "• Hyper-personalized ระดับบุคคล",
    "• Confidence score ทุก tag",
]
add_text_multiline(slide, new_lines, 7.2, 2.1, 5.3, 4.2,
                   font_size=13, color=DARK_GRAY)

add_rect(slide, 0, 6.9, 13.33, 0.1, fill_color=GOLD)


# ─────────────────────────────────────────────
# SLIDE 3: CORE CONCEPT
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "แนวคิดหลัก: Super Customization", 0.3, 0.15, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

# Center example customer
add_rect(slide, 4.5, 1.3, 4.33, 0.7, fill_color=GOLD)
add_text(slide, "ลูกค้า 1 คน", 4.5, 1.3, 4.33, 0.7,
         font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Tags around the customer
tags = [
    ("foodie", 1.0, 2.5, TEAL),
    ("traveler", 1.0, 3.5, SOFT_BLUE),
    ("commuter", 1.0, 4.5, ACCENT),
    ("online_shopper", 6.8, 2.5, TEAL),
    ("budget_conscious", 6.8, 3.5, SOFT_BLUE),
    ("weekend_spender", 6.8, 4.5, ACCENT),
    ("starbucks_loyalist", 3.5, 5.6, RGBColor(0x6B, 0x3A, 0x2A)),
    ("luxury_travel", 5.8, 5.6, RGBColor(0x7D, 0x3C, 0x98)),
]

for tag_name, lft, tp, color in tags:
    add_rect(slide, lft, tp, 2.6, 0.5, fill_color=color)
    add_text(slide, tag_name, lft, tp, 2.6, 0.5,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 dimensions
add_rect(slide, 0.5, 6.1, 3.8, 0.8, fill_color=RGBColor(0x0A, 0x0A, 0x1A))
add_text(slide, "WHEN — ซื้อตอนไหน", 0.5, 6.1, 3.8, 0.8,
         font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_rect(slide, 4.6, 6.1, 3.8, 0.8, fill_color=RGBColor(0x0A, 0x0A, 0x1A))
add_text(slide, "WHAT — ซื้ออะไร", 4.6, 6.1, 3.8, 0.8,
         font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_rect(slide, 8.7, 6.1, 4.1, 0.8, fill_color=RGBColor(0x0A, 0x0A, 0x1A))
add_text(slide, "CHANNEL — ซื้อที่ไหน", 8.7, 6.1, 4.1, 0.8,
         font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 4: FOUNDATION — WHAT WE HAVE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "สิ่งที่มีอยู่แล้ว — Foundation", 0.3, 0.15, 10, 0.8,
         font_size=26, bold=True, color=WHITE)

# Layer 1
add_rect(slide, 1.0, 1.5, 11.0, 1.4, fill_color=ACCENT)
add_text(slide, "Layer 1 — MCC (Merchant Category Code)",
         1.2, 1.55, 8, 0.45, font_size=14, bold=True, color=GOLD)
add_text(slide, "จัดกลุ่ม merchant ตามประเภท: ร้านอาหาร / เติมน้ำมัน / ท่องเที่ยว / Entertainment ...",
         1.2, 2.05, 9, 0.6, font_size=12, color=WHITE)

# Arrow down
add_text(slide, "↓", 6.4, 3.0, 0.8, 0.5, font_size=24, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)

# Layer 2
add_rect(slide, 1.0, 3.55, 11.0, 1.4, fill_color=TEAL)
add_text(slide, "Layer 2 — Centralized Merchant Name",
         1.2, 3.6, 8, 0.45, font_size=14, bold=True, color=WHITE)
add_text(slide, "รวมทุกสาขาของ chain เดียวกันภายใต้ชื่อเดียว: McDonald's / Starbucks / PTT ...",
         1.2, 4.1, 9, 0.6, font_size=12, color=WHITE)

# Arrow down
add_text(slide, "↓", 6.4, 5.05, 0.8, 0.5, font_size=24, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER)

# Layer 3 — what we'll build
add_rect(slide, 1.0, 5.6, 11.0, 1.35, fill_color=WHITE,
         line_color=GOLD, line_width=2)
add_text(slide, "Layer 3 — D-Tag (สิ่งที่จะสร้าง)",
         1.2, 5.65, 8, 0.45, font_size=14, bold=True, color=GOLD)
add_text(slide, "Behavioral tags ที่ derive จาก 2 layers ข้างบน — Multi-dimensional, Dynamic, Confidence-scored",
         1.2, 6.1, 9.5, 0.6, font_size=12, color=DARK_GRAY)

add_rect(slide, 0, 6.9, 13.33, 0.1, fill_color=GOLD)


# ─────────────────────────────────────────────
# SLIDE 5: TAG TAXONOMY
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "Tag Taxonomy — 5 มิติ", 0.3, 0.15, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

dims = [
    ("1  Lifestyle", ["foodie", "traveler", "commuter", "homebody", "health_conscious"], TEAL),
    ("2  Merchant Affinity", ["brand_loyalist", "explorer", "chain_preferrer"], SOFT_BLUE),
    ("3  Spending Pattern", ["weekend_spender", "online_first", "high_freq_low_ticket"], ACCENT),
    ("4  Value Tier", ["luxury_{category}", "budget_{category}", "mid_{category}"], RGBColor(0x7D, 0x3C, 0x98)),
    ("5  Engagement", ["active_12m", "seasonal", "dormant_recovering", "single_category"], RGBColor(0x1E, 0x8B, 0x4F)),
]

cols = [(0.3, 2.6), (2.95, 2.6), (5.6, 2.6), (8.25, 2.3), (10.6, 2.65)]

for i, ((title, tags, color), (lft, wid)) in enumerate(zip(dims, cols)):
    add_rect(slide, lft, 1.2, wid, 0.6, fill_color=color)
    add_text(slide, title, lft + 0.1, 1.22, wid - 0.15, 0.55,
             font_size=11, bold=True, color=WHITE)

    add_rect(slide, lft, 1.85, wid, 5.35, fill_color=RGBColor(0x0F, 0x1A, 0x30))
    for j, tag in enumerate(tags):
        y = 2.0 + j * 0.95
        add_rect(slide, lft + 0.1, y, wid - 0.2, 0.75, fill_color=color)
        add_text(slide, tag, lft + 0.1, y, wid - 0.2, 0.75,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 6: TAG MECHANICS
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "Tag Mechanics — วิธีทำงาน", 0.3, 0.15, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

mechanics = [
    ("Tag Type", "Binary\n(มี / ไม่มี)", TEAL),
    ("Confidence Score", "0.0 → 1.0\nทุก tag มี score แนบ", SOFT_BLUE),
    ("Lookback Window", "6 เดือน (default)\n12 เดือน (lifestyle)", ACCENT),
    ("Refresh Cycle", "รายเดือน\nอัตโนมัติ", RGBColor(0x7D, 0x3C, 0x98)),
    ("Threshold", "ปรับได้\nตาม use case", RGBColor(0x1E, 0x8B, 0x4F)),
]

for i, (label, value, color) in enumerate(mechanics):
    lft = 0.5 + i * 2.55
    add_rect(slide, lft, 1.5, 2.3, 0.65, fill_color=color)
    add_text(slide, label, lft, 1.5, 2.3, 0.65,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, lft, 2.2, 2.3, 1.3, fill_color=WHITE,
             line_color=color, line_width=1)
    add_text_multiline(slide, value.split("\n"), lft, 2.2, 2.3, 1.3,
                       font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Example
add_rect(slide, 0.5, 3.9, 12.3, 3.1, fill_color=WHITE,
         line_color=MID_GRAY, line_width=1)
add_rect(slide, 0.5, 3.9, 12.3, 0.55, fill_color=NAVY)
add_text(slide, "ตัวอย่าง: ลูกค้า ID 00123", 0.7, 3.93, 10, 0.45,
         font_size=13, bold=True, color=GOLD)

example_lines = [
    '  Tag: "foodie"          Score: 0.91   Lookback: 6M   Last updated: 2026-04',
    '  Tag: "traveler"        Score: 0.74   Lookback: 12M  Last updated: 2026-04',
    '  Tag: "online_first"    Score: 0.68   Lookback: 6M   Last updated: 2026-04',
    '  Tag: "luxury_dining"   Score: 0.55   Lookback: 6M   Last updated: 2026-04',
]
add_text_multiline(slide, example_lines, 0.7, 4.55, 12, 2.3,
                   font_size=12, color=DARK_GRAY)

add_rect(slide, 0, 6.9, 13.33, 0.1, fill_color=GOLD)


# ─────────────────────────────────────────────
# SLIDE 7: USE CASES
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "Use Cases — ใช้ทำอะไรได้บ้าง", 0.3, 0.15, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

use_cases = [
    ("Campaign Targeting", "เลือก audience ที่ใช่\nจาก tag combination", "Lifestyle\n+ Merchant Affinity", TEAL),
    ("Product Design", "ออกแบบ card benefit\nให้ตรงกับ behavior จริง", "Value Tier\n+ Lifestyle", SOFT_BLUE),
    ("Credit Limit", "ปรับ limit ตาม\nengagement pattern", "Engagement\n+ Spending Pattern", ACCENT),
    ("Churn Prevention", "จับสัญญาณเสี่ยง\nก่อนลูกค้าจะหาย", "Engagement\n+ dormant signals", RGBColor(0xCC, 0x44, 0x44)),
]

for i, (title, desc, tags_used, color) in enumerate(use_cases):
    lft = 0.5 + i * 3.2
    add_rect(slide, lft, 1.3, 2.95, 5.7, fill_color=RGBColor(0x0F, 0x1A, 0x30))
    add_rect(slide, lft, 1.3, 2.95, 0.65, fill_color=color)
    add_text(slide, title, lft, 1.3, 2.95, 0.65,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_multiline(slide, desc.split("\n"), lft + 0.1, 2.1, 2.75, 1.2,
                       font_size=12, color=WHITE)

    add_rect(slide, lft + 0.1, 3.5, 2.75, 0.35, fill_color=RGBColor(0x0A, 0x0A, 0x1A))
    add_text(slide, "Tags ที่ใช้:", lft + 0.1, 3.5, 2.75, 0.35,
             font_size=10, color=GOLD, bold=True)

    add_text_multiline(slide, tags_used.split("\n"), lft + 0.1, 3.95, 2.75, 1.5,
                       font_size=11, color=MID_GRAY)


# ─────────────────────────────────────────────
# SLIDE 8: ROADMAP
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

add_text(slide, "Roadmap — 4 Phases", 0.3, 0.15, 12, 0.8,
         font_size=26, bold=True, color=WHITE)

phases = [
    ("Phase 1", "MCC + Merchant Name\n→ Lifestyle & Affinity Tags", "ของที่มีอยู่แล้ว\nเริ่มได้ทันที", TEAL),
    ("Phase 2", "เพิ่ม Timing Dimension\n→ WHEN tags", "Transaction timestamp\nanalysis", SOFT_BLUE),
    ("Phase 3", "เพิ่ม Channel Dimension\n→ Online/Offline tags", "Channel data\nintegration", ACCENT),
    ("Phase 4", "ML-based Tags\nPattern ที่ rule-based จับไม่ได้", "Clustering &\nAnomaly Detection", RGBColor(0x7D, 0x3C, 0x98)),
]

# Timeline bar
add_rect(slide, 1.0, 2.1, 11.2, 0.12, fill_color=MID_GRAY)

for i, (phase, desc, data_needed, color) in enumerate(phases):
    lft = 0.7 + i * 3.0

    # Circle on timeline
    circle = slide.shapes.add_shape(9, Inches(lft + 1.1), Inches(1.85),
                                     Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # Phase card
    add_rect(slide, lft, 2.6, 2.7, 4.3, fill_color=WHITE,
             line_color=color, line_width=2)
    add_rect(slide, lft, 2.6, 2.7, 0.6, fill_color=color)
    add_text(slide, phase, lft, 2.6, 2.7, 0.6,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_multiline(slide, desc.split("\n"), lft + 0.1, 3.3, 2.5, 1.5,
                       font_size=12, color=DARK_GRAY)

    add_rect(slide, lft + 0.1, 4.9, 2.5, 0.3, fill_color=LIGHT_GRAY)
    add_text(slide, "Data:", lft + 0.1, 4.9, 2.5, 0.3,
             font_size=10, color=ACCENT, bold=True)
    add_text_multiline(slide, data_needed.split("\n"), lft + 0.1, 5.25, 2.5, 1.3,
                       font_size=11, color=DARK_GRAY)

add_rect(slide, 0, 6.9, 13.33, 0.1, fill_color=GOLD)


# ─────────────────────────────────────────────
# SLIDE 9: CLOSING
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=NAVY)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=GOLD)

# Gold accent box
add_rect(slide, 3.0, 1.5, 7.33, 4.0, fill_color=RGBColor(0x0F, 0x1A, 0x30))

add_text(slide, "D-Tag", 3.0, 1.6, 7.33, 1.2,
         font_size=48, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

summary_lines = [
    "จากลูกค้า 'กลุ่ม' → ลูกค้า 'บุคคล'",
    "Multi-dimensional  ·  Dynamic  ·  Confidence-scored",
    "Built on MCC + Merchant Name ที่มีอยู่แล้ว",
]
add_text_multiline(slide, summary_lines, 3.0, 2.9, 7.33, 2.3,
                   font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Data Science Team  |  2026",
         0.5, 6.85, 12, 0.4, font_size=11, color=MID_GRAY,
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
out = "/Users/adisornj/Desktop/Thena/lab/D-Tag_Framework.pptx"
prs.save(out)
print(f"Saved: {out}")
