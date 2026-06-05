"""
build_business_deck.py
CustomerTagging_Business.pptx — 6 slides, WH framework, Business audience
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_Business.pptx"

C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
C_DARK   = RGBColor(0x0A, 0x14, 0x20)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def wh_tag(slide, label, color):
    """WH label top-right"""
    rect(slide, 11.8, 0.15, 1.38, 0.45, color)
    txt(slide, label, 11.8, 0.15, 1.38, 0.45,
        size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

def footer(slide, n):
    rect(slide, 0, 7.15, 13.33, 0.35, C_DARK)
    txt(slide, "Data Science Team  |  เมษายน 2025",
        0.5, 7.18, 10, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

def km(slide, text):
    rect(slide, 0, 6.58, 13.33, 0.57, C_BG)
    rect(slide, 0, 6.60, 13.33, 0.44, RGBColor(0x04, 0x1A, 0x2C))
    rect(slide, 0, 6.60, 0.07,  0.44, C_ACCENT)
    txt(slide, text, 0.15, 6.62, 13.0, 0.42,
        size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

def bullet(slide, items, l, t, w, gap=0.62, size=13, dot_color=C_ACCENT):
    for i, item in enumerate(items):
        y = t + i * gap
        txt(slide, "•", l, y, 0.3, gap - 0.05,
            size=size, bold=True, color=dot_color)
        txt(slide, item, l + 0.3, y, w - 0.3, gap - 0.05,
            size=size, color=C_WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
circ = s.shapes.add_shape(9, Inches(8.5), Inches(-0.5), Inches(6.5), Inches(6.5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x52)
circ.line.fill.background()
rect(s, 0, 3.0, 0.12, 1.8, C_ACCENT)
txt(s, "Customer Tagging", 0.6, 1.7, 10, 1.1,
    size=48, bold=True, color=C_WHITE)
txt(s, "รู้ว่าลูกค้าเป็นใคร", 0.6, 2.95, 9, 0.7,
    size=26, color=C_ACCENT)
txt(s, "ทีม Data Science  |  เมษายน 2025", 0.6, 3.8, 9, 0.45,
    size=14, color=C_LIGHT)
footer(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHY: เสริม Merchant Enrichment
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
wh_tag(s, "WHY", C_YELLOW)

txt(s, "เรามี Merchant Enrichment อยู่แล้ว — แต่ยังมองลูกค้าไม่ครบ",
    0.5, 0.18, 11.2, 0.55, size=22, bold=True, color=C_WHITE)
rect(s, 0.5, 0.78, 12.3, 0.04, C_ACCENT)

# Left: Enrichment (เดิม)
rect(s, 0.5, 1.0, 5.8, 5.4, C_CARD)
rect(s, 0.5, 1.0, 5.8, 0.55, RGBColor(0x00, 0x2A, 0x3A))
txt(s, "Merchant Enrichment  (มีอยู่แล้ว)", 0.7, 1.07, 5.4, 0.4,
    size=13, bold=True, color=C_ACCENT)

txt(s, "บอกได้ว่า...", 0.7, 1.68, 5.4, 0.35, size=11, color=C_LIGHT)
enrich_items = [
    "ร้านนี้คือ Japanese Restaurant ระดับ Premium",
    "โรงแรมนี้คือ City Hotel / Business",
    "MCC นี้คือหมวด Healthcare",
    "ร้านนี้อยู่ใน Mall / Standalone",
]
for i, item in enumerate(enrich_items):
    y = 2.1 + i * 0.72
    rect(s, 0.7, y, 5.4, 0.6, RGBColor(0x0A, 0x28, 0x38))
    rect(s, 0.7, y, 0.06, 0.6, C_ACCENT)
    txt(s, "✓  " + item, 0.88, y + 0.12, 5.1, 0.36, size=11, color=C_WHITE)

rect(s, 0.7, 5.05, 5.4, 0.25, RGBColor(0x22, 0x12, 0x05))
txt(s, "แต่ไม่บอกว่า: ลูกค้าไปบ่อยแค่ไหน? กำลังขึ้นหรือลง?",
    0.85, 5.1, 5.1, 0.18, size=9, color=C_ORANGE)

# Arrow
txt(s, "+", 6.45, 3.3, 0.5, 0.6,
    size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Right: Customer Tag (เพิ่ม)
rect(s, 7.1, 1.0, 5.7, 5.4, C_CARD)
rect(s, 7.1, 1.0, 5.7, 0.55, RGBColor(0x00, 0x2A, 0x15))
txt(s, "Customer Tag  (Tag ใหม่ — ไม่มีในระบบเดิม)", 7.3, 1.07, 5.3, 0.4,
    size=13, bold=True, color=C_GREEN)

txt(s, "Tag ที่ยังไม่มีในระบบปัจจุบัน — ต้องสร้างใหม่ทั้งหมด", 7.3, 1.68, 5.3, 0.35, size=11, color=C_LIGHT)
tag_items = [
    ("Frequency",  "ไปบ่อยแค่ไหน — heavy / casual"),
    ("Recency",    "ยังอยู่ไหม — active / lapsed"),
    ("Trending",   "กำลังขึ้นหรือลง — rising / fading"),
    ("Spend",      "ใช้จ่ายหนักแค่ไหน — high / low intensity"),
]
for i, (label, desc) in enumerate(tag_items):
    y = 2.1 + i * 0.82
    rect(s, 7.3, y, 5.3, 0.7, RGBColor(0x0A, 0x28, 0x18))
    rect(s, 7.3, y, 0.06, 0.7, C_GREEN)
    txt(s, label, 7.5, y + 0.06, 1.5, 0.3, size=11, bold=True, color=C_GREEN)
    txt(s, desc,  7.5, y + 0.36, 5.0, 0.28, size=10, color=C_LIGHT)

rect(s, 7.3, 5.05, 5.3, 0.25, RGBColor(0x0A, 0x28, 0x18))
txt(s, "ปัจจุบันยังไม่มี tag เหล่านี้เลย — Customer Tag สร้างขึ้นมาใหม่ทั้งหมด",
    7.45, 5.1, 5.0, 0.18, size=9, color=C_GREEN)

footer(s, 2)
km(s, "ทีม DS จะทำ Customer Tag ขึ้นมาเพิ่ม ช่วยเสริมความสมบูรณ์ของ Merchant Enrichment เดิม")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT: Customer Tagging คืออะไร
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
wh_tag(s, "WHAT", C_ACCENT)

txt(s, "Customer Tagging คืออะไร?",
    0.5, 0.18, 11.2, 0.55, size=22, bold=True, color=C_WHITE)
rect(s, 0.5, 0.78, 12.3, 0.04, C_ACCENT)

# Definition box
rect(s, 0.5, 1.0, 12.3, 0.9, RGBColor(0x00, 0x28, 0x3A))
rect(s, 0.5, 1.0, 0.1, 0.9, C_ACCENT)
txt(s, "การอ่านพฤติกรรมจาก Transaction แล้วแปลงเป็น \"label\" ต่อลูกค้าแต่ละคน",
    0.75, 1.12, 12.0, 0.5, size=16, bold=True, color=C_WHITE)
txt(s, "→  1 ลูกค้า อาจมีหลาย label พร้อมกัน | refresh ทุกวัน ตามพฤติกรรมจริง",
    0.75, 1.62, 12.0, 0.24, size=11, color=C_ACCENT)

# Flow: Transaction → Tag Engine → Tags
flow = [
    (C_LIGHT,   "Transaction\nData",     "ทุกครั้งที่รูดบัตร\n12 เดือนย้อนหลัง"),
    (C_ACCENT,  "Tag\nEngine",           "คำนวณ score\nต่อ tag ต่อลูกค้า"),
    (C_GREEN,   "Customer\nTag Profile", "label + score\nพร้อมใช้งานทันที"),
]
for i, (col, name, desc) in enumerate(flow):
    x = 0.8 + i * 3.6
    rect(s, x, 2.15, 3.0, 1.8, C_CARD)
    rect(s, x, 2.15, 3.0, 0.55, col)
    txt(s, name, x + 0.15, 2.2, 2.7, 0.48,
        size=16, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, desc, x + 0.15, 2.82, 2.7, 0.95,
        size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < 2:
        txt(s, "→", x + 3.05, 2.82, 0.5, 0.5,
            size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Example profile
rect(s, 0.5, 4.18, 12.3, 0.38, RGBColor(0x00, 0x28, 0x3A))
txt(s, "ตัวอย่าง: ลูกค้า A", 0.65, 4.23, 3.0, 0.28,
    size=12, bold=True, color=C_ACCENT)

tags_ex = [
    ("frequent_diner",   "0.82", C_ACCENT),
    ("wealth_builder",   "0.61", C_GREEN),
    ("weekend_warrior",  "0.55", C_YELLOW),
    ("marketplace_shopper","0.48", C_ORANGE),
]
tag_w = 2.85
for i, (tag, score, col) in enumerate(tags_ex):
    x = 0.5 + i * (tag_w + 0.12)
    rect(s, x, 4.62, tag_w, 1.6, C_CARD)
    rect(s, x, 4.62, tag_w, 0.08, col)
    txt(s, tag,   x + 0.15, 4.74, tag_w - 0.3, 0.4, size=11, bold=True, color=col)
    txt(s, "score", x + 0.15, 5.18, 1.0, 0.28, size=9, color=C_LIGHT)
    txt(s, score,   x + 1.8,  5.12, 0.9, 0.4,  size=20, bold=True, color=col, align=PP_ALIGN.RIGHT)
    bar = float(score) * (tag_w - 0.3)
    rect(s, x + 0.15, 5.52, tag_w - 0.3, 0.14, RGBColor(0x0A, 0x18, 0x28))
    rect(s, x + 0.15, 5.52, bar, 0.14, col)

txt(s, "323 tags  |  13 มิติพฤติกรรม  |  ทุก tag มีนิยามชัดเจน ไม่ใช่การเดา",
    0.5, 6.3, 12.3, 0.3, size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

footer(s, 3)
km(s, "ทุก transaction ที่รูด คือข้อมูลที่บอกว่าลูกค้าเป็นใคร")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW: ทำงานยังไง
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
wh_tag(s, "HOW", C_ORANGE)

txt(s, "ระบบทำงานยังไง?",
    0.5, 0.18, 11.2, 0.55, size=22, bold=True, color=C_WHITE)
rect(s, 0.5, 0.78, 12.3, 0.04, C_ACCENT)

# 3 big boxes
boxes = [
    ("INPUT",   "Transaction Data",
     ["12 เดือนย้อนหลัง", "ทุกครั้งที่รูดบัตร", "ร้านไหน / เท่าไหร่ / เมื่อไหร่"],
     C_LIGHT),
    ("PROCESS", "Tag Engine",
     ["คำนวณ score ต่อ tag", "0.4×Freq + 0.4×Recency\n+ 0.2×Breadth", "score ≥ 0.3 → ได้ tag"],
     C_ACCENT),
    ("OUTPUT",  "Customer Tag Profile",
     ["tag + score + tier", "+ trending + recency", "SQL-ready ใช้งานทันที"],
     C_GREEN),
]
for i, (badge, name, points, col) in enumerate(boxes):
    x = 0.5 + i * 4.3
    rect(s, x, 1.1, 4.0, 5.2, C_CARD)
    rect(s, x, 1.1, 4.0, 0.75, col)
    txt(s, badge, x + 0.15, 1.15, 1.2, 0.6,
        size=11, bold=True, color=C_BG)
    txt(s, name, x + 0.15, 1.52, 3.7, 0.6,
        size=18, bold=True, color=C_BG, align=PP_ALIGN.LEFT)
    for j, pt in enumerate(points):
        y = 2.1 + j * 1.35
        rect(s, x + 0.2, y, 3.6, 1.18, RGBColor(0x0D, 0x22, 0x35))
        rect(s, x + 0.2, y, 0.07, 1.18, col)
        txt(s, pt, x + 0.4, y + 0.12, 3.3, 0.95, size=13, color=C_WHITE)
    if i < 2:
        txt(s, "→", x + 4.05, 3.3, 0.2, 0.6,
            size=26, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Bottom: refresh note
rect(s, 0.5, 6.28, 12.3, 0.25, RGBColor(0x00, 0x20, 0x30))
txt(s, "⟳  รันทุกคืน — profile ของลูกค้าทุกคนอัพเดทอัตโนมัติ ทันพฤติกรรมจริง",
    0.7, 6.31, 12.0, 0.2, size=10, color=C_ACCENT)

footer(s, 4)
km(s, "ระบบอ่าน transaction ทุกคืน → profile อัพเดทอัตโนมัติทุกวัน")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHAT YOU GET: ได้อะไร
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
wh_tag(s, "WHAT YOU GET", C_GREEN)

txt(s, "Output ที่ได้ — หน้าตาเป็นยังไง?",
    0.5, 0.18, 11.2, 0.55, size=22, bold=True, color=C_WHITE)
rect(s, 0.5, 0.78, 12.3, 0.04, C_ACCENT)

# Left: profile card
rect(s, 0.5, 1.05, 6.0, 5.35, C_CARD)
rect(s, 0.5, 1.05, 6.0, 0.5, RGBColor(0x00, 0x30, 0x45))
txt(s, "ลูกค้า A — Tag Profile", 0.7, 1.12, 5.6, 0.35,
    size=13, bold=True, color=C_ACCENT)

profile_rows = [
    ("frequent_diner",  "0.74", "heavy",  "rising",  "active",  C_ACCENT),
    ("wealth_builder",  "0.61", "casual", "stable",  "active",  C_GREEN),
    ("weekend_warrior", "0.55", "casual", "fading",  "lapsed",  C_YELLOW),
]
col_labels = ["tag", "score", "tier", "trending", "recency"]
col_x      = [0.65, 2.2, 3.15, 4.0, 5.0]
col_w      = [1.5,  0.9, 0.8,  0.9, 0.9]

# Header row
rect(s, 0.5, 1.6, 6.0, 0.35, RGBColor(0x06, 0x20, 0x30))
for j, (label, cx, cw) in enumerate(zip(col_labels, col_x, col_w)):
    txt(s, label, cx, 1.65, cw, 0.25, size=9, bold=True, color=C_LIGHT)

for i, (tag, score, tier, trend, recency, col) in enumerate(profile_rows):
    y = 2.0 + i * 1.28
    row_bg = RGBColor(0x12, 0x24, 0x38) if i % 2 == 0 else C_CARD
    rect(s, 0.5, y, 6.0, 1.18, row_bg)
    vals = [tag, score, tier, trend, recency]
    for j, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
        c = col if j == 0 else (C_YELLOW if j in [2,3,4] else C_WHITE)
        txt(s, val, cx, y + 0.38, cw, 0.38, size=10, bold=(j==0), color=c)
    bar = float(score) * 1.8
    rect(s, col_x[1], y + 0.72, 1.8, 0.14, RGBColor(0x0A, 0x18, 0x28))
    rect(s, col_x[1], y + 0.72, bar, 0.14, col)

# Right: what you can do
txt(s, "สิ่งที่ทำได้ทันที", 7.0, 1.05, 5.8, 0.42, size=15, bold=True, color=C_WHITE)
queries = [
    (C_ACCENT,  "Heavy diner ที่ยังอยู่",
                "tier = 'heavy'  AND  recency = 'active'"),
    (C_GREEN,   "กลุ่มที่กำลัง active ขึ้น",
                "trending = 'rising'  AND  tier = 'heavy'"),
    (C_YELLOW,  "Win-back — เคยบ่อยแต่หายไป",
                "trending = 'fading'  AND  recency = 'lapsed'"),
    (C_ORANGE,  "พฤติกรรมใหม่ — เพิ่งเริ่ม",
                "tag_age = 'new'  AND  tier = 'casual'"),
]
for i, (col, label, query) in enumerate(queries):
    y = 1.55 + i * 1.22
    rect(s, 7.0, y, 6.0, 1.1, C_CARD)
    rect(s, 7.0, y, 0.08, 1.1, col)
    txt(s, label, 7.2, y + 0.08, 5.7, 0.35, size=12, bold=True, color=col)
    txt(s, query, 7.2, y + 0.5,  5.7, 0.52, size=10, color=C_LIGHT)

footer(s, 5)
km(s, "1 ลูกค้า = 1 profile — บอกว่าเขาเป็นใคร พร้อมใช้ทันที")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW TO USE: นำไปใช้อย่างไร
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
wh_tag(s, "HOW TO USE", C_ACCENT)

txt(s, "Campaign Manager นำไปใช้งานยังไง?",
    0.5, 0.18, 11.2, 0.55, size=22, bold=True, color=C_WHITE)
rect(s, 0.5, 0.78, 12.3, 0.04, C_ACCENT)

# Before / After
for i, (side, col, bg_col, title, items) in enumerate([
    ("BEFORE", C_ORANGE, RGBColor(0x20, 0x10, 0x05),
     "วิธีเดิม — Mass Campaign",
     ["แบ่งลูกค้าด้วย spend amount",
      "ส่ง offer เดียวกันทั้งกลุ่ม",
      "ไม่รู้ว่าใคร active หรือ fading",
      "วัดผลได้แค่ response rate รวม"]),
    ("AFTER", C_GREEN, RGBColor(0x05, 0x20, 0x10),
     "วิธีใหม่ — Tag-based Targeting",
     ["เลือก tag ที่ต้องการ (เช่น frequent_diner)",
      "กรองด้วย dimension (heavy + active + rising)",
      "ส่ง offer ที่ตรงกับ behavior จริง",
      "วัดผลรายกลุ่ม — ปรับได้แม่นขึ้น"]),
]):
    x = 0.5 + i * 6.5
    rect(s, x, 1.05, 6.1, 4.7, bg_col)
    rect(s, x, 1.05, 6.1, 0.52, col)
    txt(s, side, x + 0.18, 1.1, 1.2, 0.38, size=13, bold=True, color=C_BG)
    txt(s, title, x + 1.5, 1.12, 4.5, 0.38, size=12, bold=True, color=C_BG)
    for j, item in enumerate(items):
        y = 1.72 + j * 0.9
        rect(s, x + 0.2, y, 5.7, 0.78, RGBColor(0x0D, 0x22, 0x18) if i else RGBColor(0x22, 0x10, 0x05))
        rect(s, x + 0.2, y, 0.07, 0.78, col)
        txt(s, item, x + 0.4, y + 0.14, 5.3, 0.5, size=12, color=C_WHITE)

# Arrow between
txt(s, "→", 6.55, 3.1, 0.5, 0.6,
    size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Example
rect(s, 0.5, 5.92, 12.3, 0.58, RGBColor(0x00, 0x28, 0x3A))
rect(s, 0.5, 5.92, 0.1, 0.58, C_ACCENT)
txt(s, "ตัวอย่าง:", 0.72, 5.98, 1.4, 0.42, size=11, bold=True, color=C_ACCENT)
txt(s, "frequent_diner  +  tier=heavy  +  trending=rising  →  ส่ง dining offer ก่อน 11.11",
    2.2, 5.98, 10.4, 0.42, size=12, color=C_WHITE)

footer(s, 6)
km(s, "tag ที่ใช่ → offer ที่ใช่ → คนที่ใช่ → เวลาที่เหมาะ")

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  CustomerTagging_Business.pptx — {len(prs.slides)} slides")
print(f"   1 Title  |  WHY · WHAT · HOW · WHAT YOU GET · HOW TO USE")
