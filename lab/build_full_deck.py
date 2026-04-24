"""
build_full_deck.py
รวม MemoV2 (13 slides) + Layer1 Detail (4 slides) + Tag Appendix
→ CustomerTagging_FullDeck.pptx
"""
import subprocess
import copy
from pptx import Presentation
from pptx.util import Inches

BASE = "/Users/adisornj/Desktop/Thena/lab"
OUT  = f"{BASE}/CustomerTagging_FullDeck.pptx"

# ── Step 1: Regenerate source files ───────────────────────────────────────────
print("Regenerating source files...")
subprocess.run(["python3", f"{BASE}/build_memo_pptx.py"],      check=True, capture_output=True)
subprocess.run(["python3", f"{BASE}/build_layer1_detail.py"],  check=True, capture_output=True)
print("  ✓ build_memo_pptx.py")
print("  ✓ build_layer1_detail.py")

# ── Step 2: Merge slides ───────────────────────────────────────────────────────
def copy_slides(src_path, dst_prs):
    src = Presentation(src_path)
    for slide in src.slides:
        blank    = dst_prs.slide_layouts[6]
        new_slide = dst_prs.slides.add_slide(blank)
        for shape in slide.shapes:
            new_slide.shapes._spTree.append(copy.deepcopy(shape.element))
    return len(src.slides)

print("\nMerging slides...")
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

n1 = copy_slides(f"{BASE}/CustomerTagging_MemoV2_Presentation.pptx", prs)
print(f"  ✓ MemoV2          — {n1} slides")

n2 = copy_slides(f"{BASE}/CustomerTagging_Layer1_Detail.pptx", prs)
print(f"  ✓ Layer1 Detail   — {n2} slides")

prs.save(OUT)
print(f"  → Saved: {n1 + n2} slides")

# ── Step 3: Append Tag Library ────────────────────────────────────────────────
print("\nAppending Tag Library...")
result = subprocess.run(
    ["python3", f"{BASE}/append_tags_pptx.py", OUT],
    check=True, capture_output=True, text=True
)
print(f"  ✓ {result.stdout.strip()}")

# ── Done ──────────────────────────────────────────────────────────────────────
final = Presentation(OUT)
print(f"\n✅  CustomerTagging_FullDeck.pptx — {len(final.slides)} slides total")
print(f"   {n1} MemoV2  +  {n2} Layer1  +  {len(final.slides) - n1 - n2} Appendix")
