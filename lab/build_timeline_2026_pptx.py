"""
build_timeline_2026_pptx.py
One-page 2026 Annual Timeline — light/white theme
12 placeholder tasks across 12 months
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

OUT = "/Users/adisornj/Desktop/Thena/lab/Timeline_2026.pptx"

# ── Palette (light theme) ──────────────────────────────────────────────────────
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white background
C_HDR_BG  = RGBColor(0xF0, 0xF4, 0xF8)  # light blue-gray header
C_BORDER  = RGBColor(0xE8, 0xED, 0xF2)  # very faint border
C_TEXT    = RGBColor(0x1A, 0x2A, 0x3A)  # dark navy text
C_SUBTEXT = RGBColor(0x6A, 0x7A, 0x8A)  # muted gray text
C_ACCENT  = RGBColor(0x00, 0x7A, 0xCC)  # blue accent
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# Task colors (pastel, readable on white)
TASK_COLORS = [
    RGBColor(0x4A, 0xA5, 0xFF),  # blue
    RGBColor(0x00, 0xC5, 0x8A),  # teal
    RGBColor(0xFF, 0x8C, 0x42),  # orange
    RGBColor(0x8A, 0x7A, 0xFF),  # purple
    RGBColor(0xFF, 0x6B, 0x6B),  # red
    RGBColor(0x00, 0xC2, 0xFF),  # cyan
    RGBColor(0xFF, 0xC4, 0x3A),  # yellow
    RGBColor(0x5E, 0xD0, 0x6E),  # green
    RGBColor(0xFF, 0x7E, 0xC4),  # pink
    RGBColor(0x7A, 0xB0, 0xCC),  # steel blue
    RGBColor(0xAA, 0x8A, 0xFF),  # lavender
    RGBColor(0xFF, 0xA0, 0x60),  # peach
]

MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
          "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]

# ── Tasks: (label, month_start, month_end, color_index) ───────────────────────
# month_start/end: 1=Jan ... 12=Dec
TASKS = [
    ("Task 01 — Placeholder",   1,  2,  0),
    ("Task 02 — Placeholder",   1,  3,  1),
    ("Task 03 — Placeholder",   2,  4,  2),
    ("Task 04 — Placeholder",   3,  5,  3),
    ("Task 05 — Placeholder",   4,  6,  4),
    ("Task 06 — Placeholder",   5,  7,  5),
    ("Task 07 — Placeholder",   6,  8,  6),
    ("Task 08 — Placeholder",   7,  9,  7),
    ("Task 09 — Placeholder",   8, 10,  8),
    ("Task 10 — Placeholder",   9, 11,  9),
    ("Task 11 — Placeholder",  10, 12, 10),
    ("Task 12 — Placeholder",  11, 12, 11),
]

QUARTERS = [
    ("Q1", 1,  3,  RGBColor(0xE8, 0xF0, 0xFF)),
    ("Q2", 4,  6,  RGBColor(0xE8, 0xFB, 0xF4)),
    ("Q3", 7,  9,  RGBColor(0xFF, 0xF5, 0xE8)),
    ("Q4", 10, 12, RGBColor(0xF5, 0xE8, 0xFF)),
]

# ── Layout constants ───────────────────────────────────────────────────────────
SL_W = Inches(13.33)
SL_H = Inches(7.5)

MARGIN_L = Inches(0.4)
MARGIN_T = Inches(0.4)
MARGIN_R = Inches(0.4)

HDR_H    = Inches(0.65)   # title bar
LEGEND_H = Inches(0.45)   # month header row
TASK_H   = Inches(0.42)   # each task row height
TASK_GAP = Inches(0.06)   # gap between rows
LABEL_W  = Inches(2.2)    # left label column

CHART_L = MARGIN_L + LABEL_W + Inches(0.1)
CHART_W = SL_W - CHART_L - MARGIN_R
MONTH_W = CHART_W / 12

# vertical start of gantt area
GANTT_T = MARGIN_T + HDR_H + LEGEND_H + Inches(0.08)


# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def qn(tag):
    from pptx.oxml.ns import qn as _qn
    return _qn(tag)

def add_rect(slide, x, y, w, h, fill: RGBColor, radius=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape

def add_text(slide, x, y, w, h, text, size=9, bold=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    color = color or C_TEXT
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.name  = "Helvetica Neue"
    return txb

def cell_fill(cell, color: RGBColor):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in list(tcPr):
        tag = el.tag
        if any(tag.endswith(s) for s in ["}solidFill","}noFill", "}gradFill"]):
            tcPr.remove(el)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", rgb_hex(color))


# ── Build slide ───────────────────────────────────────────────────────────────
prs  = Presentation()
prs.slide_width  = SL_W
prs.slide_height = SL_H

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# ── Background ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, SL_W, SL_H, C_BG)

# ── Header bar ────────────────────────────────────────────────────────────────
add_rect(slide, MARGIN_L, MARGIN_T,
         SL_W - MARGIN_L - MARGIN_R, HDR_H,
         C_HDR_BG)

# Accent left strip
add_rect(slide, MARGIN_L, MARGIN_T, Inches(0.06), HDR_H, C_ACCENT)

add_text(slide,
         MARGIN_L + Inches(0.18), MARGIN_T + Inches(0.1),
         Inches(6), HDR_H,
         "2026 Annual Timeline",
         size=18, bold=True, color=C_TEXT)

add_text(slide,
         MARGIN_L + Inches(0.18), MARGIN_T + Inches(0.38),
         Inches(6), HDR_H,
         "Project Overview — Placeholder",
         size=9, bold=False, color=C_SUBTEXT)

# Stats (top right)
stats = [("12", "Tasks"), ("4", "Quarters"), ("12", "Months")]
sx = SL_W - MARGIN_R - Inches(3.5)
for val, lbl in stats:
    add_text(slide, sx, MARGIN_T + Inches(0.04),
             Inches(1.0), Inches(0.3),
             val, size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, sx, MARGIN_T + Inches(0.34),
             Inches(1.0), Inches(0.2),
             lbl, size=7, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
    sx += Inches(1.15)

# ── Quarter bands (background behind chart) ──────────────────────────────────
for q_label, m_start, m_end, q_color in QUARTERS:
    qx = CHART_L + (m_start - 1) * MONTH_W
    qw = (m_end - m_start + 1) * MONTH_W
    qy = MARGIN_T + HDR_H + Inches(0.04)
    qh = LEGEND_H + Inches(0.04) + len(TASKS) * (TASK_H + TASK_GAP)
    add_rect(slide, qx, qy, qw, qh, q_color)

# ── Month header row ─────────────────────────────────────────────────────────
legend_y = MARGIN_T + HDR_H + Inches(0.04)

# "Tasks" label header
add_rect(slide,
         MARGIN_L, legend_y,
         LABEL_W + Inches(0.08), LEGEND_H,
         C_HDR_BG)
add_text(slide,
         MARGIN_L + Inches(0.12), legend_y + Inches(0.08),
         LABEL_W, LEGEND_H,
         "Task", size=8, bold=True, color=C_SUBTEXT)

for i, month in enumerate(MONTHS):
    mx = CHART_L + i * MONTH_W
    # month label
    add_text(slide, mx, legend_y + Inches(0.08),
             MONTH_W, LEGEND_H,
             month, size=8, bold=True,
             color=C_TEXT, align=PP_ALIGN.CENTER)

# Bottom divider under month headers
add_rect(slide, MARGIN_L, legend_y + LEGEND_H - Inches(0.01),
         SL_W - MARGIN_L - MARGIN_R, Inches(0.006), C_BORDER)

# ── Task rows ────────────────────────────────────────────────────────────────
for row_idx, (label, m_start, m_end, color_idx) in enumerate(TASKS):
    row_y = GANTT_T + row_idx * (TASK_H + TASK_GAP)
    bar_color = TASK_COLORS[color_idx]

    # alternating row bg
    row_bg = RGBColor(0xF8, 0xFA, 0xFC) if row_idx % 2 == 0 else C_BG
    add_rect(slide, MARGIN_L, row_y,
             SL_W - MARGIN_L - MARGIN_R, TASK_H, row_bg)

    # task label
    add_text(slide,
             MARGIN_L + Inches(0.1), row_y + Inches(0.06),
             LABEL_W - Inches(0.15), TASK_H,
             label, size=8, bold=False, color=C_TEXT)

    # Gantt bar
    bar_x = CHART_L + (m_start - 1) * MONTH_W + Inches(0.04)
    bar_w = (m_end - m_start + 1) * MONTH_W - Inches(0.08)
    bar_y = row_y + Inches(0.09)
    bar_h = TASK_H - Inches(0.18)

    bar_shape = slide.shapes.add_shape(
        1, bar_x, bar_y, bar_w, bar_h)
    bar_shape.line.fill.background()
    bar_shape.line.width = 0
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = bar_color

    # duration label on bar
    n_months = m_end - m_start + 1
    dur_txt  = f"{n_months}M"
    add_text(slide,
             bar_x + Inches(0.05), bar_y,
             bar_w - Inches(0.1), bar_h,
             dur_txt, size=7, bold=True,
             color=C_WHITE, align=PP_ALIGN.LEFT)


# ── Quarter labels (above month headers) ─────────────────────────────────────
for q_label, m_start, m_end, _ in QUARTERS:
    qx = CHART_L + (m_start - 1) * MONTH_W
    qw = (m_end - m_start + 1) * MONTH_W
    add_text(slide,
             qx, legend_y - Inches(0.02),
             qw, Inches(0.22),
             q_label, size=7, bold=True,
             color=C_SUBTEXT, align=PP_ALIGN.CENTER)

# ── Legend (task colors) ──────────────────────────────────────────────────────
legend_bottom_y = GANTT_T + len(TASKS) * (TASK_H + TASK_GAP) + Inches(0.12)

add_text(slide,
         MARGIN_L, legend_bottom_y,
         Inches(1.0), Inches(0.25),
         "Legend", size=7, bold=True, color=C_SUBTEXT)

lx = MARGIN_L + Inches(0.8)
for i, (label, _, _, color_idx) in enumerate(TASKS):
    if lx + Inches(1.8) > SL_W - MARGIN_R:
        break
    add_rect(slide, lx, legend_bottom_y + Inches(0.05),
             Inches(0.18), Inches(0.15), TASK_COLORS[color_idx])
    short = label.split("—")[0].strip()
    add_text(slide,
             lx + Inches(0.22), legend_bottom_y + Inches(0.02),
             Inches(1.4), Inches(0.22),
             short, size=6.5, color=C_SUBTEXT)
    lx += Inches(1.65)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  Timeline_2026.pptx — 1 slide")
print(f"   12 tasks | 12 months | light theme")
