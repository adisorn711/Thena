"""
build_tag_library_pptx.py
Customer Tag Library — PPTX version
Cover · TOC · 1 slide/category (split ถ้า tags เยอะ)
Color-coded rows: Existing / Overlap / New  +  S/V/M stability
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = "/Users/adisornj/Desktop/Thena/lab/customer_tag_library_v2_expansion.md"
OUT = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_TagLibrary.pptx"

# ── Palette ────────────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_HDR    = RGBColor(0x06, 0x0F, 0x18)
C_DARK   = RGBColor(0x0A, 0x14, 0x20)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_LGRAY  = RGBColor(0x6A, 0x8A, 0x9A)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
C_PURPLE = RGBColor(0x8A, 0x7A, 0xFF)
C_STEEL  = RGBColor(0x7A, 0xB0, 0xCC)

ROW_EVEN = {"Existing": RGBColor(0x0D,0x1B,0x2A), "Overlap": RGBColor(0x00,0x1C,0x2A), "New": RGBColor(0x07,0x1A,0x10)}
ROW_ODD  = {"Existing": RGBColor(0x11,0x22,0x35), "Overlap": RGBColor(0x00,0x24,0x38), "New": RGBColor(0x09,0x22,0x16)}
TAG_CLR  = {"Existing": C_STEEL,  "Overlap": C_ACCENT, "New": C_GREEN}
SRC_CLR  = {"Existing": C_STEEL,  "Overlap": C_ACCENT, "New": C_GREEN}
SRC_SHORT = {"Existing": "Ex", "Overlap": "Ov", "New": "New"}
STAB_CLR = {"S": C_ACCENT, "V": C_YELLOW, "M": C_ORANGE}
PATH_CLR = {"A": C_PURPLE, "B": C_ORANGE}

MAX_ROWS = 24   # data rows per slide

# ── Parse markdown ─────────────────────────────────────────────────────────────
categories, cur = [], None

with open(SRC, encoding="utf-8") as f:
    for line in f:
        line = line.rstrip("\n")
        m = re.match(r"###\s+(.+)", line)
        if m:
            head = m.group(1).strip()
            em = re.match(r"^([\U0001F000-\U0010FFFF\u2600-\u27FF\u2B00-\u2BFF]+)\s*", head)
            emoji = em.group(1) if em else ""
            name  = head[len(em.group(0)):].strip() if em else head
            cur   = {"name": name, "emoji": emoji, "tags": []}
            categories.append(cur)
            continue
        if not cur:
            continue
        m = re.match(
            r"\|\s*`([^`]+)`\s*\|"
            r"\s*([^|]+)\|\s*([^|]+)\|"
            r"\s*([AB])\s*\|\s*(Existing|Overlap|New)\s*\|"
            r"\s*([^|]*)\|\s*([SVM?])\s*\|", line)
        if m:
            cur["tags"].append({
                "tag": m.group(1).strip(), "defn": m.group(2).strip(),
                "signal": m.group(3).strip(), "path": m.group(4).strip(),
                "source": m.group(5).strip(), "stab": m.group(7).strip(),
            })

all_tags   = [t for c in categories for t in c["tags"]]
total_tags = len(all_tags)
print(f"Parsed: {len(categories)} categories, {total_tags} tags")

# ── pptx helpers ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=11, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb

def multirun(slide, runs, l, t, w, h, align=PP_ALIGN.LEFT):
    """Textbox with multiple colored runs."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for text, size, bold, color in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def page_footer(slide, n):
    rect(slide, 0, 7.14, 13.33, 0.36, C_DARK)
    txt(slide, "Data Science Team  |  Customer Tag Library v2.2  |  เมษายน 2025",
        0.3, 7.17, 11, 0.28, size=8, color=C_LGRAY)
    txt(slide, str(n), 12.8, 7.17, 0.4, 0.28, size=8,
        color=C_LGRAY, align=PP_ALIGN.RIGHT)

def cell_fill(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for el in list(tcPr):
        tag = el.tag
        if tag.endswith("}solidFill") or tag.endswith("}noFill") or tag.endswith("}gradFill"):
            tcPr.remove(el)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf,   qn("a:srgbClr"))
    sc.set("val", str(color))

def cell_txt(cell, text, size=9, bold=False, color=C_LIGHT,
             font="Helvetica Neue", align=PP_ALIGN.LEFT):
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for r in p._p.findall(qn("a:r")):
        p._p.remove(r)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font

# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, C_BG)
circ = s.shapes.add_shape(9, Inches(8.2), Inches(-0.8), Inches(7.0), Inches(7.0))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x18, 0x28)
circ.line.fill.background()
rect(s, 0, 2.7, 0.12, 2.1, C_ACCENT)

txt(s, "Customer Tag Library", 0.5, 1.2, 11, 1.2, size=46, bold=True)
txt(s, "พฤติกรรมลูกค้า 288 tags · 18 หมวดหมู่", 0.5, 2.65, 10, 0.65, size=19, color=C_ACCENT)
txt(s, "ทีม Data Science  |  เมษายน 2025", 0.5, 3.45, 8, 0.5, size=14, color=C_LIGHT)

# Stat boxes row 1: Source
ex_n = sum(1 for t in all_tags if t["source"] == "Existing")
ov_n = sum(1 for t in all_tags if t["source"] == "Overlap")
nw_n = sum(1 for t in all_tags if t["source"] == "New")
s_n  = sum(1 for t in all_tags if t["stab"] == "S")
v_n  = sum(1 for t in all_tags if t["stab"] == "V")
m_n  = sum(1 for t in all_tags if t["stab"] == "M")
pa_n = sum(1 for t in all_tags if t["path"] == "A")
pb_n = sum(1 for t in all_tags if t["path"] == "B")

stat_rows = [
    [("SOURCE", C_LGRAY, False), (str(ex_n), C_STEEL, True), ("Existing", C_STEEL, False),
     (str(ov_n), C_ACCENT, True), ("Overlap", C_ACCENT, False),
     (str(nw_n), C_GREEN, True),  ("New", C_GREEN, False)],
    [("STABILITY", C_LGRAY, False), (str(s_n), C_ACCENT, True), ("Stable", C_ACCENT, False),
     (str(v_n), C_YELLOW, True), ("Variable", C_YELLOW, False),
     (str(m_n), C_ORANGE, True), ("Momentum", C_ORANGE, False)],
    [("PATH", C_LGRAY, False), (str(pa_n), C_PURPLE, True), ("Path A (no enrichment)", C_PURPLE, False),
     (str(pb_n), C_ORANGE, True), ("Path B (enrichment required)", C_ORANGE, False)],
]
for row_i, row in enumerate(stat_rows):
    y = 4.5 + row_i * 0.75
    rect(s, 0.5, y, 7.5, 0.62, RGBColor(0x0A, 0x18, 0x28))
    runs = []
    for i, (label, col, bold) in enumerate(row):
        sep = "   ·   " if i > 0 and i % 2 == 0 else ("   " if i > 0 else "")
        runs.append((sep + label + "  ", 11, bold, col))
    multirun(s, runs, 0.65, y + 0.12, 7.3, 0.42, align=PP_ALIGN.LEFT)

page_footer(s, 1)

# ── SLIDE 2: Table of Contents ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, C_BG)
rect(s, 0, 0, 13.33, 0.65, C_HDR)
rect(s, 0, 0, 0.08,  0.65, C_ACCENT)
txt(s, "หมวดหมู่พฤติกรรมทั้งหมด", 0.2, 0.08, 10, 0.5, size=18, bold=True)
txt(s, f"{total_tags} tags  ·  {len(categories)} categories",
    10.5, 0.12, 2.6, 0.42, size=13, color=C_ACCENT, align=PP_ALIGN.RIGHT)

COL_X = [0.25, 6.8]  # two columns
for i, cat in enumerate(categories):
    col_i = i // 9
    row_i = i % 9
    x = COL_X[col_i]
    y = 0.82 + row_i * 0.66

    s_cnt = sum(1 for t in cat["tags"] if t["source"] == "Existing")
    o_cnt = sum(1 for t in cat["tags"] if t["source"] == "Overlap")
    n_cnt = sum(1 for t in cat["tags"] if t["source"] == "New")
    total_n = len(cat["tags"])

    rect(s, x, y, 6.2, 0.58, RGBColor(0x0A, 0x18, 0x28))
    rect(s, x, y, 0.06, 0.58, C_ACCENT)
    txt(s, f"{cat['emoji']}  {cat['name']}", x+0.15, y+0.06, 3.8, 0.36, size=11, bold=True)
    txt(s, f"{total_n}", x+4.0, y+0.06, 0.7, 0.36, size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.RIGHT)
    txt(s, "tags", x+4.7, y+0.1, 0.8, 0.28, size=9, color=C_LGRAY)

    # Mini source counts
    multirun(s, [
        (f"{s_cnt} ex", 8, False, C_STEEL),
        (f"  ·  {o_cnt} ov", 8, False, C_ACCENT),
        (f"  ·  {n_cnt} new", 8, False, C_GREEN),
    ], x+0.15, y+0.38, 4.0, 0.2)

page_footer(s, 2)

# ── Category Slides ────────────────────────────────────────────────────────────
COL_W = [Inches(w) for w in [2.4, 5.6, 2.2, 0.63, 0.6, 0.8]]
# Tag | Definition | Signal | Path | Src | S/V/M  → sum = 12.23", left=0.3", total=12.53"+margins
TBL_L = Inches(0.3)
TBL_W = sum(COL_W)
TBL_T = Inches(0.90)
TBL_H = Inches(7.12) - TBL_T  # 6.22"
HDR_LABELS = ["Tag", "Definition", "Signal", "Path", "Src", "S/V/M"]
CTR = {3, 4, 5}   # center-aligned columns by index

page_n = 3

for cat in categories:
    tags = cat["tags"]
    if not tags:
        continue

    s_cnt = sum(1 for t in tags if t["source"] == "Existing")
    o_cnt = sum(1 for t in tags if t["source"] == "Overlap")
    n_cnt = sum(1 for t in tags if t["source"] == "New")

    chunks = [tags[i:i+MAX_ROWS] for i in range(0, len(tags), MAX_ROWS)]

    for chunk_idx, chunk in enumerate(chunks):
        sl = prs.slides.add_slide(BLANK)
        rect(sl, 0, 0, 13.33, 7.5, C_BG)

        # Header
        rect(sl, 0, 0, 13.33, 0.65, C_HDR)
        rect(sl, 0, 0, 0.08,  0.65, C_ACCENT)

        label = f"{cat['emoji']}  {cat['name']}"
        if len(chunks) > 1:
            label += f"  ({chunk_idx+1}/{len(chunks)})"
        txt(sl, label, 0.2, 0.07, 8.2, 0.52, size=17, bold=True)

        multirun(sl, [
            (f"{len(tags)} tags    ", 10, True,  C_ACCENT),
            (f"{s_cnt} ex   ",        10, False, C_STEEL),
            (f"{o_cnt} overlap   ",   10, False, C_ACCENT),
            (f"{n_cnt} new",          10, False, C_GREEN),
        ], 8.5, 0.12, 4.6, 0.42, align=PP_ALIGN.RIGHT)

        # Legend strip
        rect(sl, 0, 0.66, 13.33, 0.22, RGBColor(0x0A, 0x14, 0x20))
        multirun(sl, [
            ("■ Existing  ", 7.5, False, C_STEEL),
            ("■ Overlap  ",  7.5, False, C_ACCENT),
            ("■ New        ", 7.5, False, C_GREEN),
            ("|   S = Stable  ", 7.5, False, C_ACCENT),
            ("V = Variable  ", 7.5, False, C_YELLOW),
            ("M = Momentum      ", 7.5, False, C_ORANGE),
            ("|   Path A = ไม่ต้องรอ enrichment  ", 7.5, False, C_PURPLE),
            ("Path B = ต้องรอ enrichment", 7.5, False, C_ORANGE),
        ], 0.2, 0.68, 12.9, 0.18)

        # Table
        n_rows = len(chunk) + 1
        tbl_shape = sl.shapes.add_table(n_rows, len(COL_W), TBL_L, TBL_T, TBL_W, TBL_H)
        tbl = tbl_shape.table

        for j, w in enumerate(COL_W):
            tbl.columns[j].width = w

        hdr_h    = Inches(0.3)
        data_h   = int((TBL_H - hdr_h) / len(chunk))
        tbl.rows[0].height = hdr_h
        for i in range(1, n_rows):
            tbl.rows[i].height = data_h

        # Header row
        for j, lbl in enumerate(HDR_LABELS):
            c = tbl.cell(0, j)
            cell_fill(c, C_HDR)
            cell_txt(c, lbl, size=8, bold=True, color=C_LGRAY,
                     align=PP_ALIGN.CENTER if j in CTR else PP_ALIGN.LEFT)

        # Data rows
        for i, tag in enumerate(chunk):
            src  = tag["source"]
            stab = tag["stab"]
            path = tag["path"]
            bg   = ROW_EVEN[src] if i % 2 == 0 else ROW_ODD[src]

            row_data = [
                (tag["tag"],            TAG_CLR[src],              "Courier New",     True,  PP_ALIGN.LEFT),
                (tag["defn"],           C_LIGHT,                   "Helvetica Neue",  False, PP_ALIGN.LEFT),
                (tag["signal"],         C_LGRAY,                   "Helvetica Neue",  False, PP_ALIGN.LEFT),
                (path,                  PATH_CLR[path],            "Helvetica Neue",  True,  PP_ALIGN.CENTER),
                (SRC_SHORT[src],        SRC_CLR[src],              "Helvetica Neue",  True,  PP_ALIGN.CENTER),
                (stab,                  STAB_CLR.get(stab, C_LIGHT), "Helvetica Neue", True, PP_ALIGN.CENTER),
            ]

            for j, (text, color, font, bold, align) in enumerate(row_data):
                c = tbl.cell(i + 1, j)
                cell_fill(c, bg)
                cell_txt(c, text, size=8.5, bold=bold, color=color, font=font, align=align)

        page_footer(sl, page_n)
        page_n += 1

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\n✅  CustomerTagging_TagLibrary.pptx")
print(f"   {page_n-1} slides  ({len(categories)} categories, {total_tags} tags)")
print(f"   → {OUT}")
