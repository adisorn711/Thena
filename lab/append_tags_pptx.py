from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import sys
_target = sys.argv[1] if len(sys.argv) > 1 else "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_MemoV2_Presentation.pptx"
prs = Presentation(_target)
BLANK = prs.slide_layouts[6]

C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
C_GRAY   = RGBColor(0x55, 0x70, 0x80)

def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=10, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.7, RGBColor(0x0A, 0x14, 0x20))
    rect(slide, 0, 0.7, 13.33, 0.04, C_ACCENT)
    txt(slide, title, 0.4, 0.08, 11, 0.55, size=20, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.52, 11, 0.22, size=9, color=C_ACCENT)

def footer(slide, n, total):
    rect(slide, 0, 7.18, 13.33, 0.32, RGBColor(0x0A, 0x14, 0x20))
    txt(slide, "Customer Tag Library v2.0 — 303 Tags | Data Science Team | เมษายน 2025",
        0.4, 7.2, 11.5, 0.25, size=8, color=C_GRAY)
    txt(slide, f"Appendix {n}", 12.0, 7.2, 1.2, 0.25, size=8,
        color=C_GRAY, align=PP_ALIGN.RIGHT)

def path_color(p):
    return C_ACCENT if p == "A" else C_ORANGE

# ── Source classification ──────────────────────────────────────────────────────
# O = Overlap (ตรงกับ enrichment tag ของทีม)
# N = New     (เพิ่มใหม่ใน v2.1 จาก enrichment review)
# E = Existing (default)

OVERLAP_TAGS = {
    # Food & Dining
    "frequent_diner","fast_food_regular","fine_dining_enthusiast","cafe_hopper",
    "bakery_lover","dessert_lover","japanese_food_lover","korean_food_lover",
    "chinese_food_lover","thai_food_loyalist","western_food_lover","italian_food_lover",
    "seafood_lover","food_court_regular","buffet_lover","family_restaurant_goer","healthy_eater",
    # Travel
    "frequent_traveler","luxury_traveler",
    # Shopping
    "luxury_shopper","fashion_lover","sneaker_enthusiast","beauty_conscious",
    "home_decor_lover","sports_gear_buyer","convenience_store_addict","marketplace_shopper",
    "subscription_stacker","premium_brand_loyal","accessories_collector","bargain_hunter",
    "small_business_spender",
    # Entertainment
    "movie_lover","concert_goer","night_owl_entertainer","karaoke_regular","spa_lover",
    "music_festival_goer","streaming_subscriber","outdoor_adventurer","theme_park_visitor",
    "nail_salon_regular","hair_salon_loyal","beauty_clinic_visitor","aesthetic_spender",
    # Transport
    "car_owner",
    # Health
    "health_conscious","mental_wellness_seeker",
    # Life Stage
    "parent","homeowner","home_renovator","young_professional",
    # Education
    "lifelong_learner","professional_upskiller",
    # Real Estate
    "diy_homeowner",
}

NEW_TAGS = {
    "bbq_lover","hot_pot_lover","noodle_lover","steak_lover","french_food_lover",
    "american_food_lover","vietnamese_food_lover","thai_isan_lover","thai_northern_lover",
    "thai_southern_lover","alcohol_buyer","tea_house_regular","juice_bar_regular",
    "drinks_shop_regular","eco_conscious_shopper","travel_bag_enthusiast",
    "brow_treatment_regular","lash_extension_regular","waxing_regular",
    "donation_giver","tattoo_enthusiast",
}

def source_color(tag):
    if tag in NEW_TAGS:     return C_GREEN
    if tag in OVERLAP_TAGS: return C_YELLOW
    return C_WHITE

# ── Enrichment tag mapping (Overlap + New) ─────────────────────────────────────
ENRICHMENT_MAP = {
    # Food & Dining — Overlap
    "frequent_diner":           "Food",
    "fast_food_regular":        "QSR, Burgers, Fried Chicken",
    "fine_dining_enthusiast":   "Fine Dine",
    "cafe_hopper":              "Cafe/Bistro, Coffee",
    "bakery_lover":             "Baked Pastry, Bakery",
    "dessert_lover":            "Dessert, Cakes, Ice Cream",
    "japanese_food_lover":      "Japan, Sushi, Tonkatsu",
    "korean_food_lover":        "Korea",
    "chinese_food_lover":       "China, Dimsum",
    "thai_food_loyalist":       "Thai",
    "western_food_lover":       "Other_Western",
    "italian_food_lover":       "Italy, Pasta",
    "seafood_lover":            "Seafood",
    "food_court_regular":       "Food Court",
    "buffet_lover":             "Buffet",
    "family_restaurant_goer":   "Family/Kids",
    "healthy_eater":            "Salad",
    # Travel — Overlap
    "frequent_traveler":        "Travel",
    "luxury_traveler":          "Luxury",
    # Shopping — Overlap
    "luxury_shopper":           "Luxury, Social Status",
    "fashion_lover":            "Fashion, Clothing/Apparel",
    "sneaker_enthusiast":       "Shoes/footwear",
    "beauty_conscious":         "Beauty/Feminine",
    "home_decor_lover":         "Furniture & Home Decor",
    "sports_gear_buyer":        "Sport & activewear",
    "convenience_store_addict": "Convenient Seekers, Fast Consumer",
    "marketplace_shopper":      "E-Commerce",
    "subscription_stacker":     "Subscription",
    "premium_brand_loyal":      "Social Status",
    "accessories_collector":    "Handbags & accessories",
    "bargain_hunter":           "Value Economy",
    # Entertainment — Overlap
    "movie_lover":              "Entertainment",
    "concert_goer":             "Entertainment",
    "night_owl_entertainer":    "Nightlife",
    "karaoke_regular":          "Nightlife",
    "spa_lover":                "Spa & Massage",
    "music_festival_goer":      "Entertainment",
    "streaming_subscriber":     "Subscription",
    "outdoor_adventurer":       "Adventure",
    "theme_park_visitor":       "Entertainment",
    "nail_salon_regular":       "Nail",
    "hair_salon_loyal":         "Hair",
    "beauty_clinic_visitor":    "Derma/Aesthetic, Treatment",
    "aesthetic_spender":        "Derma/Aesthetic",
    # Transport — Overlap
    "car_owner":                "Car owners",
    # Health — Overlap
    "health_conscious":         "Health",
    "mental_wellness_seeker":   "Wellness and Longevity",
    # Financial — Overlap
    "small_business_spender":   "Biz Owner",
    # Life Stage — Overlap
    "parent":                   "Family/Kids",
    "homeowner":                "Home Center",
    "home_renovator":           "Home Services, Building Materials",
    "young_professional":       "Professional",
    # Education — Overlap
    "lifelong_learner":         "Skill Development & Training",
    "professional_upskiller":   "Skill Development & Training",
    # Real Estate — Overlap
    "diy_homeowner":            "Building Materials, Home Center",
    # Food & Dining — New
    "bbq_lover":                "BBQ",
    "hot_pot_lover":            "Hot Pot",
    "noodle_lover":             "Noodles",
    "steak_lover":              "Steak",
    "french_food_lover":        "France",
    "american_food_lover":      "US",
    "vietnamese_food_lover":    "Vietnam",
    "thai_isan_lover":          "Thai_Isan, Somtum",
    "thai_northern_lover":      "Thai_Northern",
    "thai_southern_lover":      "Thai_Southern",
    "alcohol_buyer":            "Alcohol",
    "tea_house_regular":        "Tea",
    "juice_bar_regular":        "Juice",
    "drinks_shop_regular":      "Drinks",
    # Shopping — New
    "eco_conscious_shopper":    "Minimal/Sustainable",
    "travel_bag_enthusiast":    "Travel & lifestyle bags",
    # Entertainment — New
    "brow_treatment_regular":   "Brow",
    "lash_extension_regular":   "Lash",
    "waxing_regular":           "Wax",
    "tattoo_enthusiast":        "Tatto",
    # Financial — New
    "donation_giver":           "Donation",
}

# ── Tag data: (tag_name, short_definition, path) ──────────────────────────────

CATEGORIES = [
    ("🍽️ Food & Dining", "47 tags", [
        ("frequent_diner",          "กินข้าวนอกบ้าน ≥ 10 txn/เดือน",               "A"),
        ("fast_food_regular",       "fast food ≥ 4 ครั้ง/เดือน",                   "A"),
        ("fine_dining_enthusiast",  "avg ticket > P75",                             "A"),
        ("cafe_hopper",             "ร้านกาแฟ ≥ 6 ครั้ง/เดือน หลาย merchant",      "A"),
        ("bubble_tea_addict",       "ชานมไข่มุก ≥ 8 ครั้ง/เดือน",                 "A"),
        ("delivery_dependent",      "≥ 50% dining ผ่าน delivery app",              "A"),
        ("brunch_person",           "dining Sat-Sun 09:00–13:00 บ่อย",             "A"),
        ("late_night_diner",        "dining หลัง 22:00 ≥ 4 ครั้ง/เดือน",          "A"),
        ("bakery_lover",            "ร้านเบเกอรี่บ่อย",                            "A"),
        ("budget_eater",            "avg ticket dining < P25",                      "A"),
        ("solo_diner",              "avg ticket ต่ำ + merchant ไม่ซ้ำ",            "A"),
        ("street_food_lover",       "ticket < 150 บาท ร้านขนาดเล็ก",              "A"),
        ("dessert_lover",           "ร้านขนม/ไอศกรีม ≥ 6 ครั้ง/เดือน",           "A"),
        ("group_diner",             "avg ticket > 2x solo avg → มักไปกับกลุ่ม",    "A"),
        ("work_lunch_regular",      "dining weekday 11:00–14:00 ≥ 15 ครั้ง/เดือน","A"),
        ("coffee_loyalist",         "≥ 80% coffee txn ที่ร้านเดิม",               "A"),
        ("weekend_foodie",          "dining spend Sat-Sun > 60% ของ dining ทั้งหมด","A"),
        ("japanese_food_lover",     "≥ 50% dining ที่ร้านญี่ปุ่น",                "B"),
        ("korean_food_lover",       "≥ 50% dining ที่ร้านเกาหลี",                 "B"),
        ("chinese_food_lover",      "≥ 50% dining ที่ร้านจีน",                    "B"),
        ("thai_food_loyalist",      "≥ 60% dining ที่ร้านไทย",                    "B"),
        ("western_food_lover",      "≥ 50% dining ที่ร้านตะวันตก",                "B"),
        ("italian_food_lover",      "≥ 50% dining ที่ร้านอิตาเลียน",              "B"),
        ("seafood_lover",           "ร้านอาหารทะเลบ่อย",                           "B"),
        ("vegetarian_friendly",     "ร้าน vegetarian / plant-based",               "B"),
        ("food_court_regular",      "food court ในห้างบ่อย",                       "B"),
        ("buffet_lover",            "บุฟเฟ่ต์ ≥ 2 ครั้ง/เดือน",                  "B"),
        ("family_restaurant_goer",  "ร้าน family-friendly บ่อย",                  "B"),
        ("healthy_eater",           "ร้าน healthy food / salad",                   "B"),
        ("halal_restaurant_goer",   "ร้านอาหาร halal certified บ่อย",             "B"),
        ("omakase_enthusiast",      "omakase ≥ 2 ครั้ง/ปี + high ticket",         "B"),
        ("rooftop_bar_goer",        "rooftop bar / sky lounge บ่อย",               "B"),
        ("plant_based_explorer",    "plant-based / vegan บ่อยขึ้น",               "B"),
        ("bbq_lover",               "ร้านบาร์บีคิว / ปิ้งย่าง ≥ 2 ครั้ง/เดือน", "B"),
        ("hot_pot_lover",           "ร้านสุกี้ / ชาบู ≥ 2 ครั้ง/เดือน",          "B"),
        ("noodle_lover",            "ร้านก๋วยเตี๋ยว / ราเมน / noodle บ่อย",       "B"),
        ("steak_lover",             "ร้านสเต็กบ่อย + avg ticket สูง",              "B"),
        ("french_food_lover",       "≥ 50% dining ที่ร้านฝรั่งเศส",               "B"),
        ("american_food_lover",     "burger / American diner บ่อย",               "B"),
        ("vietnamese_food_lover",   "ร้านอาหารเวียดนาม / pho บ่อย",               "B"),
        ("thai_isan_lover",         "ร้านอาหารอีสานบ่อย",                          "B"),
        ("thai_northern_lover",     "ร้านอาหารเหนือบ่อย",                          "B"),
        ("thai_southern_lover",     "ร้านอาหารใต้บ่อย",                            "B"),
        ("alcohol_buyer",           "bar / wine shop / bottle shop ≥ 2 ครั้ง/เดือน","B"),
        ("tea_house_regular",       "ร้านชา / specialty tea ≥ 3 ครั้ง/เดือน",    "B"),
        ("juice_bar_regular",       "juice bar / smoothie บ่อย",                   "B"),
        ("drinks_shop_regular",     "ร้านเครื่องดื่ม (non-coffee) บ่อย",          "B"),
    ]),
    ("✈️ Travel & Hotel", "25 tags", [
        ("frequent_traveler",       "airline + hotel ≥ 4 ครั้ง/ปี",               "A"),
        ("international_traveler",  "foreign currency ≥ 2 ครั้ง/ปี",             "A"),
        ("domestic_traveler",       "airline/hotel ในประเทศ ไม่มี foreign",       "A"),
        ("budget_traveler",         "LCC airline + hotel avg ต่ำ",                "A"),
        ("luxury_traveler",         "full-service airline + hotel > P75",          "A"),
        ("business_traveler",       "บินจันทร์-ศุกร์ + hotel weekday",            "A"),
        ("leisure_traveler",        "airline/hotel ช่วง weekend/วันหยุด",         "A"),
        ("hotel_loyalist",          "≥ 70% hotel txn ที่ chain เดียว",            "A"),
        ("airbnb_user",             "Airbnb ≥ 1 ครั้ง/ปี",                       "A"),
        ("frequent_flyer",          "airline ≥ 6 ครั้ง/ปี",                      "A"),
        ("backpacker_style",        "ท่องเที่ยวบ่อย + avg spend ต่ำ",            "A"),
        ("last_minute_booker",      "booking < 7 วันก่อนเดินทาง",                "A"),
        ("early_planner",           "booking > 30 วันล่วงหน้า",                  "A"),
        ("staycation_lover",        "hotel ในประเทศ ไม่มี airline ร่วม",          "A"),
        ("duty_free_shopper",       "duty free txn ≥ 2 ครั้ง/ปี",               "A"),
        ("long_stay_traveler",      "hotel ≥ 5 คืน/trip",                         "A"),
        ("solo_traveler",           "airline 1 ที่นั่ง + hotel single room",      "A"),
        ("asia_traveler",           "foreign txn ส่วนใหญ่ใน Asia",               "B"),
        ("europe_traveler",         "foreign txn ที่ Europe",                      "B"),
        ("resort_lover",            "hotel txn ที่ resort",                        "B"),
        ("city_traveler",           "hotel txn ที่ city hotel",                    "B"),
        ("japan_traveler",          "foreign txn ที่ Japan บ่อย",                 "B"),
        ("luxury_resort_seeker",    "resort + avg > P90",                          "B"),
        ("cruise_passenger",        "cruise line txn ≥ 1 ครั้ง/ปี",             "B"),
        ("workcation_traveler",     "hotel weekday ≥ 5 คืน + ไม่ใช่ business",   "B"),
    ]),
    ("🛍️ Shopping", "34 tags", [
        ("online_first_shopper",    "≥ 50% txn เป็น online",                      "A"),
        ("luxury_shopper",          "avg ticket > P90 หรือ luxury merchant",       "A"),
        ("fashion_lover",           "apparel ≥ 4 ครั้ง/เดือน",                   "A"),
        ("sneaker_enthusiast",      "ร้านรองเท้า/sneaker บ่อย",                   "A"),
        ("beauty_conscious",        "ร้านเครื่องสำอาง ≥ 3 ครั้ง/เดือน",         "A"),
        ("tech_gadget_buyer",       "electronics ≥ 2 ครั้ง/ไตรมาส",             "A"),
        ("home_decor_lover",        "furniture/home goods บ่อย",                  "A"),
        ("book_lover",              "bookstore ≥ 1 ครั้ง/เดือน",                "A"),
        ("sports_gear_buyer",       "ร้านอุปกรณ์กีฬา ≥ 2 ครั้ง/ไตรมาส",        "A"),
        ("grocery_regular",         "supermarket ≥ 3 ครั้ง/เดือน",              "A"),
        ("convenience_store_addict","convenience store ≥ 15 ครั้ง/เดือน",        "A"),
        ("department_store_loyal",  "≥ 70% dept txn ที่ chain เดิม",             "A"),
        ("marketplace_shopper",     "Shopee/Lazada ≥ 4 ครั้ง/เดือน",            "A"),
        ("subscription_stacker",    "recurring txn ≥ 3 บริการ/เดือน",           "A"),
        ("impulse_buyer",           "txn สูง หลาย merchant ไม่ consistent",       "A"),
        ("brand_switcher",          "merchant loyalty ต่ำ ชอบลองใหม่",           "A"),
        ("secondhand_buyer",        "platform secondhand เช่น Carousell",         "A"),
        ("seasonal_shopper",        "spike ช่วง 11.11, 12.12, Songkran",         "A"),
        ("premium_brand_loyal",     "repeat premium brand > 3 ครั้ง/ไตรมาส",    "A"),
        ("hypermarket_shopper",     "Makro/Lotus's ≥ 2 ครั้ง/เดือน + large basket","A"),
        ("jewelry_buyer",           "jewelry store txn ≥ 2 ครั้ง/ปี",           "A"),
        ("gold_accumulator",        "gold shop txn recurring สม่ำเสมอ",          "A"),
        ("stationery_regular",      "stationery store ≥ 1 ครั้ง/เดือน",         "A"),
        ("home_office_buyer",       "stationery + electronics combo",             "A"),
        ("fast_fashion_buyer",      "fashion + low ticket + high frequency",       "A"),
        ("accessories_collector",   "jewelry + fashion accessory diverse merchant","A"),
        ("flash_sale_hunter",       "e-commerce spike ช่วง 11.11 > 3x avg",      "A"),
        ("cross_border_shopper",    "foreign e-commerce ≥ 4 ครั้ง/ปี",          "A"),
        ("mall_shopper",            "≥ 60% retail txn ในห้าง",                   "B"),
        ("bargain_hunter",          "outlet/discount store บ่อย",                 "B"),
        ("gifting_regular",         "gift shop ช่วงเทศกาล",                       "B"),
        ("luxury_gifter",           "jewelry + dept store spike ช่วง Valentine",  "B"),
        ("eco_conscious_shopper",   "organic / eco-friendly / sustainable merchant","B"),
        ("travel_bag_enthusiast",   "luggage / travel gear ≥ 2 ครั้ง/ปี",        "B"),
    ]),
    ("🎭 Entertainment & Lifestyle", "29 tags", [
        ("movie_lover",             "cinema ≥ 2 ครั้ง/เดือน",                    "A"),
        ("concert_goer",            "ticketing ≥ 3 ครั้ง/ปี",                   "A"),
        ("gym_member",              "fitness center recurring ≥ 6 เดือน",        "A"),
        ("gamer",                   "gaming platform บ่อย",                       "A"),
        ("night_owl_entertainer",   "entertainment หลัง 22:00 บ่อย",             "A"),
        ("karaoke_regular",         "karaoke ≥ 2 ครั้ง/เดือน",                  "A"),
        ("spa_lover",               "spa/massage ≥ 2 ครั้ง/เดือน",              "A"),
        ("music_festival_goer",     "festival ticketing ≥ 2 ครั้ง/ปี",          "A"),
        ("streaming_subscriber",    "≥ 2 streaming subscriptions",               "A"),
        ("outdoor_adventurer",      "outdoor gear + travel combo",                "A"),
        ("theme_park_visitor",      "theme park ≥ 1 ครั้ง/ปี",                 "A"),
        ("nail_salon_regular",      "nail salon ≥ 2 ครั้ง/เดือน",              "A"),
        ("hair_salon_loyal",        "≥ 70% hair txn ที่ร้านเดิม monthly",       "A"),
        ("sport_event_goer",        "stadium/sport event ≥ 3 ครั้ง/ปี",        "A"),
        ("esports_fan",             "esports event + gaming spend combo",         "A"),
        ("pet_friendly_lifestyle",  "pet + travel + dining combo",                "A"),
        ("sports_fan",              "stadium/sport event ≥ 2 ครั้ง/ปี",        "B"),
        ("yoga_practitioner",       "yoga studio ≥ 2 ครั้ง/เดือน",             "B"),
        ("golf_player",             "golf course ≥ 1 ครั้ง/เดือน",             "B"),
        ("art_culture_enthusiast",  "museum/gallery/art event",                   "B"),
        ("escape_room_fan",         "escape room ≥ 2 ครั้ง/ไตรมาส",            "B"),
        ("board_game_cafe_goer",    "board game cafe บ่อย",                       "B"),
        ("bowling_regular",         "bowling ≥ 1 ครั้ง/เดือน",                 "B"),
        ("pilates_devotee",         "pilates studio ≥ 3 ครั้ง/เดือน",          "B"),
        ("muay_thai_trainee",       "muay thai gym txn recurring",                "B"),
        ("brow_treatment_regular",  "ร้านทำคิ้วบ่อย",                              "B"),
        ("lash_extension_regular",  "ต่อขนตา ≥ 2 ครั้ง/เดือน",                   "B"),
        ("waxing_regular",          "waxing salon ≥ 1 ครั้ง/เดือน",              "B"),
        ("tattoo_enthusiast",       "tattoo studio ≥ 1 ครั้ง",                    "B"),
    ]),
    ("🚗 Transport & Commute", "16 tags", [
        ("daily_commuter",          "fuel/transit weekday ≥ 20 วัน/เดือน",       "A"),
        ("car_owner",               "fuel + parking + auto service ครบ",          "A"),
        ("grab_dependent",          "Grab ≥ 15 ครั้ง/เดือน",                    "A"),
        ("taxi_user",               "taxi บ่อย + Grab น้อย",                     "A"),
        ("fuel_brand_loyal",        "≥ 70% fuel ที่ปั๊มแบรนด์เดียว",            "A"),
        ("ev_driver",               "EV charging station",                         "A"),
        ("bts_mrt_user",            "BTS/MRT บ่อย",                               "A"),
        ("motorbike_user",          "fuel < 200 บาท/ครั้ง สม่ำเสมอ",            "A"),
        ("parking_heavy",           "parking > 15 ครั้ง/เดือน",                 "A"),
        ("long_distance_driver",    "toll บ่อย / fuel ต่างจังหวัด",             "A"),
        ("car_enthusiast",          "auto accessories + car wash + detailing",    "A"),
        ("car_maintenance_regular", "service + tire + oil ≥ 4 ครั้ง/ปี",       "A"),
        ("high_fuel_spender",       "fuel spend > P75",                            "A"),
        ("ride_share_switcher",     "สลับ Grab/Bolt/InDriver บ่อย",              "A"),
        ("van_pool_user",           "minibus/van service txn ช่วง commute",       "A"),
        ("new_car_seeker",          "auto showroom txn ≥ 1 ครั้ง",             "A"),
    ]),
    ("💊 Health & Wellness", "20 tags", [
        ("health_conscious",        "gym + pharmacy + health food ≥ 2 จาก 3",    "A"),
        ("pharmacy_regular",        "pharmacy ≥ 3 ครั้ง/เดือน",                "A"),
        ("hospital_visitor",        "hospital ≥ 2 ครั้ง/ไตรมาส",              "A"),
        ("supplement_buyer",        "Watsons/Boots บ่อย",                         "A"),
        ("organic_food_lover",      "organic/health food store",                  "A"),
        ("dental_care_conscious",   "dental clinic ≥ 2 ครั้ง/ปี",             "A"),
        ("fitness_first_user",      "Fitness First recurring",                    "A"),
        ("eye_care_regular",        "optician ≥ 2 ครั้ง/ปี",                  "A"),
        ("maternal_care",           "OB/GYN + baby products ช่วงเดียวกัน",      "A"),
        ("regular_checkup_goer",    "hospital/clinic recurring ≥ 1 ครั้ง/ไตรมาส","A"),
        ("mental_wellness_seeker",  "meditation app + wellness + spa combo",      "A"),
        ("sports_nutrition_buyer",  "protein supplement + sports gear combo",     "A"),
        ("aesthetic_spender",       "beauty clinic + cosmetics high spend",       "A"),
        ("high_medical_spender",    "hospital/clinic spend > P75",               "A"),
        ("weight_management_focused","gym + low dining + supplement combo",       "A"),
        ("beauty_clinic_visitor",   "aesthetic clinic บ่อย",                     "B"),
        ("mental_health_aware",     "therapy/wellness center",                    "B"),
        ("traditional_medicine",    "ยาแผนโบราณ / แพทย์แผนจีน",               "B"),
        ("specialized_treatment_seeker","specialist clinic ≥ 4 ครั้ง/ปี",       "B"),
        ("fertility_aware",         "fertility clinic / OB + baby combo",         "B"),
    ]),
    ("💳 Financial Behavior", "25 tags", [
        ("high_spender",            "total spend > P75",                          "A"),
        ("mid_spender",             "total spend P25–P75",                        "A"),
        ("low_spender",             "total spend < P25",                          "A"),
        ("consistent_spender",      "CV monthly spend ต่ำ",                      "A"),
        ("seasonal_spender",        "spike > 2x avg ≥ 2 ครั้ง/ปี",             "A"),
        ("end_of_month_spender",    "≥ 40% spend วันที่ 25–31",                 "A"),
        ("paycheck_spender",        "spike ต้นเดือน (1–5)",                     "A"),
        ("installment_user",        "installment ≥ 1 ครั้ง/ไตรมาส",           "A"),
        ("high_utilization",        "credit utilization > 70% สม่ำเสมอ",        "A"),
        ("low_utilization",         "credit utilization < 20%",                  "A"),
        ("multi_card_user",         "spend pattern บ่งชี้หลายบัตร",             "A"),
        ("cash_advance_user",       "cash advance ≥ 1 ครั้ง/ปี",              "A"),
        ("reward_maximizer",        "≥ 50% spend ที่ high-points merchant",      "A"),
        ("single_category_user",    "≥ 80% spend ใน MCC เดียว",                "A"),
        ("investor",                "investment platform txn ≥ 4 ครั้ง/ปี",    "A"),
        ("regular_investor",        "investment txn recurring รายเดือน",        "A"),
        ("insurance_holder",        "insurance premium txn recurring",           "A"),
        ("multi_insurance_holder",  "≥ 2 insurance products recurring",          "A"),
        ("high_installment_user",   "installment txn ≥ 3 active พร้อมกัน",     "A"),
        ("financially_disciplined", "low utilization + no cash advance + consistent","A"),
        ("cash_flow_stressed",      "high utilization + cash advance + EOM heavy","A"),
        ("wealth_builder",          "investor + insurance + low impulse combo",  "A"),
        ("big_ticket_buyer",        "single txn > P95 ≥ 2 ครั้ง/ปี",          "A"),
        ("small_business_spender",  "telecom + ads + supplies + irregular high",  "A"),
        ("donation_giver",          "charity / temple / donation txn recurring",  "B"),
    ]),
    ("👨‍👩‍👧 Life Stage", "22 tags", [
        ("parent",                  "baby/kids + family restaurant บ่อย",         "A"),
        ("new_parent",              "spike baby products ใน 12 เดือน",           "A"),
        ("pet_owner",               "pet shop/vet ≥ 1 ครั้ง/เดือน",            "A"),
        ("homeowner",               "home improvement + utility ≥ 1 ครั้ง/ไตรมาส","A"),
        ("senior_lifestyle",        "pharmacy + hospital + market dominant",      "A"),
        ("empty_nester",            "kids merchant ลด + travel/leisure เพิ่ม",   "A"),
        ("teen_spender",            "gaming + fastfood + convenience dominant",   "A"),
        ("career_starter",          "low spend + professional attire",            "A"),
        ("family_breadwinner",      "high spend + diverse category",              "A"),
        ("retiree_lifestyle",       "weekday ≈ weekend ratio + pharmacy heavy",   "A"),
        ("multi_pet_owner",         "vet + grooming + food + toys diverse",       "A"),
        ("home_renovator",          "HomePro/contractor ≥ 3 ครั้ง/ปี + high",   "A"),
        ("school_fee_payer",        "school/tuition recurring รายเทอม",          "A"),
        ("extracurricular_parent",  "music school + art class + sports academy",  "A"),
        ("sandwich_generation",     "kids txn + senior care txn พร้อมกัน",       "A"),
        ("newly_married",           "jewelry + wedding + honeymoon combo",        "A"),
        ("student",                 "ใกล้มหาวิทยาลัย + stationery + low spend",  "B"),
        ("young_professional",      "coffee weekday เช้า + lunch ใกล้ office",   "B"),
        ("college_age",             "ใกล้มหาวิทยาลัย + delivery + low ticket",   "B"),
        ("digital_nomad",           "co-working space + coffee + travel combo",   "B"),
        ("expat_lifestyle",         "foreign cuisine + international market",      "B"),
        ("newly_married_b",         "wedding vendor + honeymoon enrichment",      "B"),
    ]),
    ("⏰ Time Patterns", "15 tags", [
        ("morning_person",          "≥ 40% txn ก่อน 09:00",                     "A"),
        ("night_owl_spender",       "≥ 30% txn หลัง 22:00",                     "A"),
        ("weekend_warrior",         "≥ 60% spend Sat-Sun",                        "A"),
        ("weekday_spender",         "≥ 70% spend Mon-Fri",                        "A"),
        ("lunch_hour_regular",      "spike 12:00–14:00 ทุก weekday",             "A"),
        ("after_work_spender",      "spike 18:00–21:00",                          "A"),
        ("holiday_spender",         "spend +50% ช่วงวันหยุดยาว",               "A"),
        ("payday_spike",            "spike วันที่ 25–5 ต้นเดือน",               "A"),
        ("monthly_planner",         "spend กระจาย consistent ไม่มี spike",       "A"),
        ("off_peak_shopper",        "retail 09:00–12:00 วันธรรมดา",             "A"),
        ("songkran_spender",        "spend spike ช่วง Songkran ทุกปี",           "A"),
        ("year_end_big_spender",    "spend spike ธ.ค. > 2x avg ทุกปี",          "A"),
        ("ramadan_aware_spender",   "dining pattern เปลี่ยนช่วง Ramadan",        "A"),
        ("event_day_spender",       "spend spike ตรง major event",               "A"),
        ("early_bird_shopper",      "≥ 40% retail txn ก่อน 10:00",             "A"),
    ]),
    ("📱 Digital Behavior", "15 tags", [
        ("online_native",           "≥ 70% txn เป็น online",                     "A"),
        ("contactless_payer",       "≥ 80% in-store เป็น tap/NFC",              "A"),
        ("app_heavy_user",          "≥ 50% txn ผ่าน mobile app",                "A"),
        ("digital_subscription_heavy","≥ 4 recurring digital/เดือน",           "A"),
        ("streaming_only",          "streaming ≥ 2 + cinema = 0",               "A"),
        ("social_commerce_buyer",   "TikTok Shop ≥ 4 ครั้ง",                   "A"),
        ("crypto_curious",          "crypto exchange ≥ 1 ครั้ง",               "A"),
        ("fintech_user",            "fintech platform บ่อย",                     "A"),
        ("high_telecom_spender",    "telecom spend > P75",                        "A"),
        ("device_upgrader",         "new device + high telecom pattern",          "A"),
        ("saas_subscriber",         "software/SaaS recurring ≥ 3",              "A"),
        ("ai_tool_user",            "ChatGPT Plus, Midjourney, etc.",            "A"),
        ("cashless_convert",        "contactless + online > 90% ของ txn",        "A"),
        ("social_ads_spender",      "Facebook/Google/TikTok Ads ≥ 4 ครั้ง/ปี", "A"),
        ("e_wallet_heavy",          "e-wallet > 30% ของ txn",                   "A"),
    ]),
    ("🎓 Education & Learning", "9 tags", [
        ("lifelong_learner",        "online course platform ≥ 4 ครั้ง/ปี",     "A"),
        ("language_learner",        "language school/app txn recurring",         "A"),
        ("professional_upskiller",  "certification / professional course ≥ 2/ปี","A"),
        ("tutor_seeker",            "tutoring center recurring รายเดือน",       "A"),
        ("school_supply_buyer",     "stationery + kids spike ต้นปีการศึกษา",   "A"),
        ("university_student",      "student merchant + dining + low spend",     "A"),
        ("education_investor",      "high education spend > P75 ทั้งครอบครัว",  "A"),
        ("stem_enthusiast",         "electronics + coding subscription combo",   "A"),
        ("arts_student",            "art supplies + music + performance txn",    "A"),
    ]),
    ("🏠 Real Estate & Housing", "7 tags", [
        ("utility_payer",           "utility txn recurring รายเดือน",           "A"),
        ("high_utility_user",       "utility spend > P75",                        "A"),
        ("condo_resident",          "utility + condo management fee recurring",  "A"),
        ("home_appliance_buyer",    "white goods ≥ 1 ครั้ง/ปี + high amount",  "A"),
        ("diy_homeowner",           "hardware store + home improvement บ่อย",   "A"),
        ("property_market_watcher", "real estate agent txn ≥ 1 ครั้ง/ปี",     "A"),
        ("rental_income_owner",     "irregular large incoming + utility multi",  "A"),
    ]),
    ("📊 Investment & Financial Planning", "6 tags", [
        ("stock_market_participant", "brokerage txn ≥ 4 ครั้ง/ปี",            "A"),
        ("mutual_fund_investor",    "fund platform txn recurring",               "A"),
        ("gold_investor",           "gold shop / gold ETF ≥ 4 ครั้ง/ปี",      "A"),
        ("real_estate_investor",    "property expo + large amount",              "A"),
        ("passive_income_seeker",   "investment + insurance + low consumption",  "A"),
        ("financial_planner_client","insurance + investment + consistent saving","A"),
    ]),
]

# ── Build slides ───────────────────────────────────────────────────────────────

def make_appendix_title(slide_num):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    circ = s.shapes.add_shape(9, Inches(9), Inches(0), Inches(5), Inches(5))
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x25, 0x38)
    circ.line.fill.background()
    rect(s, 0, 3.2, 0.12, 1.8, C_ACCENT)
    txt(s, "Appendix", 0.6, 2.2, 9, 0.7, size=20, color=C_ACCENT)
    txt(s, "Customer Tag Library", 0.6, 3.0, 10, 1.1, size=46, bold=True, color=C_WHITE)
    txt(s, "Version 2.1  |  323 Tags  |  13 หมวดหมู่", 0.6, 4.2, 9, 0.5, size=16, color=C_LIGHT)

    stats = [("323", "Tags ทั้งหมด"), ("254", "Path A (ทันที)"), ("69", "Path B (Enrichment)"), ("13", "หมวดหมู่")]
    for i, (n, l) in enumerate(stats):
        x = 0.6 + i * 3.1
        rect(s, x, 5.1, 2.8, 1.5, C_CARD)
        txt(s, n, x + 0.1, 5.2, 2.6, 0.85, size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        txt(s, l, x + 0.1, 6.05, 2.6, 0.4, size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    footer(s, "", 0)

def make_category_slide(cat_title, cat_count, tags, slide_num, part=None):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    subtitle = f"{cat_count}  |  Path A = ทำได้ทันที  |  Path B = รอ Enrichment"
    title_display = cat_title + (f" (ต่อ)" if part == 2 else "")
    header(s, title_display, subtitle)

    # Split into 2 columns
    half = (len(tags) + 1) // 2
    col1, col2 = tags[:half], tags[half:]

    col_x = [0.35, 6.85]
    col_w = 6.3

    for ci, col_tags in enumerate([col1, col2]):
        x = col_x[ci]
        # Column header
        rect(s, x, 0.82, col_w, 0.32, RGBColor(0x0A, 0x25, 0x35))
        txt(s, "Tag", x + 0.1, 0.84, 2.5, 0.26, size=8, bold=True, color=C_ACCENT)
        txt(s, "Definition", x + 2.7, 0.84, 3.1, 0.26, size=8, bold=True, color=C_ACCENT)
        txt(s, "Path", x + 5.9, 0.84, 0.35, 0.26, size=8, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

        for ri, (tag, defn, path) in enumerate(col_tags):
            y = 1.2 + ri * 0.38
            row_bg = RGBColor(0x12, 0x22, 0x30) if ri % 2 == 0 else C_CARD
            rect(s, x, y, col_w, 0.35, row_bg)
            txt(s, tag, x + 0.1, y + 0.04, 2.5, 0.28, size=8.5, bold=True, color=source_color(tag))
            txt(s, defn, x + 2.7, y + 0.04, 3.1, 0.28, size=8, color=C_LIGHT)
            pc = C_ACCENT if path == "A" else C_ORANGE
            rect(s, x + 5.88, y + 0.05, 0.35, 0.24, pc)
            txt(s, path, x + 5.88, y + 0.05, 0.35, 0.24, size=7, bold=True,
                color=C_BG, align=PP_ALIGN.CENTER)

    # Legend (bottom-left)
    legend_x = 0.35
    legend_y = 7.0
    for lc, lt in [(C_WHITE, "Existing"), (C_YELLOW, "Overlap"), (C_GREEN, "New")]:
        rect(s, legend_x, legend_y + 0.05, 0.12, 0.12, lc)
        txt(s, lt, legend_x + 0.16, legend_y, 0.7, 0.22, size=7, color=C_LIGHT)
        legend_x += 0.92

    footer(s, f"A-{slide_num:02d}", 0)

def make_source_summary_slides():
    # Flatten all tags from CATEGORIES preserving order
    all_tags = []
    for _, _, tags in CATEGORIES:
        for tag, _, _ in tags:
            if tag in NEW_TAGS:
                src, sc = "New", C_GREEN
            elif tag in OVERLAP_TAGS:
                src, sc = "Overlap", C_YELLOW
            else:
                src, sc = "Existing", C_WHITE
            enrichment = ENRICHMENT_MAP.get(tag, "")
            all_tags.append((tag, src, sc, enrichment))

    # Sort: Overlap → New → Existing
    order = {"Overlap": 0, "New": 1, "Existing": 2}
    all_tags.sort(key=lambda x: order[x[1]])

    n_o = sum(1 for _, s, _, _ in all_tags if s == "Overlap")
    n_n = sum(1 for _, s, _, _ in all_tags if s == "New")
    n_e = sum(1 for _, s, _, _ in all_tags if s == "Existing")
    subtitle_txt = f"Overlap: {n_o}  |  New: {n_n}  |  Existing: {n_e}  |  ทั้งหมด: {len(all_tags)} tags"

    COL_X   = [0.30, 6.90]
    TAG_W   = 2.5
    SRC_W   = 0.80
    ENR_W   = 2.35
    ROW_H   = 0.31
    COL_HDR = 0.82
    DATA_Y0 = 1.17
    MAX_R   = int((7.0 - DATA_Y0) / ROW_H)   # rows per column
    PER_SLD = MAX_R * 2

    pages  = [all_tags[i:i+PER_SLD] for i in range(0, len(all_tags), PER_SLD)]
    n_pages = len(pages)

    for pidx, page in enumerate(pages):
        s = prs.slides.add_slide(BLANK)
        bg(s)
        header(s, f"Tag Source Overview  ({pidx+1}/{n_pages})", subtitle_txt)

        for ci, cx in enumerate(COL_X):
            bw = TAG_W + SRC_W + ENR_W + 0.05
            rect(s, cx, COL_HDR, bw, 0.30, RGBColor(0x0A, 0x25, 0x35))
            txt(s, "Tag", cx + 0.05, COL_HDR + 0.02, TAG_W, 0.26,
                size=7.5, bold=True, color=C_ACCENT)
            txt(s, "Source", cx + TAG_W + 0.05, COL_HDR + 0.02, SRC_W, 0.26,
                size=7.5, bold=True, color=C_ACCENT)
            txt(s, "Enrichment Tag (ทีม)", cx + TAG_W + SRC_W + 0.08, COL_HDR + 0.02,
                ENR_W, 0.26, size=7.5, bold=True, color=C_ACCENT)

        col1 = page[:MAX_R]
        col2 = page[MAX_R:]

        for ci, col_tags in enumerate([col1, col2]):
            cx = COL_X[ci]
            for ri, (tag, src, sc, enrichment) in enumerate(col_tags):
                y = DATA_Y0 + ri * ROW_H
                row_bg = RGBColor(0x12, 0x22, 0x30) if ri % 2 == 0 else C_CARD
                bw = TAG_W + SRC_W + ENR_W + 0.05
                rect(s, cx, y, bw, ROW_H - 0.02, row_bg)

                # Tag name
                txt(s, tag, cx + 0.05, y + 0.02, TAG_W - 0.1, ROW_H - 0.04,
                    size=7.5, bold=True, color=sc)

                # Source badge
                if src == "Overlap":
                    bc, btc = C_YELLOW, C_BG
                elif src == "New":
                    bc, btc = C_GREEN, C_BG
                else:
                    bc, btc = RGBColor(0x1A, 0x2E, 0x44), C_LIGHT
                rect(s, cx + TAG_W + 0.05, y + 0.05, SRC_W - 0.10, ROW_H - 0.13, bc)
                txt(s, src,
                    cx + TAG_W + 0.05, y + 0.04, SRC_W - 0.10, ROW_H - 0.10,
                    size=6.5, bold=True, color=btc, align=PP_ALIGN.CENTER)

                # Enrichment tag (Overlap + New only)
                if enrichment:
                    txt(s, enrichment,
                        cx + TAG_W + SRC_W + 0.08, y + 0.02, ENR_W, ROW_H - 0.04,
                        size=7, color=C_LIGHT)

        # Legend
        lx, ly = 0.35, 7.0
        for lc, lt in [(C_WHITE, "Existing"), (C_YELLOW, "Overlap"), (C_GREEN, "New")]:
            rect(s, lx, ly + 0.05, 0.12, 0.12, lc)
            txt(s, lt, lx + 0.16, ly, 0.75, 0.22, size=7, color=C_LIGHT)
            lx += 0.95

        footer(s, f"A-SRC-{pidx+1:02d}", 0)

# Build appendix
make_appendix_title(10)

app_n = 1
for cat_title, cat_count, tags in CATEGORIES:
    PER_SLIDE = 34  # max tags per slide (17 per column × 2 cols)
    if len(tags) <= PER_SLIDE:
        make_category_slide(cat_title, cat_count, tags, app_n)
        app_n += 1
    else:
        make_category_slide(cat_title, cat_count, tags[:PER_SLIDE], app_n)
        app_n += 1
        make_category_slide(cat_title, cat_count, tags[PER_SLIDE:], app_n, part=2)
        app_n += 1

make_source_summary_slides()

out = _target
prs.save(out)
print(f"Saved: {out}  |  Total slides: {len(prs.slides)}")
