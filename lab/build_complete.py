"""
build_complete.py
Single-pass build → CustomerTagging_Complete.pptx
ไม่มี deepcopy, ไม่มี multiple open/save
Slide order: Title → Agenda → AgendaKM → MemoV2(2-13) → Layer1(1-4) → Layer2 → Appendix
"""
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_Complete.pptx"

# ── Colors ─────────────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
C_PINK   = RGBColor(0xFF, 0x8C, 0xD4)
C_DARK   = RGBColor(0x0A, 0x14, 0x20)
C_RED    = RGBColor(0xFF, 0x5C, 0x5C)

# ── Presentation setup ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.75, C_DARK)
    rect(slide, 0, 0.75, 13.33, 0.05, C_ACCENT)
    txt(slide, title, 0.5, 0.1, 11, 0.6, size=24, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.55, 11, 0.3, size=11, color=C_ACCENT)

def footer(slide, n, note="Data Science Team | เมษายน 2025"):
    rect(slide, 0, 7.15, 13.33, 0.35, C_DARK)
    txt(slide, note, 0.5, 7.18, 11, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

def km(slide, text):
    """Key Message bar — y=6.6, 24pt centered"""
    rect(slide, 0, 6.58, 13.33, 0.57, C_BG)
    rect(slide, 0, 6.60, 13.33, 0.44, RGBColor(0x04, 0x1A, 0x2C))
    rect(slide, 0, 6.60, 0.07,  0.44, C_ACCENT)
    txt(slide, text, 0.15, 6.62, 13.0, 0.42,
        size=24, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
circ = s.shapes.add_shape(9, Inches(8.5), Inches(-0.5), Inches(6), Inches(6))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x52)
circ.line.fill.background()
rect(s, 0, 2.8, 0.1, 2.0, C_ACCENT)
txt(s, "Customer Intelligence Platform", 0.6, 1.6, 11, 1.0,
    size=44, bold=True, color=C_WHITE)
txt(s, "แผนงานสำหรับการนำเสนอ", 0.6, 2.8, 9, 0.6, size=24, color=C_ACCENT)
txt(s, "ทีม Data Science  |  เมษายน 2025", 0.6, 3.6, 9, 0.4, size=14, color=C_LIGHT)
footer(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.33, 0.75, C_DARK)
rect(s, 0, 0.75, 13.33, 0.05, C_ACCENT)
txt(s, "Agenda", 0.5, 0.1, 11, 0.6, size=24, bold=True, color=C_WHITE)
txt(s, "สิ่งที่เราจะคุยกันวันนี้", 0.5, 0.9, 12, 0.4, size=13, color=C_LIGHT)
agenda_items = [
    ("01", "Big Picture",                   "5 Layers of Intelligence",              C_ACCENT),
    ("02", "Layer 1: Customer Tags",         "Tag Library · สร้าง/แปลง/Validate",    C_GREEN),
    ("03", "Layer 2: Personalized Trigger",  "Automation Flow · Use Cases",           C_YELLOW),
    ("04", "Roadmap & Next Steps",           "แผนการพัฒนาต่อ",                        C_LIGHT),
    ("05", "Appendix",                       "Tag Library เต็ม — 323 Tags",           RGBColor(0xA0, 0xA0, 0xA0)),
]
for i, (num, title, sub, accent) in enumerate(agenda_items):
    top = 1.4 + i * (0.95 + 0.08)
    rect(s, 0.5, top, 12.33, 0.95, C_CARD)
    rect(s, 0.5, top, 0.07, 0.95, accent)
    txt(s, num, 0.65, top + 0.18, 0.6, 0.55, size=26, bold=True, color=accent)
    txt(s, title, 1.35, top + 0.1, 5.5, 0.45, size=18, bold=True, color=C_WHITE)
    txt(s, sub,   1.35, top + 0.55, 8.0, 0.35, size=11, color=C_LIGHT)
footer(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Agenda + Key Message
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.33, 0.75, C_DARK)
rect(s, 0, 0.75, 13.33, 0.05, C_ACCENT)
txt(s, "Agenda — Key Message", 0.5, 0.1, 10, 0.6, size=24, bold=True, color=C_WHITE)
txt(s, "สิ่งสำคัญที่เราอยากให้จำได้จากแต่ละ section", 0.5, 0.55, 10, 0.28, size=11, color=C_ACCENT)
km_sections = [
    ("01", "Big Picture — 5 Layers of Intelligence",
     "ธุรกิจที่ชนะระยะยาวไม่ได้แข่งด้วย cashback — แต่แข่งด้วยความเข้าใจลูกค้าที่ลึกกว่า", C_ACCENT),
    ("02", "Layer 1: Customer Tags",
     "ทุก transaction ที่ลูกค้ารูด คือข้อมูลที่บอกว่าเขาเป็นใคร — 323 tags คือภาษากลางที่ทั้งองค์กรใช้พูดถึงลูกค้าคนเดียวกัน", C_GREEN),
    ("03", "Layer 2: Personalized Trigger",
     "Campaign Manager ตั้ง rule ครั้งเดียว — ระบบทำงานแทนทุกวัน", C_YELLOW),
    ("04", "Roadmap & Next Steps",
     "Layer 1-2 คือรากฐาน — ยิ่งเราเก็บ data นานขึ้น intelligence ยิ่งลึกขึ้น และ competitor ตามยิ่งยาก", C_LIGHT),
    ("05", "Appendix: Tag Library",
     "323 tags — reference สำหรับทีมที่อยากดูรายละเอียด", RGBColor(0x70, 0x70, 0x70)),
]
for i, (num, title, keymsg, accent) in enumerate(km_sections):
    top = 0.9 + i * (1.02 + 0.07)
    rect(s, 0.4, top, 12.33, 1.02, C_CARD)
    rect(s, 0.4, top, 0.07, 1.02, accent)
    txt(s, num, 0.55, top + 0.1, 0.55, 0.38, size=20, bold=True, color=accent)
    txt(s, title, 1.22, top + 0.08, 4.2, 0.38, size=13, bold=True, color=C_WHITE)
    rect(s, 1.22, top + 0.5, 11.3, 0.02, RGBColor(0x2A, 0x40, 0x58))
    txt(s, f"Key Message: {keymsg}", 1.22, top + 0.56, 11.3, 0.42, size=11, color=C_LIGHT)
footer(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Big Picture: 5 Layers
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ภาพใหญ่ — เราจะไปถึงไหน",
       "Hyundai Card UNIVERSE Platform — 5 Layers of Intelligence")
txt(s, "ธุรกิจบัตรเครดิตที่แข็งแกร่งที่สุดไม่ได้แข่งด้วย Cashback แต่แข่งด้วยความเข้าใจลูกค้า",
    0.5, 1.05, 12.3, 0.45, size=14, color=C_LIGHT)
layers = [
    ("1","Customer Tags",         "รู้ว่าลูกค้าเป็นใคร",          "Foundation ทั้งหมด",             C_ACCENT, True),
    ("2","Personalized Trigger",  "เสนอ offer ที่ใช่ ให้คนที่ใช่","ลด CAC, เพิ่ม Conversion",       RGBColor(0x4D,0xD0,0xFF), False),
    ("3","Personalized Products", "Card ที่ benefit ตรงแต่ละคน",  "ลด Churn, เพิ่ม Spend per Card", RGBColor(0x7B,0xC8,0xFF), False),
    ("4","Merchant Intelligence", "ขายความแม่นยำให้ Merchant",    "Revenue Stream ใหม่จาก Data",    RGBColor(0xA8,0xD8,0xEA), False),
    ("5","Lifestyle Platform",    "ส่วนหนึ่งของชีวิตลูกค้า",      "Brand Moat คู่แข่ง copy ไม่ได้", RGBColor(0xB0,0xC4,0xDE), False),
]
for i, (num, name, goal, result, col, active) in enumerate(layers):
    y = 1.6 + i * 0.98
    bg_col = RGBColor(0x00, 0x30, 0x45) if active else C_CARD
    rect(s, 0.5, y, 12.3, 0.83, bg_col)
    if active:
        rect(s, 0.5, y, 0.1, 0.83, C_ACCENT)
    rect(s, 0.5, y, 0.9, 0.83, col)
    txt(s, num, 0.5, y + 0.15, 0.9, 0.5, size=22, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name, 1.55, y + 0.08, 3.2, 0.38, size=13, bold=True, color=col)
    txt(s, goal, 4.9, y + 0.1, 3.9, 0.62, size=11, color=C_WHITE)
    txt(s, "→  " + result, 8.9, y + 0.1, 4.0, 0.62, size=11, color=C_LIGHT)
    if active:
        txt(s, "◀ เราอยู่ตรงนี้", 1.55, y + 0.5, 3.0, 0.28, size=9, color=C_ACCENT)
footer(s, 4)
km(s, "ธุรกิจที่ชนะระยะยาวไม่แข่งด้วย cashback — แต่แข่งด้วยความเข้าใจลูกค้า")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Customer Tag Profile
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1: สิ่งที่เราได้ — Customer Tag Profile",
       "Tag Engine รัน daily → Profile 1 ชุดต่อลูกค้า 1 คน บอกว่า \"เขาเป็นใคร\"")
txt(s, "Output คืออะไร?", 0.5, 1.0, 5.5, 0.4, size=14, bold=True, color=C_WHITE)
key_points = [
    ("323 Tags",  "ครอบ 13 มิติพฤติกรรม"),
    ("Score 0–1", "ต่อ Tag ต่อลูกค้า — ยิ่งสูงยิ่ง active"),
    ("12 เดือน",  "lookback window — ไม่ใช่แค่ปัจจุบัน"),
    ("Daily",     "refresh — profile เปลี่ยนตาม behavior จริง"),
]
for j, (kw, desc) in enumerate(key_points):
    y = 1.48 + j * 0.84
    rect(s, 0.5, y, 5.5, 0.72, C_CARD)
    rect(s, 0.5, y, 0.08, 0.72, C_ACCENT)
    txt(s, kw,   0.72, y + 0.07, 1.6, 0.28, size=13, bold=True, color=C_ACCENT)
    txt(s, desc, 0.72, y + 0.38, 4.8, 0.28, size=10, color=C_LIGHT)
txt(s, "→  ลูกค้าคนเดียวกันมีได้หลาย Tags พร้อมกัน", 0.5, 4.88, 5.5, 0.4, size=12, color=C_WHITE)
rect(s, 6.5, 1.0, 6.4, 0.45, RGBColor(0x00, 0x30, 0x45))
txt(s, "ตัวอย่าง: ลูกค้า A", 6.65, 1.07, 6.0, 0.32, size=13, bold=True, color=C_ACCENT)
profile = [
    ("flash_sale_hunter",        "ช้อปหนักช่วง sale event", "0.82", C_ACCENT),
    ("japanese_food_lover",      "ชอบอาหารญี่ปุ่น",         "0.74", C_YELLOW),
    ("wealth_builder",           "กำลังสร้างความมั่งคั่ง",   "0.61", C_GREEN),
    ("weekend_warrior",          "ใช้จ่ายหนักช่วง weekend", "0.55", C_ORANGE),
    ("convenience_store_addict", "7-Eleven บ่อย",           "0.49", C_LIGHT),
]
for j, (tag, meaning, score, col) in enumerate(profile):
    y = 1.55 + j * 0.98
    row_bg = RGBColor(0x12, 0x24, 0x38) if j % 2 == 0 else C_CARD
    rect(s, 6.5, y, 6.4, 0.88, row_bg)
    txt(s, tag,     6.65, y + 0.07, 3.0, 0.32, size=10, bold=True, color=col)
    txt(s, meaning, 6.65, y + 0.48, 3.0, 0.32, size=9,  color=C_LIGHT)
    txt(s, score,   9.75, y + 0.16, 1.0, 0.45, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    bar_fill = min(float(score), 1.0) * 2.6
    rect(s, 10.85, y + 0.32, 2.6,  0.16, RGBColor(0x0A, 0x18, 0x28))
    if bar_fill > 0:
        rect(s, 10.85, y + 0.32, bar_fill, 0.16, col)
footer(s, 5)
km(s, "ทุก transaction ที่รูด คือข้อมูลที่บอกว่าลูกค้าเป็นใคร")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Automation Flow
# ══════════════════════════════════════════════════════════════════════════════
KM_FLOW = "tag ที่ถูกต้อง → trigger offer ที่ใช่ → ให้คนที่ใช่ → ในเวลาที่เหมาะ"
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "จาก Tag สู่ Action — Automation Flow",
       "Tags ไม่ใช่แค่ข้อมูล — พวกมันคือ trigger ของทุก campaign และ offer")
steps = [
    ("①","Tag\nEngine",      "รัน daily\nรู้ว่าลูกค้า\nเป็นใคร",      C_ACCENT),
    ("②","Trigger\nEngine",  "ประเมินทุกวัน\nเลือกว่าใคร\nควรได้อะไร",  RGBColor(0x4D,0xD0,0xFF)),
    ("③","Campaign\nManager","เลือก offer\nปรับ message\nตาม tag",       C_YELLOW),
    ("④","Delivery\nChannel","Push / SMS\nEmail\nIn-app",                C_GREEN),
    ("⑤","Response\nTracker","วัดผล\nfeedback กลับ\nปรับ model",         C_ORANGE),
]
box_w = 2.2; gap = 0.28; start_x = 0.45
for i, (num, name, desc, col) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    rect(s, x, 1.3, box_w, 4.5, C_CARD)
    rect(s, x, 1.3, box_w, 0.52, col)
    txt(s, num,  x + 0.05, 1.33, box_w - 0.1, 0.45, size=20, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name, x + 0.1,  1.95, box_w - 0.2, 0.8,  size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, desc, x + 0.1,  2.85, box_w - 0.2, 1.4,  size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < 4:
        ax = x + box_w + 0.03
        txt(s, "→", ax, 3.1, gap + 0.1, 0.5, size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
rect(s, 0.45, 5.95, 12.4, 0.58, RGBColor(0x0A, 0x20, 0x30))
rect(s, 0.45, 5.95, 0.08, 0.58, C_ACCENT)
txt(s, "ทีม DS ออกแบบ rule — ทุก offer มาจาก Tag ของลูกค้า ไม่ว่าจะ automate หรือ manual ก็ตาม",
    0.65, 6.04, 12.0, 0.4, size=11, color=C_WHITE)
footer(s, 6)
km(s, KM_FLOW)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 3 Use Cases
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ตัวอย่าง — Tag ได้แล้วเอาไปใช้ยังไง",
       "3 use cases จาก tag เดียว → trigger → campaign → outcome")
cases = [
    {"tag":"flash_sale_hunter","col":C_ACCENT,"icon":"🛒",
     "trigger":"5 วันก่อน 11.11 / 12.12","action":"Push: \"ผ่อน 0% 3 เดือน\"\npre-approval limit พิเศษ",
     "channel":"Push Notification","outcome":"capture spend ก่อนลูกค้าเลือกบัตรอื่น"},
    {"tag":"japanese_food_lover","col":C_YELLOW,"icon":"🍣",
     "trigger":"ร้านญี่ปุ่นใหม่เปิดใน 5km","action":"In-app: \"แคชแบ็ก 15%\"\nที่ร้านใหม่นี้เฉพาะคุณ",
     "channel":"In-app Banner","outcome":"drive trial + merchant partnership revenue"},
    {"tag":"wealth_builder","col":C_GREEN,"icon":"📈",
     "trigger":"ช่วง Q4 / ต้นปีใหม่","action":"Email: แนะนำ\ninvestment product / premium card",
     "channel":"Email","outcome":"upsell + เพิ่ม LTV ระยะยาว"},
]
card_w = 3.9
row_cols  = [RGBColor(0x00,0x20,0x35), RGBColor(0x0A,0x25,0x18), RGBColor(0x20,0x18,0x00), RGBColor(0x20,0x08,0x08)]
lab_cols  = [C_ACCENT, C_GREEN, C_YELLOW, C_ORANGE]
for i, c in enumerate(cases):
    x = 0.45 + i * (card_w + 0.3)
    col = c["col"]
    rect(s, x, 1.05, card_w, 5.45, C_CARD)
    rect(s, x, 1.05, card_w, 0.6,  col)
    txt(s, c["icon"] + "  " + c["tag"], x + 0.15, 1.1, card_w - 0.3, 0.45, size=11, bold=True, color=C_BG)
    for j, (label, content) in enumerate([("Trigger",c["trigger"]),("Action",c["action"]),("Channel",c["channel"]),("Outcome",c["outcome"])]):
        y = 1.8 + j * 1.17
        rect(s, x + 0.15, y, card_w - 0.3, 1.03, row_cols[j])
        txt(s, label,   x + 0.25, y + 0.05, card_w - 0.5, 0.28, size=9, bold=True, color=lab_cols[j])
        txt(s, content, x + 0.25, y + 0.36, card_w - 0.5, 0.62, size=10, color=C_WHITE)
footer(s, 7)
km(s, KM_FLOW)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 5 Building Blocks
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "สิ่งที่ต้องมีเพื่อให้ Automate ได้จริง",
       "5 building blocks ที่ต้องพร้อมพร้อมกัน — นี่คือ roadmap ที่เราจะสร้าง")
blocks = [
    ("①","Tag Engine",                 "รัน daily pipeline สร้าง customer_tag_profile",
     "Path A: 254 tags | Path B: 69 tags | refresh ทุกคืน",         C_ACCENT),
    ("②","Trigger Rule Engine",        "ประเมิน condition ทุกวัน → สร้าง eligible customer list",
     "rule: IF tag = X AND score ≥ Y AND condition Z → FIRE",        RGBColor(0x4D,0xD0,0xFF)),
    ("③","CRM / Campaign Tool",        "รับ list + tag จาก Trigger Engine → personalize + schedule",
     "ต้องรองรับ tag-based targeting และ dynamic message template",   C_YELLOW),
    ("④","Frequency Cap & Consent",    "ป้องกันส่ง message ถี่เกินไป + เคารพ opt-out ของลูกค้า",
     "max 2 campaigns/week/customer | honor opt-out flag",           C_GREEN),
    ("⑤","Response Tracking Pipeline","วัดว่า campaign ได้ผลไหม → feedback กลับปรับ model",
     "click / convert / ignore → reinforce หรือ adjust threshold",   C_ORANGE),
]
for i, (num, name, desc, detail, col) in enumerate(blocks):
    y = 1.1 + i * 1.06
    rect(s, 0.45, y, 12.4, 0.93, C_CARD)
    rect(s, 0.45, y, 0.72, 0.93, col)
    txt(s, num, 0.45, y + 0.18, 0.72, 0.52, size=20, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name,   1.3, y + 0.07, 3.5, 0.36, size=12, bold=True, color=col)
    txt(s, desc,   1.3, y + 0.5,  4.8, 0.36, size=10, color=C_WHITE)
    txt(s, detail, 6.3, y + 0.07, 6.4, 0.78, size=9,  color=C_LIGHT)
footer(s, 8)
km(s, "automate ได้จริง ต้องมี 5 ส่วนพร้อมกัน — ขาดอันเดียวระบบหยุด")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Layer 1: Problem & Proposal
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1: Customer Tags — ที่เราอยู่ตอนนี้")
rect(s, 0.5, 1.0, 12.3, 1.42, RGBColor(0x2A, 0x10, 0x10))
rect(s, 0.5, 1.0, 0.1, 1.42, C_ORANGE)
txt(s, "ปัญหาที่เราเจออยู่", 0.75, 1.06, 11, 0.38, size=13, bold=True, color=C_ORANGE)
txt(s, "ปัจจุบันเราแบ่งลูกค้าเป็นกลุ่มใหญ่ๆ แล้วส่ง Campaign เดียวกันทั้งกลุ่ม\nทั้งที่ลูกค้าที่ใช้จ่ายเท่ากันอาจมีพฤติกรรมต่างกันอย่างสิ้นเชิง",
    0.75, 1.5, 11.8, 0.82, size=12, color=C_LIGHT)
rect(s, 0.5, 2.6, 12.3, 3.2, C_CARD)
rect(s, 0.5, 2.6, 0.1, 3.2, C_ACCENT)
txt(s, "สิ่งที่เราจะทำ", 0.75, 2.67, 11, 0.38, size=13, bold=True, color=C_ACCENT)
txt(s, "เราจะสร้างระบบ Customer Tagging ที่เปลี่ยนจาก \"ลูกค้าอยู่กลุ่มไหน\" มาเป็น \"ลูกค้าเป็นคนแบบไหน\"\nระบบนี้อ่านพฤติกรรมจาก Transaction ที่มีอยู่แล้ว ไม่ต้องลงทุนซื้อข้อมูลใหม่ และ Refresh ทุกวัน",
    0.75, 3.1, 11.8, 0.72, size=12, color=C_WHITE)
for i, (num, label) in enumerate([("254 Tags","ทำได้ทันที"), ("69 Tags","ต้องใช้ Enrichment")]):
    x = 1.5 + i * 5.5
    col = C_ACCENT if i == 0 else C_ORANGE
    rect(s, x, 4.05, 4.5, 1.55, RGBColor(0x0D, 0x2A, 0x3A) if i == 0 else RGBColor(0x2A, 0x1A, 0x0D))
    txt(s, num, x + 0.2, 4.18, 4.0, 0.72, size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, label, x + 0.2, 4.95, 4.0, 0.45, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
txt(s, "รวม 323 Tags ครอบ 13 มิติพฤติกรรม", 0.5, 5.85, 12.3, 0.4, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
footer(s, 9)
km(s, "เปลี่ยนจาก 'ลูกค้าอยู่กลุ่มไหน' เป็น 'ลูกค้าเป็นคนแบบไหน'")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 254 Tags
# ══════════════════════════════════════════════════════════════════════════════
KM_323 = "323 tags ครอบทุกมิติพฤติกรรม — 254 เริ่มได้เลย ไม่ต้องลงทุนเพิ่ม"
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "254 Tags — ทำได้ทันที",
       "ใช้ข้อมูล Transaction ที่มีอยู่แล้ว ไม่รอ ไม่ต้องลงทุนเพิ่ม")
categories = [
    ("🍽️","Food & Dining","17",["frequent_diner","cafe_hopper","delivery_dependent","late_night_diner","budget_eater"]),
    ("✈️","Travel","13",["frequent_traveler","luxury_traveler","business_traveler","budget_traveler"]),
    ("🛍️","Shopping","19",["marketplace_shopper","luxury_shopper","convenience_store_addict","online_first_shopper"]),
    ("🎭","Entertainment","13",["movie_lover","gym_member","streaming_subscriber","gamer"]),
    ("🚗","Transport","16",["car_owner","grab_dependent","ev_driver","daily_commuter"]),
    ("💊","Health","15",["health_conscious","pharmacy_regular","hospital_visitor"]),
    ("💳","Financial","25",["high_spender","installment_user","reward_maximizer","consistent_spender"]),
    ("👨‍👩‍👧","Life Stage","15",["parent","pet_owner","homeowner","new_parent"]),
    ("⏰","Time Patterns","15",["night_owl_spender","weekend_warrior","morning_person"]),
    ("📱","Digital","15",["online_native","crypto_curious","fintech_user"]),
]
cols = 5
for i, (icon, cat, count, tags) in enumerate(categories):
    col = i % cols; row = i // cols
    x = 0.4 + col * 2.58; y = 1.1 + row * 2.72
    rect(s, x, y, 2.45, 2.5, C_CARD)
    rect(s, x, y, 2.45, 0.5, RGBColor(0x0A, 0x30, 0x45))
    txt(s, icon + " " + cat, x + 0.1, y + 0.05, 2.1, 0.4, size=10, bold=True, color=C_ACCENT)
    txt(s, count + " tags", x + 0.1, y + 0.56, 2.2, 0.3, size=10, color=C_YELLOW, bold=True)
    tag_text = "\n".join(["• " + t for t in tags[:3]])
    txt(s, tag_text, x + 0.1, y + 0.88, 2.3, 1.5, size=8, color=C_LIGHT)
footer(s, 10)
km(s, KM_323)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — 69 Tags (Enrichment)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "69 Tags — ต้องใช้ Merchant Enrichment",
       "เพิ่มความเข้าใจร้านค้าในเชิงลึก เพื่อ Tags ที่ละเอียดขึ้นอีกชั้น")
txt(s, "Merchant Enrichment คืออะไร?", 0.5, 1.05, 12, 0.4, size=15, bold=True, color=C_WHITE)
txt(s, "ปัจจุบัน Transaction บอกได้ว่าลูกค้าไปร้านไหน และหมวดหมู่กว้างๆ เป็นอะไร\nEnrichment คือการเพิ่ม Label ให้ร้านแต่ละแห่งในเชิงลึก เพื่อให้เราแยกแยะได้มากขึ้น",
    0.5, 1.5, 12.3, 0.58, size=12, color=C_LIGHT)
for i, (bg_col, border, title, tcol, rows) in enumerate([
    (RGBColor(0x1A,0x10,0x10), RGBColor(0x3A,0x10,0x10), "❌  ก่อน Enrichment", C_ORANGE,
     [("ร้าน Fuji","→ รู้แค่ \"ร้านอาหาร\""),("Marriott","→ รู้แค่ \"โรงแรม\""),("ผล:","แยกได้แค่ว่า \"ชอบกินข้าวนอกบ้าน\"")]),
    (RGBColor(0x0A,0x1E,0x10), RGBColor(0x0A,0x2E,0x15), "✅  หลัง Enrichment", C_GREEN,
     [("ร้าน Fuji","→ \"อาหารญี่ปุ่น / ในห้าง / กลาง\""),("Marriott","→ \"City Hotel / Luxury / ธุรกิจ\""),("ผล:","แยกได้ว่า \"ชอบอาหารญี่ปุ่นโดยเฉพาะ\"")]),
]):
    x = 0.5 + i * 6.4 + (0.4 if i else 0)
    rect(s, x, 2.25, 5.9, 3.1, bg_col)
    rect(s, x, 2.25, 5.9, 0.42, border)
    txt(s, title, x + 0.2, 2.3, 5.5, 0.32, size=12, bold=True, color=tcol)
    for j, (a, b) in enumerate(rows):
        y = 2.82 + j * 0.78
        txt(s, a, x + 0.2, y, 1.6, 0.42, size=11, bold=True, color=C_WHITE)
        txt(s, b, x + 1.9, y, 3.8, 0.42, size=11, color=C_LIGHT)
txt(s, "→", 6.3, 3.5, 0.6, 0.5, size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
rect(s, 0.5, 5.5, 12.3, 0.72, C_CARD)
txt(s, "Tags ที่ unlock ได้หลัง Enrichment:", 0.7, 5.56, 4.0, 0.3, size=11, bold=True, color=C_ACCENT)
txt(s, "japanese_food_lover  |  korean_food_lover  |  resort_lover  |  mall_shopper  |  golf_player  |  yoga_practitioner  |  young_professional  |  ...",
    0.7, 5.88, 12.0, 0.28, size=10, color=C_LIGHT)
footer(s, 11)
km(s, KM_323)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Tags ชีวิตประจำวัน
# ══════════════════════════════════════════════════════════════════════════════
KM_TAGS = "ทุก tag มีนิยามชัด วัดจากพฤติกรรมจริง — ไม่ใช่การเดา"
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ตัวอย่าง Tags — ชีวิตประจำวัน",
       "แต่ละ Tag มาพร้อม Score 0–1 และ Refresh ทุกวันตามพฤติกรรมจริง")
groups = [
    ("🍽️ Food & Dining", C_ACCENT, [
        ("frequent_diner","กินข้าวนอกบ้าน ≥ 10 ครั้ง/เดือน"),
        ("cafe_hopper","ร้านกาแฟ ≥ 6 ครั้ง/เดือน หลาย merchant"),
        ("delivery_dependent","≥ 50% ของ dining ผ่าน GrabFood/LINE MAN"),
        ("late_night_diner","กินหลัง 22:00 ≥ 4 ครั้ง/เดือน"),
        ("fine_dining_enthusiast","avg ticket ร้านอาหาร > P75"),
    ]),
    ("✈️ Travel", RGBColor(0x7B,0xC8,0xFF), [
        ("frequent_traveler","airline + hotel ≥ 4 ครั้ง/ปี"),
        ("luxury_traveler","full-service airline + hotel avg > P75"),
        ("business_traveler","บินวันจันทร์-ศุกร์ + hotel weekday"),
        ("budget_traveler","LCC airline + hotel avg ต่ำ"),
        ("international_traveler","foreign currency txn ≥ 2 ครั้ง/ปี"),
    ]),
    ("🛍️ Shopping", C_YELLOW, [
        ("marketplace_shopper","Shopee/Lazada ≥ 4 ครั้ง/เดือน"),
        ("luxury_shopper","avg ticket > P90 หรือ luxury brand"),
        ("convenience_store_addict","7-Eleven/FamilyMart ≥ 15 ครั้ง/เดือน"),
        ("online_first_shopper","≥ 50% ของ txn ทั้งหมดเป็น online"),
        ("grocery_regular","Tops/BigC/Lotus ≥ 3 ครั้ง/เดือน"),
    ]),
    ("🎭 Entertainment", C_GREEN, [
        ("movie_lover","โรงหนัง ≥ 2 ครั้ง/เดือน"),
        ("streaming_subscriber","Netflix/Disney+ recurring ≥ 2 บริการ"),
        ("gym_member","ค่า fitness รายเดือนสม่ำเสมอ ≥ 6 เดือน"),
        ("gamer","Steam/PS Store ≥ 12 ครั้ง/ปี"),
        ("concert_goer","ticketing event ≥ 3 ครั้ง/ปี"),
    ]),
]
for i, (title, col, tags) in enumerate(groups):
    x = 0.4 + (i % 2) * 6.45; y = 1.1 + (i // 2) * 2.72
    rect(s, x, y, 6.2, 2.52, C_CARD)
    rect(s, x, y, 6.2, 0.46, RGBColor(0x0A, 0x25, 0x35))
    txt(s, title, x + 0.15, y + 0.06, 5.8, 0.35, size=12, bold=True, color=col)
    for j, (tag, desc) in enumerate(tags):
        ty = y + 0.54 + j * 0.39
        txt(s, tag, x + 0.2, ty, 2.2, 0.35, size=9, bold=True, color=col)
        txt(s, desc, x + 2.5, ty, 3.5, 0.35, size=9, color=C_LIGHT)
footer(s, 12)
km(s, KM_TAGS)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Tags ไลฟ์สไตล์
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ตัวอย่าง Tags — ไลฟ์สไตล์และการเงิน",
       "ครอบคลุม 13 มิติ — ทุก Tag มาจากสิ่งที่ลูกค้าทำอยู่แล้ว ไม่ใช่การเดา")
groups2 = [
    ("🚗 Transport", C_ORANGE, [
        ("car_owner","ครบ 3 อย่าง: fuel + parking + auto service"),
        ("grab_dependent","Grab ≥ 15 ครั้ง/เดือน"),
        ("ev_driver","มี txn ที่ EV Charging Station"),
        ("daily_commuter","fuel/transit weekday ≥ 20 วัน/เดือน"),
        ("fuel_brand_loyal","≥ 70% fuel txn ที่ปั๊มแบรนด์เดียวกัน"),
    ]),
    ("💊 Health & Wellness", C_PINK, [
        ("health_conscious","gym + pharmacy + health food รวมกัน"),
        ("pharmacy_regular","ร้านขายยา ≥ 3 ครั้ง/เดือน"),
        ("hospital_visitor","โรงพยาบาล ≥ 2 ครั้ง/ไตรมาส"),
        ("supplement_buyer","Watsons/Boots ≥ 1 ครั้ง/เดือน"),
        ("gym_member","fitness center recurring ≥ 6 เดือน"),
    ]),
    ("💳 Financial Behavior", C_ACCENT, [
        ("high_spender","total spend > P75 ของ portfolio"),
        ("installment_user","มี installment ≥ 1 ครั้ง/ไตรมาส"),
        ("reward_maximizer","≥ 50% spend ที่ merchant points สูง"),
        ("consistent_spender","monthly spend variation ต่ำ สม่ำเสมอ"),
        ("paycheck_spender","spend spike ต้นเดือนทุกเดือน"),
    ]),
    ("👨‍👩‍👧 Life Stage", C_GREEN, [
        ("parent","baby/kids products + family restaurant บ่อย"),
        ("pet_owner","pet shop/vet ≥ 1 ครั้ง/เดือน"),
        ("homeowner","HomePro/SCG + utility ≥ 1 ครั้ง/ไตรมาส"),
        ("new_parent","spike ที่ baby products ใน 12 เดือนล่าสุด"),
        ("teen_spender","gaming + fastfood + convenience dominant"),
    ]),
]
for i, (title, col, tags) in enumerate(groups2):
    x = 0.4 + (i % 2) * 6.45; y = 1.1 + (i // 2) * 2.72
    rect(s, x, y, 6.2, 2.52, C_CARD)
    rect(s, x, y, 6.2, 0.46, RGBColor(0x0A, 0x25, 0x35))
    txt(s, title, x + 0.15, y + 0.06, 5.8, 0.35, size=12, bold=True, color=col)
    for j, (tag, desc) in enumerate(tags):
        ty = y + 0.54 + j * 0.39
        txt(s, tag, x + 0.2, ty, 2.2, 0.35, size=9, bold=True, color=col)
        txt(s, desc, x + 2.5, ty, 3.5, 0.35, size=9, color=C_LIGHT)
footer(s, 13)
km(s, KM_TAGS)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Plan & Ask
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "แผนการทำงานและสิ่งที่ขออนุมัติ")
txt(s, "เราเสนอให้เริ่ม Phase 1 ด้วย 10 Tags ที่ชัดเจนและวัดผลได้ก่อน ทดสอบกับ Campaign จริง เทียบผลกับ approach เดิม และ scale เมื่อเห็นตัวเลขที่ดีขึ้น",
    0.5, 1.05, 12.3, 0.52, size=12, color=C_LIGHT)
phases = [
    ("Phase 1\nทันที",      C_ACCENT,              ["10 Quick Win Tags","Campaign A/B Test","วัดผลเทียบ approach เดิม"]),
    ("Phase 2\nหลัง Pilot", RGBColor(0x4D,0xD0,0xFF),["ขยาย 254 Tags ครบ","Personalized Trigger","Tags เป็น ML Features"]),
    ("Phase 3\nFull Intel", RGBColor(0x7B,0xC8,0xFF),["69 Tags Enrichment","Merchant Partnership","Revenue Stream ใหม่"]),
    ("Phase 4\nLong-term",  RGBColor(0xB0,0xC4,0xDE),["Personalized Products","Lifestyle Platform","Brand Moat"]),
]
for i, (phase, col, bullets) in enumerate(phases):
    x = 0.4 + i * 3.15
    rect(s, x, 1.75, 3.0, 3.4, C_CARD)
    rect(s, x, 1.75, 3.0, 0.65, col)
    txt(s, phase, x + 0.15, 1.8, 2.7, 0.58, size=12, bold=True, color=C_BG)
    for j, b in enumerate(bullets):
        by = 2.55 + j * 0.78
        rect(s, x + 0.2, by, 2.6, 0.66, RGBColor(0x0D, 0x2A, 0x3A))
        txt(s, "→  " + b, x + 0.35, by + 0.08, 2.3, 0.5, size=10, color=C_WHITE)
txt(s, "สิ่งที่ขอการตัดสินใจ 3 เรื่อง", 0.5, 5.42, 12, 0.38, size=13, bold=True, color=C_WHITE)
asks = [
    ("01","Approve Phase 1", "เริ่ม implement 10 Tags แรก"),
    ("02","กำหนด KPI",       "วัดความสำเร็จก่อน scale — Conversion, Cost per Acquisition"),
    ("03","ระบุ Campaign",   "Campaign แรกที่จะใช้ทดสอบ"),
]
for i, (num, title, desc) in enumerate(asks):
    x = 0.5 + i * 4.3
    rect(s, x, 5.88, 4.1, 0.68, C_CARD)
    rect(s, x, 5.88, 0.56, 0.68, C_ACCENT)
    txt(s, num, x, 5.94, 0.56, 0.54, size=16, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, title, x + 0.68, 5.91, 3.2, 0.3, size=11, bold=True, color=C_ACCENT)
    txt(s, desc,  x + 0.68, 6.22, 3.2, 0.3, size=9, color=C_LIGHT)
footer(s, 14)
km(s, "เริ่มจาก 10 tags — พิสูจน์ผล — แล้วค่อย scale")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Closing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
circ = s.shapes.add_shape(9, Inches(-1.5), Inches(3.5), Inches(6), Inches(6))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x52)
circ.line.fill.background()
lines = [
    ("เรามีข้อมูลลูกค้าอยู่แล้ว", 30, C_WHITE, True),
    ("คำถามคือ  เราอ่านมันได้ดีแค่ไหน", 20, C_LIGHT, False),
    ("", 8, C_WHITE, False),
    ("Customer Tagging คือการเปลี่ยนจาก", 16, C_LIGHT, False),
    ("\"รู้ว่าลูกค้าใช้จ่ายเท่าไหร่\"", 22, C_WHITE, True),
    ("เป็น", 16, C_LIGHT, False),
    ("\"รู้ว่าลูกค้าเป็นใคร\"", 34, C_ACCENT, True),
]
y = 1.4
for text, size, color, bold in lines:
    h = size / 30.0 * 0.6
    txt(s, text, 1.5, y, 10.5, h + 0.1, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER)
    y += h + 0.16
txt(s, "พร้อมนำเสนอรายละเอียดเพิ่มเติมเมื่อสะดวก", 1.5, 6.1, 10.5, 0.38, size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
footer(s, 15)
km(s, "เรามีข้อมูลอยู่แล้ว — คำถามคือเราอ่านมันได้ดีแค่ไหน")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Layer1: Score Formula
# ══════════════════════════════════════════════════════════════════════════════
KM_SCORE = "score บอกว่า active แค่ไหน — dimension บอกว่า active แบบไหน"
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 1: สร้าง Tag Profile",
       "Transaction ทุกครั้งที่รูดบัตร บอกเราว่าลูกค้าเป็นใคร")
flow = [
    ("INPUT",   "Transaction Data",     "12 เดือน | ร้านไหน / เท่าไหร่ / เมื่อไหร่ / ช่องทางไหน", C_LIGHT),
    ("PROCESS", "Tag Engine",           "Path A: 254 tags (ทำได้ทันที)\nPath B: 69 tags (รอ Enrichment)", C_ACCENT),
    ("OUTPUT",  "Tag Score ต่อลูกค้า", "score 0–1 ต่อ tag | refresh ทุกวัน 01:00", C_GREEN),
]
for i, (badge, name, desc, col) in enumerate(flow):
    y = 1.05 + i * 1.72
    rect(s, 0.5, y, 5.9, 1.5, C_CARD)
    rect(s, 0.5, y, 1.25, 1.5, col)
    txt(s, badge, 0.5, y + 0.52, 1.25, 0.44, size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name, 1.9, y + 0.1, 4.3, 0.44, size=14, bold=True, color=col)
    txt(s, desc, 1.9, y + 0.6, 4.2, 0.8, size=11, color=C_LIGHT)
    if i < 2:
        txt(s, "↓", 2.9, y + 1.52, 0.5, 0.28, size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
rect(s, 7.2, 1.05, 5.6, 2.65, C_CARD)
rect(s, 7.2, 1.05, 5.6, 0.42, RGBColor(0x00, 0x30, 0x45))
txt(s, "Score Formula (Scored Tags)", 7.35, 1.1, 5.3, 0.32, size=11, bold=True, color=C_ACCENT)
for i, (formula, meaning) in enumerate([
    ("0.4  ×  Frequency","ใช้บ่อยแค่ไหน"),
    ("0.4  ×  Recency",  "ล่าสุดนานแค่ไหน"),
    ("0.2  ×  Breadth",  "กี่ merchant ที่ต่างกัน"),
]):
    y = 1.6 + i * 0.65
    txt(s, formula, 7.35, y, 2.8, 0.42, size=12, bold=True, color=C_WHITE)
    txt(s, meaning, 10.2, y, 2.5, 0.42, size=10, color=C_LIGHT)
rect(s, 7.2, 3.0, 5.6, 0.5, RGBColor(0x00, 0x1E, 0x30))
txt(s, "threshold ≥ 0.3  →  ได้รับ Tag  (ปรับได้หลัง pilot)",
    7.35, 3.07, 5.3, 0.35, size=10, color=C_ACCENT)
for i, (tag_type, example, desc, col) in enumerate([
    ("Scored Tag","ev_driver: 0.82","วัด intensity — ยิ่งสูงยิ่ง active",C_ACCENT),
    ("Binary Tag","ev_driver: 1 หรือ 0","มี/ไม่มีพฤติกรรม — ไม่มี intensity",C_YELLOW),
]):
    y = 3.7 + i * 1.35
    rect(s, 7.2, y, 5.6, 1.18, RGBColor(0x12, 0x24, 0x38))
    rect(s, 7.2, y, 0.08, 1.18, col)
    txt(s, tag_type, 7.4, y + 0.1, 3.5, 0.36, size=12, bold=True, color=col)
    txt(s, example,  7.4, y + 0.5, 5.2, 0.3,  size=10, color=C_YELLOW)
    txt(s, desc,     7.4, y + 0.82, 5.2, 0.28, size=9, color=C_LIGHT)
rect(s, 0.5, 6.05, 6.0, 0.38, RGBColor(0x00, 0x3A, 0x50))
txt(s, "Tag Library: 323 Tags  |  13 มิติ  |  Path A: 254  |  Path B: 69",
    0.65, 6.1, 5.7, 0.28, size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
footer(s, 16)
km(s, KM_SCORE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Layer1: Dimensions
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 2: แปลง Score ให้ใช้งานได้",
       "จาก 0–1 สู่ label ที่ campaign manager ใช้ได้ทันที")
rect(s, 0.5, 0.9, 12.3, 0.8, RGBColor(0x00, 0x30, 0x45))
rect(s, 0.5, 0.9, 0.08, 0.8, C_ACCENT)
txt(s, 'Core:  score ≥ threshold  →  "japanese_food_lover"  ✓',
    0.72, 1.02, 12.0, 0.55, size=16, bold=True, color=C_WHITE)
txt(s, "อยู่ใน score แล้ว", 0.5, 1.88, 4.5, 0.35, size=11, bold=True, color=C_LIGHT)
txt(s, "ข้อมูลใหม่ที่ score ไม่บอก", 6.9, 1.88, 6.0, 0.35, size=11, bold=True, color=C_GREEN)
in_score = [
    ("Tier",          "heavy / casual",  "เข้มแค่ไหน",   C_ACCENT),
    ("Recency Status","active / lapsed", "ยังอยู่ไหม",   RGBColor(0x4D,0xD0,0xFF)),
    ("Breadth Type",  "wide / narrow",   "กว้างหรือแคบ", RGBColor(0x7B,0xC8,0xFF)),
]
for i, (name, values, question, col) in enumerate(in_score):
    y = 2.32 + i * 1.38
    rect(s, 0.5, y, 6.1, 1.2, C_CARD)
    rect(s, 0.5, y, 0.08, 1.2, col)
    txt(s, name,     0.72, y + 0.1,  3.5, 0.38, size=13, bold=True, color=col)
    txt(s, values,   0.72, y + 0.54, 3.0, 0.32, size=10, color=C_YELLOW)
    txt(s, question, 3.9,  y + 0.35, 2.6, 0.42, size=10, color=C_LIGHT, align=PP_ALIGN.RIGHT)
new_dims = [
    ("Trending",       "rising / stable / fading", "กำลังขึ้นหรือลง",               C_GREEN),
    ("Spend Intensity","high / low",                "ใช้จ่ายหนักแค่ไหนใน category",  C_ORANGE),
    ("Consistency",    "habitual / sporadic",       "สม่ำเสมอหรือเป็นพักๆ",          C_YELLOW),
    ("Tag Age",        "new / established",         "พฤติกรรมใหม่หรือเก่า",           C_PINK),
]
for i, (name, values, question, col) in enumerate(new_dims):
    y = 2.32 + i * 1.03
    rect(s, 6.9, y, 6.0, 0.9, C_CARD)
    rect(s, 6.9, y, 0.08, 0.9, col)
    txt(s, name,     7.1, y + 0.05, 3.0, 0.35, size=12, bold=True, color=col)
    txt(s, values,   7.1, y + 0.5,  2.8, 0.3,  size=9,  color=C_YELLOW)
    txt(s, question, 10.0, y + 0.24, 2.8, 0.42, size=9, color=C_LIGHT)
footer(s, 17)
km(s, KM_SCORE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Layer1: SQL Schema
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 3: Output สำหรับ Campaign Manager",
       "SQL-ready — filter ได้ทันที ไม่ต้องรู้ logic ข้างใน")
rect(s, 0.5, 0.95, 6.1, 0.4, RGBColor(0x00, 0x30, 0x45))
txt(s, "customer_tag_profile", 0.65, 1.0, 5.8, 0.3, size=11, bold=True, color=C_ACCENT)
columns = [
    ("customer_id",    "VARCHAR","รหัสลูกค้า",                     C_LIGHT),
    ("tag",            "VARCHAR","'frequent_diner'",               C_LIGHT),
    ("score",          "FLOAT",  "0.0 – 1.0",                     C_LIGHT),
    ("tier",           "VARCHAR","'heavy'  |  'casual'",           C_ACCENT),
    ("trending",       "VARCHAR","'rising'  |  'stable'  |  'fading'", C_GREEN),
    ("recency_status", "VARCHAR","'active'  |  'lapsed'",          RGBColor(0x4D,0xD0,0xFF)),
    ("breadth_type",   "VARCHAR","'wide'  |  'narrow'",            RGBColor(0x7B,0xC8,0xFF)),
    ("spend_intensity","VARCHAR","'high'  |  'low'",               C_ORANGE),
    ("consistency",    "VARCHAR","'habitual'  |  'sporadic'",      C_YELLOW),
    ("tag_age",        "VARCHAR","'new'  |  'established'",        C_PINK),
    ("as_of_date",     "DATE",   "refresh date",                   C_LIGHT),
]
for i, (col_name, dtype, desc, col) in enumerate(columns):
    y = 1.4 + i * 0.48
    row_bg = RGBColor(0x12, 0x24, 0x38) if i % 2 == 0 else C_CARD
    rect(s, 0.5, y, 6.1, 0.44, row_bg)
    txt(s, col_name, 0.62, y + 0.06, 2.2, 0.3, size=9,  bold=True, color=col)
    txt(s, dtype,    2.9,  y + 0.06, 0.9, 0.3, size=8,  color=C_YELLOW)
    txt(s, desc,     3.85, y + 0.06, 2.7, 0.3, size=8,  color=C_LIGHT)
txt(s, "ตัวอย่าง Query", 7.1, 0.95, 6.0, 0.38, size=12, bold=True, color=C_WHITE)
queries = [
    ("Heavy diner ที่ยัง active",
     "WHERE tag = 'frequent_diner'\n  AND tier = 'heavy'\n  AND recency_status = 'active'", C_ACCENT),
    ("Win-back — เคยกินบ่อยแต่หายไป",
     "WHERE tag = 'frequent_diner'\n  AND trending = 'fading'\n  AND recency_status = 'lapsed'", C_GREEN),
    ("Fine diner — ใช้จ่ายสูงต่อครั้ง",
     "WHERE tag = 'frequent_diner'\n  AND spend_intensity = 'high'", C_ORANGE),
    ("Life event trigger — พฤติกรรมใหม่",
     "WHERE tag = 'car_owner'\n  AND tag_age = 'new'", C_PINK),
]
for i, (label, query, col) in enumerate(queries):
    y = 1.4 + i * 1.3
    rect(s, 7.1, y, 6.0, 1.18, C_CARD)
    rect(s, 7.1, y, 6.0, 0.32, RGBColor(0x00, 0x28, 0x3A))
    rect(s, 7.1, y, 0.08, 1.18, col)
    txt(s, label, 7.28, y + 0.05, 5.7, 0.24, size=9, bold=True, color=col)
    txt(s, query, 7.28, y + 0.36, 5.7, 0.78, size=9, color=C_LIGHT)
footer(s, 18)
km(s, "campaign manager query ได้ตรง — ไม่ต้องรู้ logic ข้างใน")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Layer1: Validate Tag
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 4: Validate Tag",
       "tag ที่สร้างขึ้นมา มีประโยชน์จริงไหม — ก่อน keep ต้องผ่าน 3 ด่าน")
stages = [
    {"num":"1","title":"Sanity Check","subtitle":"ทำได้ทันที ไม่รอ campaign","col":C_ACCENT,
     "checks":[("Coverage","% ลูกค้าที่ได้ tag นี้","< 1%  หรือ  > 80%"),
               ("Score Distribution","histogram ของ score","กระจุกที่ 0 หรือ 1 ผิดปกติ"),
               ("Stability","% ที่ยังมี tag เดือนถัดไป","< 60%  (เปลี่ยนบ่อยเกินไป)")]},
    {"num":"2","title":"Tag Space Quality","subtitle":"ดูว่า tag ซ้ำซ้อนกันไหม","col":C_YELLOW,
     "checks":[("Tag Correlation","correlation matrix\nระหว่าง tag ทุกตัว","corr > 0.8\nกับ tag อื่น"),
               ("Co-occurrence","% ที่มี tag A และ B\nพร้อมกัน","overlap สูงมาก\n→ อาจ merge ได้")]},
    {"num":"3","title":"Business Validation","subtitle":"ทำหลัง deploy — วัดผลจริง","col":C_GREEN,
     "checks":[("Business Signal","spend / churn rate\ntag holder vs non-holder","ไม่ต่างกันเลย"),
               ("Feature Importance","train model: Naive Features + Tag scores\n→ ดู importance ต่อ topline KPI","importance ต่ำมาก"),
               ("Campaign Lift","A/B test: targeted\nvs random","lift ไม่ดีกว่า\nbaseline")]},
]
col_w = 3.9; col_gap = 0.17; sx = 0.5
CT = 1.85; CB = 5.78
for ci, stage in enumerate(stages):
    x = sx + ci * (col_w + col_gap); col = stage["col"]
    n = len(stage["checks"]); gap = 0.1
    ch = (CB - CT - gap * (n - 1)) / n
    rect(s, x, 1.0, col_w, 0.75, col)
    txt(s, f"ด่าน {stage['num']}: {stage['title']}", x + 0.15, 1.05, col_w - 0.2, 0.38, size=12, bold=True, color=C_BG)
    txt(s, stage["subtitle"], x + 0.15, 1.48, col_w - 0.2, 0.24, size=8, color=C_BG)
    for ri, (name, measure, flag) in enumerate(stage["checks"]):
        y = CT + ri * (ch + gap)
        rect(s, x, y, col_w, ch, C_CARD)
        rect(s, x, y, col_w, 0.06, col)
        txt(s, name,    x + 0.15, y + 0.1, col_w - 0.3, 0.34, size=11, bold=True, color=col)
        txt(s, measure, x + 0.15, y + 0.48, col_w - 0.3, ch * 0.36, size=9, color=C_WHITE)
        flag_y = y + ch - 0.38
        rect(s, x + 0.12, flag_y, col_w - 0.27, 0.32, RGBColor(0x2A, 0x10, 0x10))
        txt(s, "⚑  " + flag, x + 0.22, flag_y + 0.04, col_w - 0.4, 0.24, size=8, color=C_ORANGE)
decisions = [("ผ่านทั้ง 3 ด่าน","KEEP  ✓",C_GREEN),("ด่าน 1 ไม่ผ่าน","ปรับ Threshold",C_ACCENT),
             ("ด่าน 2 ไม่ผ่าน","MERGE",C_YELLOW),("ด่าน 3 ไม่ผ่าน","DELETE",C_ORANGE)]
for i, (condition, action, col) in enumerate(decisions):
    x = 0.5 + i * 3.22
    rect(s, x, 5.9, 3.1, 0.62, C_CARD)
    rect(s, x, 5.9, 0.08, 0.62, col)
    txt(s, condition, x + 0.22, 5.93, 2.75, 0.25, size=8,  color=C_LIGHT)
    txt(s, action,    x + 0.22, 6.2,  2.75, 0.28, size=11, bold=True, color=col)
footer(s, 19)
km(s, "tag ที่ดีต้องผ่าน 3 ด่าน — sanity, quality, business signal")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Layer 2: Automation System
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 2: Personalized Trigger — ระบบ Automation",
       "Tag ของ Customer × Tag ของ Offer → Campaign อัตโนมัติ")
FLOW_L = 0.4; FLOW_W = 6.8; BOX_H = 1.0; GAP = 0.36
TOPS = [1.0, 1.0 + BOX_H + GAP, 1.0 + 2*(BOX_H+GAP), 1.0 + 3*(BOX_H+GAP)]
flow_items = [
    (C_YELLOW, "01  กำหนด Rule (ครั้งเดียว)",   "Campaign Manager ระบุ: tag + condition + offer ที่จะ push"),
    (C_ACCENT, "02  Matching Engine  (ทุกวัน)",  "customer_tag_profile  ×  offer_tag_profile  →  eligible list"),
    (C_GREEN,  "03  Auto-Trigger → Campaign Tool","ส่ง offer ผ่าน SMS / Push / Email — ไม่ต้อง query เอง"),
    (C_LIGHT,  "04  Feedback Loop",               "Redemption data → Tag Validation · วัด Campaign Lift"),
]
for i, (accent, title, sub) in enumerate(flow_items):
    top = TOPS[i]
    rect(s, FLOW_L, top, FLOW_W, BOX_H, C_CARD)
    rect(s, FLOW_L, top, 0.08, BOX_H, accent)
    txt(s, title, FLOW_L + 0.18, top + 0.08, FLOW_W - 0.3, 0.4, size=13, bold=True, color=C_WHITE)
    txt(s, sub,   FLOW_L + 0.18, top + 0.52, FLOW_W - 0.3, 0.4, size=10, color=C_LIGHT)
    if i < 3:
        txt(s, "↓", FLOW_L + FLOW_W/2 - 0.15, top + BOX_H + 0.04, 0.4, 0.28, size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)
TBL_L = 7.8; TBL_W = 5.1
txt(s, "Campaign Manager: Before vs After", TBL_L, 1.0, TBL_W, 0.38, size=12, bold=True, color=C_ACCENT)
col_w = (TBL_W - 1.6) / 2; COL1 = TBL_L + 1.6; COL2 = COL1 + col_w
rect(s, TBL_L, 1.42, 1.55,  0.35, RGBColor(0x0A, 0x14, 0x20))
rect(s, COL1,  1.42, col_w, 0.35, RGBColor(0x33, 0x11, 0x11))
rect(s, COL2,  1.42, col_w, 0.35, RGBColor(0x0A, 0x2A, 0x1A))
txt(s, "งาน",    TBL_L + 0.1, 1.47, 1.4,   0.25, size=10, bold=True, color=C_LIGHT)
txt(s, "Before", COL1  + 0.1, 1.47, col_w, 0.25, size=10, bold=True, color=C_RED)
txt(s, "After",  COL2  + 0.1, 1.47, col_w, 0.25, size=10, bold=True, color=C_GREEN)
rows = [("หา Segment","Query SQL เอง","ระบบทำให้"),("Match Offer","Manual","Auto (tag ↔ tag)"),
        ("Timing","Set เองทุก campaign","Rule ทำงานทุกวัน"),("Effort","ทำซ้ำทุกครั้ง","ตั้งครั้งเดียว")]
for i, (label, before, after) in enumerate(rows):
    top = 1.8 + i * 0.54
    bg_r = C_CARD if i % 2 == 0 else RGBColor(0x16, 0x27, 0x3B)
    rect(s, TBL_L, top, 1.55,  0.5, bg_r)
    rect(s, COL1,  top, col_w, 0.5, bg_r)
    rect(s, COL2,  top, col_w, 0.5, bg_r)
    txt(s, label,  TBL_L + 0.1, top + 0.1, 1.4,   0.3, size=10, bold=True, color=C_LIGHT)
    txt(s, before, COL1  + 0.1, top + 0.1, col_w, 0.3, size=10, color=RGBColor(0xFF,0xAA,0xAA))
    txt(s, after,  COL2  + 0.1, top + 0.1, col_w, 0.3, size=10, color=C_GREEN)
txt(s, "* offer_tag_profile ใช้ taxonomy เดียวกับ customer_tag_profile",
    TBL_L, 4.1, TBL_W, 0.3, size=8, color=RGBColor(0x66, 0x88, 0xAA))
footer(s, 20)
km(s, "ตั้ง rule ครั้งเดียว — ระบบทำงานแทนทุกวัน")

# ══════════════════════════════════════════════════════════════════════════════
# Save + Append Tag Library
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
n_main = len(prs.slides)
print(f"✓ Saved {n_main} slides → {OUT}")

print("Appending Tag Library...")
result = subprocess.run(
    ["python3", "/Users/adisornj/Desktop/Thena/lab/append_tags_pptx.py", OUT],
    check=True, capture_output=True, text=True
)
print(f"  {result.stdout.strip()}")

from pptx import Presentation as P2
final = P2(OUT)
print(f"\n✅  CustomerTagging_Complete.pptx — {len(final.slides)} slides")
print(f"   {n_main} content  +  {len(final.slides)-n_main} appendix")
