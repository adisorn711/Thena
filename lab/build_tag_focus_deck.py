from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x2E)
NAVY       = RGBColor(0x16, 0x21, 0x3E)
BLUE       = RGBColor(0x0F, 0x3D, 0x8C)
TEAL       = RGBColor(0x00, 0x86, 0x8A)
GREEN      = RGBColor(0x1A, 0x7A, 0x4A)
AMBER      = RGBColor(0xD4, 0x7E, 0x00)
RED        = RGBColor(0xBF, 0x2B, 0x2B)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
MID_GRAY   = RGBColor(0xCC, 0xD3, 0xE0)
MUTED      = RGBColor(0x55, 0x65, 0x80)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw or Pt(0.75)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, size, bold=False, color=NEAR_BLACK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box

def header(slide, title, subtitle="", slide_num=None):
    rect(slide, 0, 0, W, Inches(1.1), NAVY)
    rect(slide, 0, Inches(1.1), W, Inches(0.05), BLUE)
    tb(slide, Inches(0.45), Inches(0.18), Inches(11), Inches(0.72),
       title, 26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, Inches(0.45), Inches(0.78), Inches(9), Inches(0.35),
           subtitle, 11, color=MID_GRAY)
    if slide_num:
        tb(slide, Inches(12.5), Inches(0.22), Inches(0.7), Inches(0.5),
           str(slide_num), 13, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

def footer(slide, note=""):
    rect(slide, 0, Inches(7.2), W, Inches(0.04), MID_GRAY)
    tb(slide, Inches(0.45), Inches(7.25), Inches(12), Inches(0.22),
       note, 9, color=MUTED)

def card(slide, l, t, w, h, accent_color=BLUE):
    rect(slide, l, t, w, h, WHITE, MID_GRAY, Pt(0.75))
    rect(slide, l, t, Inches(0.07), h, accent_color)

def tag_pill(slide, l, t, label, fill=LIGHT_GRAY, text_color=NEAR_BLACK):
    rect(slide, l, t, Inches(2.6), Inches(0.3), fill)
    tb(slide, l + Inches(0.08), t + Inches(0.04), Inches(2.45), Inches(0.25),
       label, 9, bold=True, color=text_color)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank); bg(s1, WHITE)
rect(s1, 0, 0, Inches(4.4), H, NAVY)
rect(s1, Inches(4.4), 0, Inches(0.08), H, BLUE)

tb(s1, Inches(4.7), Inches(1.6), Inches(8.2), Inches(1.0),
   "Customer Tag Strategy", 38, bold=True, color=NAVY)
tb(s1, Inches(4.7), Inches(2.7), Inches(8.2), Inches(0.65),
   "3 Key Business Drivers", 26, bold=False, color=BLUE)
rect(s1, Inches(4.7), Inches(3.5), Inches(5.5), Inches(0.05), BLUE)
tb(s1, Inches(4.7), Inches(3.7), Inches(8.2), Inches(0.45),
   "Hypermarket  ·  Dining  ·  Petrol", 16, color=MUTED)
tb(s1, Inches(4.7), Inches(4.3), Inches(8.2), Inches(0.4),
   "Tag Prioritization + EV Intention Framework", 13, color=MUTED)

tb(s1, Inches(0.4), Inches(5.8), Inches(3.5), Inches(0.4),
   "TAG LIBRARY v2.2 — FOCUS EDITION", 10, bold=True, color=WHITE)
rect(s1, Inches(0.4), Inches(6.15), Inches(2.4), Inches(0.04), BLUE)
tb(s1, Inches(0.4), Inches(6.3), Inches(3.5), Inches(0.4),
   "Base: 342 tags  |  2026-05-20", 11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 3 Key Drivers Overview
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank); bg(s2, LIGHT_GRAY)
header(s2, "3 Key Business Drivers — Tag Coverage Summary",
       "สรุปภาพรวม tag ที่มีและ gap สำหรับแต่ละ category", 2)

cols = [Inches(0.4), Inches(4.65), Inches(8.9)]
colors = [TEAL, BLUE, AMBER]
icons  = ["🛒", "🍽️", "⛽"]
labels = ["Hypermarket", "Dining", "Petrol"]
existing = ["2 tags", "47 tags", "7 tags"]
gap_counts = ["6 new", "3 new", "6 new"]
notes = [
    "hypermarket_shopper\ngrocery_regular\n+ brand / channel / basket gaps",
    "10 core tags (Path A)\n+ 3 enrichment-ready\n+ 3 new gaps",
    "car_owner · fuel_brand_loyal\nhigh_fuel_spender + more\n+ brand + EV risk",
]

for i, (cx, col, ico, lbl, ex, gp, note) in enumerate(
        zip(cols, colors, icons, labels, existing, gap_counts, notes)):
    card(s2, cx, Inches(1.45), Inches(3.9), Inches(5.6), col)
    tb(s2, cx + Inches(0.18), Inches(1.55), Inches(3.6), Inches(0.55),
       f"{ico}  {lbl}", 20, bold=True, color=col)
    rect(s2, cx + Inches(0.18), Inches(2.15), Inches(3.4), Inches(0.04), col)

    # existing badge
    rect(s2, cx + Inches(0.18), Inches(2.3), Inches(1.55), Inches(0.38), LIGHT_GRAY)
    tb(s2, cx + Inches(0.22), Inches(2.34), Inches(1.45), Inches(0.3),
       f"✓ {ex} existing", 11, color=GREEN, bold=True)

    # gap badge
    rect(s2, cx + Inches(1.85), Inches(2.3), Inches(1.55), Inches(0.38), LIGHT_GRAY)
    tb(s2, cx + Inches(1.89), Inches(2.34), Inches(1.45), Inches(0.3),
       f"+ {gp} to build", 11, color=AMBER, bold=True)

    tb(s2, cx + Inches(0.18), Inches(2.85), Inches(3.55), Inches(2.5),
       note, 12, color=MUTED)

footer(s2, "Tag library v2.2 — 342 total tags  |  Focus: 3 key categories")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Hypermarket
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank); bg(s3, LIGHT_GRAY)
header(s3, "🛒 Hypermarket — Tag Recommendations",
       "Merchant name แยก brand ได้ · Online channel แยกได้ · Path A ทั้งหมด", 3)

# Left: existing
card(s3, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.7), TEAL)
tb(s3, Inches(0.6), Inches(1.45), Inches(5.5), Inches(0.45),
   "✓  Existing — Implement ได้เลย", 13, bold=True, color=TEAL)
existing_hyp = [
    ("hypermarket_shopper",  "Makro/Lotus's ≥ 2 ครั้ง/เดือน + large basket",   "Core tag"),
    ("grocery_regular",      "Supermarket ≥ 3 ครั้ง/เดือน (รวม Tops, BigC)",    "Broader signal"),
]
for j, (tag, defn, note) in enumerate(existing_hyp):
    ty = Inches(2.0) + j * Inches(1.15)
    rect(s3, Inches(0.55), ty, Inches(5.6), Inches(0.95), LIGHT_GRAY, MID_GRAY)
    tb(s3, Inches(0.7), ty + Inches(0.07), Inches(5.3), Inches(0.32),
       tag, 11, bold=True, color=TEAL)
    tb(s3, Inches(0.7), ty + Inches(0.38), Inches(5.3), Inches(0.28),
       defn, 10, color=NEAR_BLACK)
    tb(s3, Inches(0.7), ty + Inches(0.64), Inches(5.3), Inches(0.25),
       note, 9, italic=True, color=MUTED)

# Right: new gaps
card(s3, Inches(6.75), Inches(1.35), Inches(6.2), Inches(5.7), AMBER)
tb(s3, Inches(6.95), Inches(1.45), Inches(5.8), Inches(0.45),
   "+  New — ควรเพิ่ม (Gap Tags)", 13, bold=True, color=AMBER)
new_hyp = [
    ("makro_loyalist",           "≥ 70% hypermarket txn ที่ Makro"),
    ("lotus_loyalist",           "≥ 70% hypermarket txn ที่ Lotus's"),
    ("hypermarket_bulk_buyer",   "avg basket > P75 + ≤ 8 ครั้ง/เดือน"),
    ("hypermarket_top_up_buyer", "≥ 12 ครั้ง/เดือน + avg basket ต่ำ"),
    ("online_grocery_buyer",     "Makro Click / Lotus's online ≥ 2 ครั้ง/เดือน"),
    ("hypermarket_weekend_only", "≥ 80% txn วันเสาร์-อาทิตย์"),
]
for j, (tag, defn) in enumerate(new_hyp):
    ty = Inches(2.0) + j * Inches(0.82)
    rect(s3, Inches(6.9), ty, Inches(5.9), Inches(0.7), LIGHT_GRAY, MID_GRAY)
    tb(s3, Inches(7.05), ty + Inches(0.05), Inches(5.6), Inches(0.28),
       tag, 10, bold=True, color=AMBER)
    tb(s3, Inches(7.05), ty + Inches(0.34), Inches(5.6), Inches(0.28),
       defn, 9, color=NEAR_BLACK)

footer(s3, "grocery_regular ใช้เป็น broader signal — ใช้ hypermarket_shopper + brand tags เป็น core สำหรับ campaign")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dining
# ══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank); bg(s4, LIGHT_GRAY)
header(s4, "🍽️ Dining — Tag Recommendations",
       "47 tags ใน library — คัดมา 3 tiers ตาม priority", 4)

# Priority 1 — Core (left block)
card(s4, Inches(0.4), Inches(1.35), Inches(4.1), Inches(5.7), BLUE)
tb(s4, Inches(0.6), Inches(1.43), Inches(3.8), Inches(0.42),
   "Priority 1 — Core (ทำก่อน)", 13, bold=True, color=BLUE)
p1_tags = [
    "frequent_diner",
    "fast_food_regular",
    "fine_dining_enthusiast",
    "delivery_dependent",
    "food_court_regular",
    "work_lunch_regular",
    "budget_eater",
    "cafe_hopper",
    "weekend_foodie",
    "late_night_diner",
]
for j, t in enumerate(p1_tags):
    ty = Inches(1.95) + j * Inches(0.48)
    rect(s4, Inches(0.55), ty, Inches(3.8), Inches(0.39), LIGHT_GRAY)
    tb(s4, Inches(0.7), ty + Inches(0.06), Inches(3.5), Inches(0.28),
       t, 10, color=NEAR_BLACK)

# Priority 2 — Secondary (middle)
card(s4, Inches(4.85), Inches(1.35), Inches(3.5), Inches(3.1), TEAL)
tb(s4, Inches(5.05), Inches(1.43), Inches(3.2), Inches(0.42),
   "Priority 2 — Venue / Pattern", 12, bold=True, color=TEAL)
p2_tags = [
    ("buffet_lover",      "Path B — รอ enrichment"),
    ("brunch_person",     "Weekend dining signal"),
    ("coffee_loyalist",   "Brand loyalty signal"),
]
for j, (t, note) in enumerate(p2_tags):
    ty = Inches(2.0) + j * Inches(0.88)
    rect(s4, Inches(5.0), ty, Inches(3.2), Inches(0.72), LIGHT_GRAY, MID_GRAY)
    tb(s4, Inches(5.15), ty + Inches(0.05), Inches(2.9), Inches(0.28),
       t, 10, bold=True, color=TEAL)
    tb(s4, Inches(5.15), ty + Inches(0.38), Inches(2.9), Inches(0.25),
       note, 9, italic=True, color=MUTED)

# Priority 3 — Cuisine (middle bottom)
card(s4, Inches(4.85), Inches(4.7), Inches(3.5), Inches(2.35), MID_GRAY)
tb(s4, Inches(5.05), Inches(4.78), Inches(3.2), Inches(0.42),
   "Priority 3 — Cuisine (รอ enrichment)", 12, bold=True, color=MUTED)
tb(s4, Inches(5.05), Inches(5.25), Inches(3.2), Inches(1.55),
   "japanese · korean · thai\nwestern · chinese · italian\n(Path B ทั้งหมด)", 11, color=MUTED)

# New gaps (right)
card(s4, Inches(8.7), Inches(1.35), Inches(4.25), Inches(5.7), AMBER)
tb(s4, Inches(8.9), Inches(1.43), Inches(3.9), Inches(0.42),
   "+  New — Gap Tags", 13, bold=True, color=AMBER)
new_dining = [
    ("dining_high_frequency",  "≥ 25 txn/เดือน — super heavy user\n(frequent_diner ≥ 10 เท่านั้น)"),
    ("dining_chain_loyal",     "≥ 50% dining txn ที่ chain เดิม\nbrand loyalty cross-category"),
    ("dining_spend_tier_high", "avg dining ticket > P75\nทุก venue type — spend tier แยก"),
]
for j, (t, defn) in enumerate(new_dining):
    ty = Inches(2.0) + j * Inches(1.55)
    rect(s4, Inches(8.85), ty, Inches(4.0), Inches(1.35), LIGHT_GRAY, MID_GRAY)
    tb(s4, Inches(9.0), ty + Inches(0.08), Inches(3.8), Inches(0.32),
       t, 10, bold=True, color=AMBER)
    tb(s4, Inches(9.0), ty + Inches(0.45), Inches(3.8), Inches(0.75),
       defn, 9, color=NEAR_BLACK)

footer(s4, "47 tags ใน library — Priority 1 (10 tags) ทำก่อน เพราะ Path A ทั้งหมด ไม่รอ enrichment")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Petrol
# ══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank); bg(s5, LIGHT_GRAY)
header(s5, "⛽ Petrol — Tag Recommendations",
       "Merchant name แยก brand ได้ · Path A ทั้งหมด", 5)

# Existing
card(s5, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.7), AMBER)
tb(s5, Inches(0.6), Inches(1.43), Inches(5.5), Inches(0.42),
   "✓  Existing — Implement ได้เลย", 13, bold=True, color=AMBER)
existing_petrol = [
    ("car_owner",           "fuel + parking + auto service ครบ"),
    ("fuel_brand_loyal",    "≥ 70% fuel ที่ปั๊มแบรนด์เดียว"),
    ("high_fuel_spender",   "fuel spend > P75 ของ portfolio"),
    ("daily_commuter",      "fuel/transit weekday ≥ 20 วัน/เดือน"),
    ("ev_driver",           "EV charging station (EA Anywhere, PEA Volta)"),
    ("long_distance_driver","toll บ่อย + fuel ต่างจังหวัด"),
    ("motorbike_user",      "fuel < 200 บาท/ครั้ง สม่ำเสมอ"),
]
for j, (t, defn) in enumerate(existing_petrol):
    ty = Inches(2.0) + j * Inches(0.72)
    rect(s5, Inches(0.55), ty, Inches(5.6), Inches(0.6), LIGHT_GRAY, MID_GRAY)
    tb(s5, Inches(0.7), ty + Inches(0.05), Inches(5.3), Inches(0.25),
       t, 10, bold=True, color=AMBER)
    tb(s5, Inches(0.7), ty + Inches(0.32), Inches(5.3), Inches(0.22),
       defn, 9, color=NEAR_BLACK)

# New gaps
card(s5, Inches(6.75), Inches(1.35), Inches(6.2), Inches(5.7), BLUE)
tb(s5, Inches(6.95), Inches(1.43), Inches(5.8), Inches(0.42),
   "+  New — ควรเพิ่ม (Gap Tags)", 13, bold=True, color=BLUE)
new_petrol = [
    ("ptt_loyalist",                  "≥ 70% fuel txn ที่ PTT",                    ""),
    ("shell_loyalist",                "≥ 70% fuel txn ที่ Shell",                  ""),
    ("caltex_loyalist",               "≥ 70% fuel txn ที่ Caltex/Chevron",         ""),
    ("petrol_station_ancillary_buyer","non-fuel txn ที่ปั๊ม ≥ 4 ครั้ง/เดือน\n(Amazon Coffee, car wash, ร้านสะดวก)", "cross-sell สูง"),
    ("high_frequency_fuel",           "fuel txn ≥ 10 ครั้ง/เดือน",                "fleet / heavy driver"),
    ("petrol_to_ev_shifting",         "fuel ลด ≥ 25% + EV charging ปรากฏ\nใน 6 เดือน", "⚠ churn risk"),
]
for j, (t, defn, note) in enumerate(new_petrol):
    ty = Inches(1.95) + j * Inches(0.85)
    rect(s5, Inches(6.9), ty, Inches(5.9), Inches(0.73), LIGHT_GRAY, MID_GRAY)
    tb(s5, Inches(7.05), ty + Inches(0.05), Inches(3.8), Inches(0.28),
       t, 10, bold=True, color=BLUE)
    if note:
        tb(s5, Inches(10.5), ty + Inches(0.05), Inches(2.1), Inches(0.28),
           note, 9, bold=True, color=RED if "churn" in note else AMBER,
           align=PP_ALIGN.RIGHT)
    tb(s5, Inches(7.05), ty + Inches(0.38), Inches(5.7), Inches(0.38),
       defn, 9, color=NEAR_BLACK)

footer(s5, "petrol_to_ev_shifting = early churn signal — ดูรายละเอียด EV Framework ใน slide ถัดไป")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EV: Propensity vs Intent
# ══════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank); bg(s6, LIGHT_GRAY)
header(s6, "⚡ EV Churn Risk — Propensity vs Intent",
       "ตอนนี้ทำได้ระดับไหน และจะ upgrade ยังไงเมื่อข้อมูลพร้อม", 6)

# Comparison table header
rect(s6, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.5), NAVY)
cols_x = [Inches(0.45), Inches(4.3), Inches(8.55)]
cols_w = [Inches(3.7), Inches(4.0), Inches(4.1)]
headers = ["", "Intent Detection", "Propensity Profiling"]
h_colors = [WHITE, TEAL, AMBER]
for cx, cw, ht, hc in zip(cols_x, cols_w, headers, h_colors):
    tb(s6, cx, Inches(1.37), cw, Inches(0.44), ht, 12, bold=True, color=hc)

rows = [
    ("Signal",        "เห็นว่าไปโชว์รูม BYD จริง",     "profile คล้ายคนที่ซื้อ EV"),
    ("Precision",     "สูง — เห็นพฤติกรรมตรง",         "ปานกลาง — เดาจาก pattern"),
    ("ทำได้ตอนนี้?",   "ยังไม่ได้\n(รอ dealership รับบัตร)", "ได้เลย"),
    ("Use case",      "Act ก่อนซื้อ 1–2 เดือน",         "Segment for awareness"),
]
row_fills = [WHITE, LIGHT_GRAY, WHITE, LIGHT_GRAY]
for ri, (label, intent, prop) in enumerate(rows):
    ry = Inches(1.9) + ri * Inches(0.72)
    rf = row_fills[ri]
    for cx, cw in zip(cols_x, cols_w):
        rect(s6, cx, ry, cw, Inches(0.66), rf, MID_GRAY, Pt(0.5))
    tb(s6, cols_x[0] + Inches(0.08), ry + Inches(0.12),
       cols_w[0] - Inches(0.1), Inches(0.5), label, 11, bold=True, color=NEAR_BLACK)
    tc = TEAL if "ได้" in intent or "Act" in intent or "สูง" in intent else RED
    tb(s6, cols_x[1] + Inches(0.08), ry + Inches(0.08),
       cols_w[1] - Inches(0.1), Inches(0.55), intent, 10, color=tc if ri != 0 else NEAR_BLACK)
    tc2 = GREEN if "ได้เลย" in prop or "awareness" in prop else NEAR_BLACK
    tb(s6, cols_x[2] + Inches(0.08), ry + Inches(0.08),
       cols_w[2] - Inches(0.1), Inches(0.55), prop, 10, color=tc2)

# Phase arrow
rect(s6, Inches(0.4), Inches(5.85), Inches(5.5), Inches(0.55), AMBER)
tb(s6, Inches(0.55), Inches(5.9), Inches(5.2), Inches(0.42),
   "ตอนนี้: ใช้ Propensity  →  ev_likely_profile", 13, bold=True, color=WHITE)
rect(s6, Inches(6.3), Inches(5.85), Inches(6.6), Inches(0.55), TEAL)
tb(s6, Inches(6.45), Inches(5.9), Inches(6.3), Inches(0.42),
   "อนาคต: Upgrade เป็น Intent  →  ev_showroom_visitor", 13, bold=True, color=WHITE)

rect(s6, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5), LIGHT_GRAY, MID_GRAY)
tb(s6, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4),
   "ev_likely_profile จะกลายเป็น top-of-funnel  |  ev_showroom_visitor เป็น bottom-of-funnel  |  ใช้คู่กันได้ทันทีที่ข้อมูลพร้อม",
   10, color=MUTED)

footer(s6, "ข้อมูล dealership พร้อมเมื่อไหร่ → backtest lag correlation บน ev_driver cohort ทันที")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ev_likely_profile Tag Design
# ══════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank); bg(s7, LIGHT_GRAY)
header(s7, "⚡ ev_likely_profile — Tag Design",
       "Combo tag สำหรับ EV propensity — ทำได้เลยด้วยข้อมูลที่มีอยู่", 7)

# Rule box
card(s7, Inches(0.4), Inches(1.35), Inches(7.4), Inches(5.7), BLUE)
tb(s7, Inches(0.6), Inches(1.43), Inches(7.0), Inches(0.45),
   "Rule: TRUE ถ้าตรง ≥ 3 จาก 5 conditions", 14, bold=True, color=BLUE)

conditions = [
    ("car_owner",                "มีรถอยู่แล้ว — มีของให้แทนด้วย EV",      TEAL),
    ("high_spender / big_ticket_buyer", "afford EV ได้ (> P75 spend)",    GREEN),
    ("tech_gadget_buyer",        "tech early adopter — EV = tech product", BLUE),
    ("fuel_declining",           "fuel spend ลด ≥ 25% ต่อเนื่อง 3 เดือน\n(ยังไม่มี EV txn)", AMBER),
    ("eco_conscious_shopper",    "lifestyle align กับ EV narrative",        TEAL),
]
for j, (cond, reason, col) in enumerate(conditions):
    cy = Inches(2.05) + j * Inches(0.93)
    rect(s7, Inches(0.55), cy, Inches(7.1), Inches(0.82), WHITE, MID_GRAY)
    rect(s7, Inches(0.55), cy, Inches(0.25), Inches(0.82), col)
    tb(s7, Inches(0.88), cy + Inches(0.06), Inches(6.6), Inches(0.3),
       cond, 11, bold=True, color=col)
    tb(s7, Inches(0.88), cy + Inches(0.42), Inches(6.6), Inches(0.35),
       reason, 9, color=NEAR_BLACK)

# Warning box
card(s7, Inches(8.15), Inches(1.35), Inches(4.8), Inches(2.9), RED)
tb(s7, Inches(8.35), Inches(1.43), Inches(4.4), Inches(0.42),
   "⚠  fuel_declining ต้องกรอง", 13, bold=True, color=RED)
tb(s7, Inches(8.35), Inches(1.95), Inches(4.4), Inches(1.8),
   "fuel ลดลง ≠ EV intent เสมอไป\n\nอาจเป็น:\n• Grab ใช้บ่อยขึ้น\n• Work from home\n• ย้ายบ้าน\n\n→ กรอง: grab_dependent = FALSE\n         bts_mrt_user = FALSE",
   10, color=NEAR_BLACK)

# Limitation box
card(s7, Inches(8.15), Inches(4.5), Inches(4.8), Inches(2.55), MUTED)
tb(s7, Inches(8.35), Inches(4.58), Inches(4.4), Inches(0.42),
   "Honest Limitation", 12, bold=True, color=MUTED)
tb(s7, Inches(8.35), Inches(5.08), Inches(4.4), Inches(1.75),
   "นี่คือ Propensity ไม่ใช่ Intent\n\nบอกได้ว่า 'น่าจะสนใจ EV'\nไม่ใช่ 'กำลังจะซื้อ'\n\nใช้สำหรับ awareness campaign\nและ protect share-of-wallet",
   10, color=NEAR_BLACK)

footer(s7, "Validate: เปรียบเทียบ ev_likely_profile = TRUE vs FALSE — อัตรา conversion เป็น ev_driver ต่างกันไหม?")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Build Roadmap
# ══════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank); bg(s8, LIGHT_GRAY)
header(s8, "Build Roadmap — 3 Phases",
       "เรียงตาม data readiness และ business impact", 8)

phases = [
    {
        "label": "Phase 1",
        "sub":   "Path A · Existing · ทำได้เลย",
        "color": GREEN,
        "items": [
            "🛒  hypermarket_shopper, grocery_regular",
            "🍽️  10 core dining tags\n    (frequent_diner → late_night_diner)",
            "⛽  car_owner, fuel_brand_loyal,\n    high_fuel_spender, daily_commuter",
            "⚡  ev_likely_profile (combo tag)",
        ],
    },
    {
        "label": "Phase 2",
        "sub":   "Gap Tags · New Design",
        "color": AMBER,
        "items": [
            "🛒  makro_loyalist, lotus_loyalist,\n    online_grocery_buyer",
            "🍽️  dining_high_frequency,\n    dining_chain_loyal",
            "⛽  ptt/shell/caltex_loyalist,\n    petrol_station_ancillary_buyer,\n    petrol_to_ev_shifting",
        ],
    },
    {
        "label": "Phase 3",
        "sub":   "Enrichment + Future Data",
        "color": BLUE,
        "items": [
            "🍽️  Cuisine tags (Path B)\n    รอ enrichment พร้อม",
            "🛒  bulk_buyer / top_up_buyer\n    รอ basket-level data",
            "⚡  ev_showroom_visitor\n    รอ dealership รับบัตร",
            "⚡  Backtest lag correlation\n    บน ev_driver cohort จริง",
        ],
    },
]

for i, ph in enumerate(phases):
    px = Inches(0.4) + i * Inches(4.3)
    col = ph["color"]
    card(s8, px, Inches(1.35), Inches(4.0), Inches(5.7), col)
    rect(s8, px, Inches(1.35), Inches(4.0), Inches(0.65), col)
    tb(s8, px + Inches(0.18), Inches(1.38), Inches(3.7), Inches(0.38),
       ph["label"], 18, bold=True, color=WHITE)
    tb(s8, px + Inches(0.18), Inches(2.08), Inches(3.7), Inches(0.32),
       ph["sub"], 10, italic=True, color=col)
    rect(s8, px + Inches(0.18), Inches(2.42), Inches(3.2), Inches(0.04), col)
    for j, item in enumerate(ph["items"]):
        iy = Inches(2.6) + j * Inches(1.05)
        rect(s8, px + Inches(0.18), iy, Inches(3.65), Inches(0.9), LIGHT_GRAY, MID_GRAY)
        tb(s8, px + Inches(0.28), iy + Inches(0.07), Inches(3.4), Inches(0.78),
           item, 10, color=NEAR_BLACK)

# Total count
rect(s8, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.28), NAVY)
tb(s8, Inches(0.55), Inches(7.12), Inches(12), Inches(0.24),
   "Phase 1: ~17 tags (implement)   |   Phase 2: ~11 tags (design+build)   |   Phase 3: ~7 tags (pending data)   |   Total new focus: ~35 tags",
   10, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════
out = "/Users/adisornj/Desktop/Thena/lab/CustomerTag_FocusStrategy.pptx"
prs.save(out)
print(f"Saved: {out}")
