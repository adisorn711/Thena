"""
build_automation_slide.py
สร้าง Layer 2 — Automation System slide แล้วแทรกที่ position 20
(หลัง Layer 1 Detail, ก่อน Appendix)
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE   = "/Users/adisornj/Desktop/Thena/lab"
TARGET = f"{BASE}/CustomerTagging_FullDeck.pptx"

C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_RED    = RGBColor(0xFF, 0x5C, 0x5C)
C_DARK   = RGBColor(0x0A, 0x14, 0x20)

def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.75, C_DARK)
    rect(slide, 0, 0.75, 13.33, 0.05, C_ACCENT)
    txt(slide, title, 0.5, 0.1, 11, 0.6, size=24, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.55, 11, 0.3, size=11, color=C_ACCENT)

def footer(slide, n, note="Data Science Team | เมษายน 2025"):
    rect(slide, 0, 7.15, 13.33, 0.35, C_DARK)
    txt(slide, note, 0.5, 7.18, 11, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

def arrow_down(slide, cx, top, h=0.28):
    """วาด arrow ลงตรงกลาง cx"""
    # shaft
    shaft_w = 0.04
    rect(slide, cx - shaft_w/2, top, shaft_w, h * 0.65, RGBColor(0x44, 0x66, 0x88))
    # arrowhead triangle (ใช้ shape type 5 = isoceles triangle, flip vertical)
    from pptx.util import Inches as I
    tip = slide.shapes.add_shape(
        5,
        I(cx - 0.12), I(top + h * 0.55),
        I(0.24), I(h * 0.45)
    )
    tip.fill.solid(); tip.fill.fore_color.rgb = RGBColor(0x44, 0x66, 0x88)
    tip.line.fill.background()

# ── Build slide on temp prs ─────────────────────────────────────────────────
tmp  = Presentation()
tmp.slide_width  = Inches(13.33)
tmp.slide_height = Inches(7.5)
s = tmp.slides.add_slide(tmp.slide_layouts[6])

bg(s)
header(s, "Layer 2: Personalized Trigger — ระบบ Automation",
          "Tag ของ Customer × Tag ของ Offer → Campaign อัตโนมัติ")
footer(s, 20)

# ══════════════════════════════════════════════════════════════════════════════
# LEFT: Flow diagram (4 boxes + arrows)  x: 0.4 – 7.4
# ══════════════════════════════════════════════════════════════════════════════
FLOW_L = 0.4
FLOW_W = 6.8
BOX_H  = 1.0
GAP    = 0.38
TOPS   = [1.0, 1.0 + BOX_H + GAP, 1.0 + 2*(BOX_H + GAP), 1.0 + 3*(BOX_H + GAP)]

flow_items = [
    (C_YELLOW,  "01  กำหนด Rule (ครั้งเดียว)",
                "Campaign Manager ระบุ: tag + condition + offer ที่จะ push"),
    (C_ACCENT,  "02  Matching Engine  (ทุกวัน)",
                "customer_tag_profile  ×  offer_tag_profile  →  eligible list"),
    (C_GREEN,   "03  Auto-Trigger → Campaign Tool",
                "ส่ง offer ผ่าน SMS / Push / Email — ไม่ต้อง query เอง"),
    (C_LIGHT,   "04  Feedback Loop",
                "Redemption data → Tag Validation · วัด Campaign Lift"),
]

for i, (accent, title, sub) in enumerate(flow_items):
    top = TOPS[i]
    rect(s, FLOW_L, top, FLOW_W, BOX_H, C_CARD)
    rect(s, FLOW_L, top, 0.08, BOX_H, accent)
    txt(s, title, FLOW_L + 0.18, top + 0.08, FLOW_W - 0.3, 0.42,
        size=14, bold=True, color=C_WHITE)
    txt(s, sub,   FLOW_L + 0.18, top + 0.52, FLOW_W - 0.3, 0.42,
        size=11, color=C_LIGHT)
    # arrow between boxes
    if i < 3:
        arrow_cx  = FLOW_L + FLOW_W / 2
        arrow_top = top + BOX_H
        arrow_down(s, arrow_cx, arrow_top, GAP)

# ══════════════════════════════════════════════════════════════════════════════
# RIGHT: Before vs After table  x: 7.8 – 12.9
# ══════════════════════════════════════════════════════════════════════════════
TBL_L = 7.8
TBL_W = 5.1

# Section label
txt(s, "Campaign Manager: Before vs After", TBL_L, 1.0, TBL_W, 0.4,
    size=13, bold=True, color=C_ACCENT)

# Column headers
col_w  = (TBL_W - 1.6) / 2
COL1   = TBL_L + 1.6
COL2   = COL1 + col_w

rect(s, TBL_L,  1.45, 1.55,  0.38, RGBColor(0x0A, 0x14, 0x20))
rect(s, COL1,   1.45, col_w, 0.38, RGBColor(0x33, 0x11, 0x11))
rect(s, COL2,   1.45, col_w, 0.38, RGBColor(0x0A, 0x2A, 0x1A))

txt(s, "งาน",    TBL_L + 0.1, 1.5, 1.4,    0.28, size=11, bold=True, color=C_LIGHT)
txt(s, "Before", COL1  + 0.1, 1.5, col_w,  0.28, size=11, bold=True, color=C_RED)
txt(s, "After",  COL2  + 0.1, 1.5, col_w,  0.28, size=11, bold=True, color=C_GREEN)

rows = [
    ("หา Segment",   "Query SQL เอง",         "ระบบทำให้"),
    ("Match Offer",  "Manual",                 "Auto (tag ↔ tag)"),
    ("Timing",       "Set เองทุก campaign",    "Rule ทำงานทุกวัน"),
    ("Effort",       "ทำซ้ำทุกครั้ง",          "ตั้งครั้งเดียว"),
]

ROW_H   = 0.55
ROW_TOP = 1.83

for i, (label, before, after) in enumerate(rows):
    top  = ROW_TOP + i * ROW_H
    bg_r = C_CARD if i % 2 == 0 else RGBColor(0x16, 0x27, 0x3B)

    rect(s, TBL_L, top, 1.55,  ROW_H, bg_r)
    rect(s, COL1,  top, col_w, ROW_H, bg_r)
    rect(s, COL2,  top, col_w, ROW_H, bg_r)

    txt(s, label,  TBL_L + 0.1, top + 0.12, 1.4,   ROW_H - 0.1, size=11, bold=True, color=C_LIGHT)
    txt(s, before, COL1  + 0.1, top + 0.12, col_w, ROW_H - 0.1, size=11, color=RGBColor(0xFF, 0xAA, 0xAA))
    txt(s, after,  COL2  + 0.1, top + 0.12, col_w, ROW_H - 0.1, size=11, color=C_GREEN)

# bottom note
txt(s, "* offer_tag_profile ใช้ taxonomy เดียวกับ customer_tag_profile",
    TBL_L, 5.6, TBL_W, 0.35, size=9, color=RGBColor(0x66, 0x88, 0xAA))

# ══════════════════════════════════════════════════════════════════════════════
# Merge into FullDeck at position 20 (index 19)
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation(TARGET)

blank     = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank)
sp_tree   = new_slide.shapes._spTree
for el in list(sp_tree):
    sp_tree.remove(el)
for shape in s.shapes:
    sp_tree.append(copy.deepcopy(shape.element))

# Move from last → index 19 (position 20)
xml_slides = prs.slides._sldIdLst
entries    = list(xml_slides)
last       = entries[-1]
xml_slides.remove(last)
xml_slides.insert(19, last)

prs.save(TARGET)
print(f"✅  Automation slide inserted at position 20")
print(f"   Total slides: {len(prs.slides)}")
