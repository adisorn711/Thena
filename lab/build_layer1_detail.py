from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)
C_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
C_PINK   = RGBColor(0xFF, 0x8C, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.75, RGBColor(0x0A, 0x14, 0x20))
    rect(slide, 0, 0.75, 13.33, 0.05, C_ACCENT)
    txt(slide, title, 0.5, 0.1, 11, 0.6, size=24, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.55, 11, 0.3, size=11, color=C_ACCENT)

def footer(slide, n, note="Data Science Team | เมษายน 2025"):
    rect(slide, 0, 7.15, 13.33, 0.35, RGBColor(0x0A, 0x14, 0x20))
    txt(slide, note, 0.5, 7.18, 11, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

# ─── Slide 1: สร้าง Tag Profile ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 1: สร้าง Tag Profile",
       "Transaction ทุกครั้งที่รูดบัตร บอกเราว่าลูกค้าเป็นใคร")

# Left: Input → Process → Output
flow = [
    ("INPUT",   "Transaction Data",      "12 เดือน | ร้านไหน / เท่าไหร่ / เมื่อไหร่ / ช่องทางไหน", C_LIGHT),
    ("PROCESS", "Tag Engine",            "Path A: 254 tags (ทำได้ทันที)\nPath B: 69 tags (รอ Enrichment)", C_ACCENT),
    ("OUTPUT",  "Tag Score ต่อลูกค้า",  "score 0–1 ต่อ tag | refresh ทุกวัน 01:00", C_GREEN),
]
for i, (badge, name, desc, col) in enumerate(flow):
    y = 1.05 + i * 1.85
    rect(s, 0.5, y, 5.9, 1.6, C_CARD)
    rect(s, 0.5, y, 1.3, 1.6, col)
    txt(s, badge, 0.5, y + 0.55, 1.3, 0.5,
        size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txt(s, name, 1.95, y + 0.1, 4.3, 0.5, size=15, bold=True, color=col)
    txt(s, desc, 1.95, y + 0.65, 4.2, 0.85, size=12, color=C_LIGHT)
    if i < 2:
        txt(s, "↓", 2.9, y + 1.62, 0.5, 0.3,
            size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Right top: Formula
rect(s, 7.2, 1.05, 5.6, 2.8, C_CARD)
rect(s, 7.2, 1.05, 5.6, 0.45, RGBColor(0x00, 0x30, 0x45))
txt(s, "Score Formula (Scored Tags)", 7.35, 1.1, 5.3, 0.35,
    size=12, bold=True, color=C_ACCENT)

formula_lines = [
    ("0.4  ×  Frequency", "ใช้บ่อยแค่ไหน"),
    ("0.4  ×  Recency",   "ล่าสุดนานแค่ไหน"),
    ("0.2  ×  Breadth",   "กี่ merchant ที่ต่างกัน"),
]
for i, (formula, meaning) in enumerate(formula_lines):
    y = 1.65 + i * 0.7
    txt(s, formula, 7.35, y, 2.8, 0.45, size=13, bold=True, color=C_WHITE)
    txt(s, meaning, 10.2, y, 2.5, 0.45, size=11, color=C_LIGHT)

rect(s, 7.2, 3.15, 5.6, 0.55, RGBColor(0x00, 0x1E, 0x30))
txt(s, "threshold ≥ 0.3  →  ได้รับ Tag  (ปรับได้หลัง pilot)",
    7.35, 3.22, 5.3, 0.38, size=11, color=C_ACCENT)

# Right bottom: Scored vs Binary
for i, (tag_type, example, desc, col) in enumerate([
    ("Scored Tag", "ev_driver: 0.82", "วัด intensity — ยิ่งสูงยิ่ง active", C_ACCENT),
    ("Binary Tag", "ev_driver: 1 หรือ 0", "มี/ไม่มีพฤติกรรม — ไม่มี intensity", C_YELLOW),
]):
    y = 3.95 + i * 1.45
    rect(s, 7.2, y, 5.6, 1.25, RGBColor(0x12, 0x24, 0x38))
    rect(s, 7.2, y, 0.08, 1.25, col)
    txt(s, tag_type, 7.4, y + 0.1, 3.5, 0.4, size=13, bold=True, color=col)
    txt(s, example,  7.4, y + 0.55, 5.2, 0.35, size=11, color=C_YELLOW)
    txt(s, desc,     7.4, y + 0.88, 5.2, 0.3,  size=10, color=C_LIGHT)

# Bottom stat bar
rect(s, 0.5, 6.62, 6.0, 0.42, RGBColor(0x00, 0x3A, 0x50))
txt(s, "Tag Library: 323 Tags  |  13 มิติพฤติกรรม  |  Path A: 254  |  Path B: 69",
    0.65, 6.67, 5.7, 0.32, size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

footer(s, 1)

# ─── Slide 2: แปลง Score ให้ใช้งานได้ ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 2: แปลง Score ให้ใช้งานได้",
       "จาก 0–1 สู่ label ที่ campaign manager ใช้ได้ทันที")

# Core concept banner
rect(s, 0.5, 0.9, 12.3, 0.85, RGBColor(0x00, 0x30, 0x45))
rect(s, 0.5, 0.9, 0.08, 0.85, C_ACCENT)
txt(s, 'Core:  score ≥ threshold  →  "japanese_food_lover"  ✓',
    0.75, 1.02, 12.0, 0.6, size=17, bold=True, color=C_WHITE)

# Column headers
txt(s, "อยู่ใน score แล้ว",
    0.5, 1.95, 4.5, 0.38, size=12, bold=True, color=C_LIGHT)
txt(s, "(แยกไว้เพื่อ usability)",
    4.3, 2.0, 2.5, 0.3, size=9, color=RGBColor(0x60, 0x80, 0x9A))
txt(s, "ข้อมูลใหม่ที่ score ไม่บอก",
    6.9, 1.95, 6.0, 0.38, size=12, bold=True, color=C_GREEN)

# Left column: อยู่ใน score (3 items)
in_score = [
    ("Tier",           "heavy / casual",    "เข้มแค่ไหน",              C_ACCENT),
    ("Recency Status", "active / lapsed",   "ยังอยู่ไหม",              RGBColor(0x4D,0xD0,0xFF)),
    ("Breadth Type",   "wide / narrow",     "กว้างหรือแคบ",            RGBColor(0x7B,0xC8,0xFF)),
]
for i, (name, values, question, col) in enumerate(in_score):
    y = 2.45 + i * 1.52
    rect(s, 0.5, y, 6.1, 1.3, C_CARD)
    rect(s, 0.5, y, 0.08, 1.3, col)
    txt(s, name,     0.75, y + 0.1,  3.5, 0.42, size=14, bold=True, color=col)
    txt(s, values,   0.75, y + 0.58, 3.0, 0.35, size=11, color=C_YELLOW)
    txt(s, question, 3.9,  y + 0.38, 2.6, 0.45, size=11, color=C_LIGHT, align=PP_ALIGN.RIGHT)

# Right column: ข้อมูลใหม่ (4 items)
new_dims = [
    ("Trending",       "rising / stable / fading",  "กำลังขึ้นหรือลง",              C_GREEN),
    ("Spend Intensity","high / low",                 "ใช้จ่ายหนักแค่ไหนใน category", C_ORANGE),
    ("Consistency",    "habitual / sporadic",        "สม่ำเสมอหรือเป็นพักๆ",        C_YELLOW),
    ("Tag Age",        "new / established",          "พฤติกรรมใหม่หรือเก่า",         C_PINK),
]
for i, (name, values, question, col) in enumerate(new_dims):
    y = 2.45 + i * 1.14
    rect(s, 6.9, y, 6.0, 0.98, C_CARD)
    rect(s, 6.9, y, 0.08, 0.98, col)
    txt(s, name,     7.1, y + 0.05, 3.0, 0.38, size=13, bold=True, color=col)
    txt(s, values,   7.1, y + 0.55, 2.8, 0.32, size=10, color=C_YELLOW)
    txt(s, question, 10.0, y + 0.25, 2.8, 0.45, size=10, color=C_LIGHT)

footer(s, 2)

# ─── Slide 3: Output สำหรับ Campaign Manager ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 3: Output สำหรับ Campaign Manager",
       "SQL-ready — filter ได้ทันที ไม่ต้องรู้ logic ข้างใน")

# Left: table schema
rect(s, 0.5, 0.95, 6.1, 0.42, RGBColor(0x00, 0x30, 0x45))
txt(s, "customer_tag_profile", 0.65, 1.0, 5.8, 0.32,
    size=12, bold=True, color=C_ACCENT)

columns = [
    ("customer_id",     "VARCHAR", "รหัสลูกค้า",                    C_LIGHT),
    ("tag",             "VARCHAR", "'frequent_diner'",              C_LIGHT),
    ("score",           "FLOAT",   "0.0 – 1.0",                    C_LIGHT),
    ("tier",            "VARCHAR", "'heavy'  |  'casual'",          C_ACCENT),
    ("trending",        "VARCHAR", "'rising'  |  'stable'  |  'fading'", C_GREEN),
    ("recency_status",  "VARCHAR", "'active'  |  'lapsed'",         RGBColor(0x4D,0xD0,0xFF)),
    ("breadth_type",    "VARCHAR", "'wide'  |  'narrow'",           RGBColor(0x7B,0xC8,0xFF)),
    ("spend_intensity", "VARCHAR", "'high'  |  'low'",              C_ORANGE),
    ("consistency",     "VARCHAR", "'habitual'  |  'sporadic'",     C_YELLOW),
    ("tag_age",         "VARCHAR", "'new'  |  'established'",       C_PINK),
    ("as_of_date",      "DATE",    "refresh date",                  C_LIGHT),
]
for i, (col_name, dtype, desc, col) in enumerate(columns):
    y = 1.42 + i * 0.52
    row_bg = RGBColor(0x12, 0x24, 0x38) if i % 2 == 0 else C_CARD
    rect(s, 0.5, y, 6.1, 0.47, row_bg)
    txt(s, col_name, 0.62, y + 0.06, 2.2, 0.33, size=10, bold=True, color=col)
    txt(s, dtype,    2.9,  y + 0.06, 0.9, 0.33, size=9,  color=C_YELLOW)
    txt(s, desc,     3.85, y + 0.06, 2.7, 0.33, size=9,  color=C_LIGHT)

# Right: example queries
txt(s, "ตัวอย่าง Query", 7.1, 0.95, 6.0, 0.4,
    size=13, bold=True, color=C_WHITE)

queries = [
    ("Heavy diner ที่ยัง active",
     "WHERE tag = 'frequent_diner'\n  AND tier = 'heavy'\n  AND recency_status = 'active'",
     C_ACCENT),
    ("Win-back — เคยกินบ่อยแต่หายไป",
     "WHERE tag = 'frequent_diner'\n  AND trending = 'fading'\n  AND recency_status = 'lapsed'",
     C_GREEN),
    ("Fine diner — ใช้จ่ายสูงต่อครั้ง",
     "WHERE tag = 'frequent_diner'\n  AND spend_intensity = 'high'",
     C_ORANGE),
    ("Life event trigger — พฤติกรรมใหม่",
     "WHERE tag = 'car_owner'\n  AND tag_age = 'new'",
     C_PINK),
]
for i, (label, query, col) in enumerate(queries):
    y = 1.42 + i * 1.42
    rect(s, 7.1, y, 6.0, 1.28, C_CARD)
    rect(s, 7.1, y, 6.0, 0.35, RGBColor(0x00, 0x28, 0x3A))
    rect(s, 7.1, y, 0.08, 1.28, col)
    txt(s, label, 7.3, y + 0.06, 5.7, 0.25, size=10, bold=True, color=col)
    txt(s, query, 7.3, y + 0.42, 5.7, 0.82, size=10, color=C_LIGHT)

footer(s, 3)

# ─── Slide 4: Validate Tag ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Layer 1 — Part 3: Validate Tag",
       "tag ที่สร้างขึ้นมา มีประโยชน์จริงไหม — ก่อน keep ต้องผ่าน 3 ด่าน")

stages = [
    {
        "num": "1", "title": "Sanity Check",
        "subtitle": "ทำได้ทันที ไม่รอ campaign",
        "col": C_ACCENT,
        "checks": [
            ("Coverage",           "% ลูกค้าที่ได้ tag นี้",        "< 1%  หรือ  > 80%"),
            ("Score Distribution", "histogram ของ score",           "กระจุกที่ 0 หรือ 1 ผิดปกติ"),
            ("Stability",          "% ที่ยังมี tag เดือนถัดไป",      "< 60%  (เปลี่ยนบ่อยเกินไป)"),
        ],
    },
    {
        "num": "2", "title": "Tag Space Quality",
        "subtitle": "ดูว่า tag ซ้ำซ้อนกันไหม",
        "col": C_YELLOW,
        "checks": [
            ("Tag Correlation", "correlation matrix\nระหว่าง tag ทุกตัว", "corr > 0.8\nกับ tag อื่น"),
            ("Co-occurrence",   "% ที่มี tag A และ B\nพร้อมกัน",          "overlap สูงมาก\n→ อาจ merge ได้"),
        ],
    },
    {
        "num": "3", "title": "Business Validation",
        "subtitle": "ทำหลัง deploy — วัดผลจริง",
        "col": C_GREEN,
        "checks": [
            ("Business Signal",   "spend / churn rate\ntag holder vs non-holder",           "ไม่ต่างกันเลย"),
            ("Feature Importance","train model: Naive Features + Tag scores\n→ ดู importance ของแต่ละ tag ต่อ topline KPI", "importance ต่ำมาก\n→ tag ไม่ได้เพิ่มอะไร"),
            ("Campaign Lift",     "A/B test: targeted\nvs random",                           "lift ไม่ดีกว่า\nbaseline"),
        ],
    },
]

col_w   = 3.9
col_gap = 0.17
sx      = 0.5
CONTENT_TOP = 1.85
CONTENT_BOT = 6.05

for ci, stage in enumerate(stages):
    x   = sx + ci * (col_w + col_gap)
    col = stage["col"]
    n   = len(stage["checks"])
    gap = 0.12
    ch  = (CONTENT_BOT - CONTENT_TOP - gap * (n - 1)) / n  # dynamic height per check

    # Stage header
    rect(s, x, 1.0, col_w, 0.78, col)
    txt(s, f"ด่าน {stage['num']}: {stage['title']}", x + 0.15, 1.05,
        col_w - 0.2, 0.4, size=13, bold=True, color=C_BG)
    txt(s, stage["subtitle"], x + 0.15, 1.5,
        col_w - 0.2, 0.25, size=9, color=C_BG)

    # Checks
    for ri, (name, measure, flag) in enumerate(stage["checks"]):
        y = CONTENT_TOP + ri * (ch + gap)
        rect(s, x, y, col_w, ch, C_CARD)
        rect(s, x, y, col_w, 0.07, col)
        txt(s, name,    x + 0.15, y + 0.12, col_w - 0.3, 0.38,
            size=12, bold=True, color=col)
        txt(s, measure, x + 0.15, y + 0.55, col_w - 0.3, ch * 0.38,
            size=10, color=C_WHITE)
        # Flag bar
        flag_y = y + ch - 0.42
        rect(s, x + 0.12, flag_y, col_w - 0.27, 0.36, RGBColor(0x2A, 0x10, 0x10))
        txt(s, "⚑  " + flag, x + 0.22, flag_y + 0.04, col_w - 0.4, 0.28,
            size=9, color=C_ORANGE)

# Decision strip
decisions = [
    ("ผ่านทั้ง 3 ด่าน",  "KEEP  ✓",          C_GREEN),
    ("ด่าน 1 ไม่ผ่าน",   "ปรับ Threshold",    C_ACCENT),
    ("ด่าน 2 ไม่ผ่าน",   "MERGE",             C_YELLOW),
    ("ด่าน 3 ไม่ผ่าน",   "DELETE",            C_ORANGE),
]
for i, (condition, action, col) in enumerate(decisions):
    x = 0.5 + i * 3.22
    rect(s, x, 6.2, 3.1, 0.78, C_CARD)
    rect(s, x, 6.2, 0.08, 0.78, col)
    txt(s, condition, x + 0.22, 6.24, 2.75, 0.28, size=9,  color=C_LIGHT)
    txt(s, action,    x + 0.22, 6.55, 2.75, 0.35, size=13, bold=True, color=col)

footer(s, 4)

# ─── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_Layer1_Detail.pptx"
prs.save(out)
print(f"Saved: {out}  (4 slides)")
