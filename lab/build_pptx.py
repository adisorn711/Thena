from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# === Color Palette ===
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # Navy dark
C_ACCENT   = RGBColor(0x00, 0xC2, 0xFF)   # Cyan accent
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xB0, 0xC4, 0xDE)   # Light steel blue
C_CARD     = RGBColor(0x1A, 0x2E, 0x44)   # Card bg
C_YELLOW   = RGBColor(0xFF, 0xD7, 0x00)   # Highlight

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # Blank layout

# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=C_BG, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return txb

def bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)

def accent_bar(slide, y=0.55, h=0.04):
    add_rect(slide, 0.5, y, 2.5, h, C_ACCENT)

def slide_number(slide, n):
    add_text(slide, str(n), 12.6, 7.1, 0.5, 0.3,
             size=10, color=C_LIGHT, align=PP_ALIGN.RIGHT)

def footer(slide, text="Data Science Team | เมษายน 2025"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, RGBColor(0x0A, 0x14, 0x20))
    add_text(slide, text, 0.5, 7.15, 12, 0.3,
             size=9, color=C_LIGHT, align=PP_ALIGN.LEFT)

# ─── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
# Big accent circle decoration
circ = s.shapes.add_shape(9, Inches(9.5), Inches(-1), Inches(6), Inches(6))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x00, 0x5F, 0x7A)
circ.line.fill.background()

add_rect(s, 0, 3.2, 0.08, 1.6, C_ACCENT)
add_text(s, "รู้จักลูกค้าดีขึ้น", 0.6, 1.8, 9, 1.2,
         size=52, bold=True, color=C_WHITE)
add_text(s, "ด้วย Customer Tagging", 0.6, 3.0, 9, 0.9,
         size=40, bold=False, color=C_ACCENT)
add_text(s, "Data Science Team  |  เมษายน 2025", 0.6, 4.3, 8, 0.5,
         size=16, color=C_LIGHT)
footer(s)
slide_number(s, 1)

# ─── Slide 2: Problem ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "ปัญหาที่เรามีอยู่", 0.5, 0.2, 10, 0.6,
         size=28, bold=True, color=C_WHITE)

add_rect(s, 0.5, 1.0, 12.3, 1.4, RGBColor(0x1A, 0x10, 0x10))
add_text(s, "❝  เราแบ่งลูกค้าเป็นกลุ่มใหญ่ๆ แล้วยิง Campaign\n   เดียวกันทั้งกลุ่ม  ❞",
         0.8, 1.05, 11.5, 1.3, size=20, bold=True, color=C_YELLOW, align=PP_ALIGN.LEFT)

points = [
    ("🎯", "ลูกค้าที่ใช้จ่ายเท่ากัน ไม่ได้มีพฤติกรรมเหมือนกัน"),
    ("📉", "Campaign เดียวกันตรงสำหรับบางคน และพลาดสำหรับคนส่วนใหญ่"),
    ("🔁", "Segment แบบเดิม Static — ไม่สะท้อนพฤติกรรมที่เปลี่ยนไป"),
]
for i, (icon, text) in enumerate(points):
    y = 2.7 + i * 0.95
    add_rect(s, 0.5, y, 12.3, 0.8, C_CARD)
    add_text(s, icon, 0.7, y + 0.08, 0.6, 0.65, size=22)
    add_text(s, text, 1.4, y + 0.1, 11, 0.65, size=18, color=C_WHITE)

footer(s); slide_number(s, 2)

# ─── Slide 3: Solution ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "สิ่งที่เราเสนอ", 0.5, 0.2, 10, 0.6, size=28, bold=True, color=C_WHITE)

add_text(s, "เปลี่ยนจาก  \"ลูกค้าอยู่กลุ่มไหน\"  →  \"ลูกค้าเป็นคนแบบไหน\"",
         0.5, 0.9, 12.3, 0.6, size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Two customer cards
cards = [
    ("ลูกค้า A", ["ชอบกินข้าวนอกบ้าน", "ใช้ Grab บ่อย", "ช้อปออนไลน์", "มีรถยนต์"]),
    ("ลูกค้า B", ["นักเดินทางบ่อย", "ใช้จ่ายสูง", "ซื้อของ Luxury", "สมาชิกฟิตเนส"]),
]
for i, (name, tags) in enumerate(cards):
    x = 0.5 + i * 6.5
    add_rect(s, x, 1.8, 6.2, 4.8, C_CARD)
    add_rect(s, x, 1.8, 6.2, 0.55, C_ACCENT)
    add_text(s, name, x + 0.2, 1.85, 5.8, 0.45, size=18, bold=True,
             color=C_BG, align=PP_ALIGN.LEFT)
    for j, tag in enumerate(tags):
        ty = 2.55 + j * 0.9
        add_rect(s, x + 0.3, ty, 5.5, 0.65, RGBColor(0x0D, 0x3A, 0x52))
        add_text(s, "✦  " + tag, x + 0.5, ty + 0.08, 5.2, 0.5,
                 size=16, color=C_WHITE)

add_text(s, "→ เราเรียกสิ่งนี้ว่า  Customer Tags",
         0.5, 6.85, 12.3, 0.4, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
footer(s); slide_number(s, 3)

# ─── Slide 4: No New Data Needed ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "ทำได้จากข้อมูลที่มีอยู่แล้ว", 0.5, 0.2, 10, 0.6,
         size=28, bold=True, color=C_WHITE)
add_text(s, "ไม่ต้องลงทุนซื้อข้อมูลใหม่", 0.5, 0.85, 12, 0.5,
         size=20, color=C_ACCENT)

items = [
    ("Merchant Tag", "รู้ว่าลูกค้าใช้จ่ายที่ไหน ประเภทอะไร"),
    ("Amount", "วิเคราะห์พฤติกรรมการใช้จ่าย"),
    ("Date / Time", "รู้ pattern เวลา — กลางคืน / เช้า / วันหยุด"),
    ("Channel", "รู้ความชอบช่องทาง Online vs In-store"),
    ("Currency", "รู้การใช้จ่ายต่างประเทศ"),
]
for i, (col, desc) in enumerate(items):
    y = 1.6 + i * 0.98
    add_rect(s, 0.5, y, 3.5, 0.78, C_ACCENT)
    add_text(s, col, 0.65, y + 0.1, 3.2, 0.6, size=16, bold=True,
             color=C_BG, align=PP_ALIGN.LEFT)
    add_rect(s, 4.1, y, 8.7, 0.78, C_CARD)
    add_text(s, desc, 4.3, y + 0.1, 8.3, 0.6, size=16, color=C_WHITE)

add_rect(s, 0.5, 6.55, 12.3, 0.45, RGBColor(0x00, 0x3A, 0x50))
add_text(s, "Lookback: 12 เดือนย้อนหลัง  |  Refresh: ทุกวัน  |  เริ่มได้ทันที 117 Tags",
         0.7, 6.6, 12, 0.35, size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

footer(s); slide_number(s, 4)

# ─── Slide 5: Scale of Opportunity ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "ขนาดของโอกาส", 0.5, 0.2, 10, 0.6, size=28, bold=True, color=C_WHITE)

metrics = [
    ("151", "Tags", "ครอบคลุมทุกมิติพฤติกรรม"),
    ("117", "Tags", "ทำได้ทันที ไม่รอข้อมูลเพิ่ม"),
    ("Daily", "Refresh", "ตามพฤติกรรมจริง"),
    ("0–1", "Score", "ทุก Tag มีความมั่นใจ"),
]
for i, (num, label, desc) in enumerate(metrics):
    x = 0.4 + i * 3.15
    add_rect(s, x, 1.2, 3.0, 3.8, C_CARD)
    add_rect(s, x, 1.2, 3.0, 0.08, C_ACCENT)
    add_text(s, num, x + 0.1, 1.5, 2.8, 1.4,
             size=54, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, label, x + 0.1, 2.95, 2.8, 0.55,
             size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.1, 3.6, 2.8, 0.9,
             size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_text(s, "10 มิติพฤติกรรมที่เราจะรู้จักลูกค้า",
         0.5, 5.25, 12.3, 0.5, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
dims = "🍽️ การกิน  |  ✈️ การเดินทาง  |  🛍️ ช้อปปิ้ง  |  🎭 ไลฟ์สไตล์  |  🚗 การเดินทาง\n💊 สุขภาพ  |  💳 การเงิน  |  👨‍👩‍👧 ช่วงชีวิต  |  ⏰ เวลา  |  📱 Digital"
add_text(s, dims, 0.5, 5.75, 12.3, 0.9, size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

footer(s); slide_number(s, 5)

# ─── Slide 6: Workflow Diagram ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "วิธีที่ระบบทำงาน — ภาพรวม", 0.5, 0.2, 12, 0.6,
         size=28, bold=True, color=C_WHITE)

# Pipeline steps (top row)
steps = [
    ("Transaction\nData",           "12 เดือนย้อนหลัง",    C_LIGHT),
    ("Tag\nEngine",                 "Refresh ทุกวัน 01:00", C_ACCENT),
    ("Customer\nTag Profile",       "Score 0–1 ต่อ Tag",    C_ACCENT),
    ("Personalized\nTrigger",       "เลือก offer ที่ใช่",   RGBColor(0x4D,0xD0,0xFF)),
    ("Delivery\nChannel",           "Email / SMS / Push",   C_LIGHT),
]
box_w = 2.1
gap   = 0.22
sx    = 0.45
for i, (name, sub, col) in enumerate(steps):
    x = sx + i * (box_w + gap)
    add_rect(s, x, 1.05, box_w, 1.6, C_CARD)
    add_rect(s, x, 1.05, box_w, 0.08, col)
    add_text(s, name, x + 0.1, 1.2, box_w - 0.2, 0.7,
             size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, sub,  x + 0.1, 1.9, box_w - 0.2, 0.6,
             size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ax = x + box_w + 0.02
        add_text(s, "→", ax, 1.6, gap + 0.1, 0.4,
                 size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Trigger rule examples (bottom row)
rules = [
    ("frequent_diner ≥ 0.7\n+ high_spender ≥ 0.7",  "→  เสนอ Dining Cashback Offer",   C_ACCENT),
    ("car_owner = 1\n+ Tag เพิ่งเกิดใหม่",           "→  เสนอ EV Privilege Welcome",     RGBColor(0x4D,0xD0,0xFF)),
    ("frequent_diner score\nลดลง > 30% ใน 30 วัน",   "→  เสนอ Retention Offer",          C_YELLOW),
]
add_text(s, "ตัวอย่าง Personalized Trigger Rules", 0.5, 2.9, 12, 0.4,
         size=14, bold=True, color=C_WHITE)
for i, (cond, action, col) in enumerate(rules):
    x = 0.45 + i * 4.3
    add_rect(s, x, 3.4, 4.05, 1.8, C_CARD)
    add_rect(s, x, 3.4, 0.08, 1.8, col)
    add_text(s, cond,   x + 0.2, 3.5,  3.7, 0.75, size=11, color=C_LIGHT)
    add_text(s, action, x + 0.2, 4.25, 3.7, 0.75, size=12, bold=True, color=col)

# Callout
add_rect(s, 0.5, 5.45, 12.3, 0.55, RGBColor(0x00, 0x3A, 0x50))
add_text(s, "ไม่ว่าจะ automate หรือ manual — ทุก offer มาจาก Tag ของลูกค้า ไม่ใช่การเดา",
         0.7, 5.52, 12, 0.4, size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

footer(s); slide_number(s, 6)

# ─── Slide 7: Use Cases ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "นำไปใช้ได้ทันทีกับอะไร", 0.5, 0.2, 10, 0.6,
         size=28, bold=True, color=C_WHITE)

cases = [
    ("🎯  Targeted Campaign",
     "แทนที่จะส่ง Cashback ให้ทุกคน\n→ ส่งเฉพาะคนที่มี Tag \"frequent_diner\" + \"high_spender\"\n→ Conversion สูงขึ้น, Cost ต่ำลง"),
    ("💡  Product Offer Matching",
     "ลูกค้ามี Tag \"ev_driver\" + \"high_spender\"\n→ นำเสนอ EV Charging Privilege ก่อนใคร\n→ Offer ตรง ลูกค้าพึงพอใจสูง"),
    ("🔔  Early Churn Detection",
     "ลูกค้า Tag \"frequent_diner\" หายไปจากเดือนก่อน\n→ สัญญาณ Churn ก่อนที่จะเลิกใช้บัตร\n→ Retention ก่อนสาย"),
]
for i, (title, body) in enumerate(cases):
    y = 1.1 + i * 1.95
    add_rect(s, 0.5, y, 12.3, 1.75, C_CARD)
    add_rect(s, 0.5, y, 0.08, 1.75, C_ACCENT)
    add_text(s, title, 0.75, y + 0.08, 11.5, 0.5, size=17, bold=True, color=C_ACCENT)
    add_text(s, body, 0.75, y + 0.55, 11.5, 1.1, size=14, color=C_LIGHT)

footer(s); slide_number(s, 7)

# ─── Slide 8: End Game ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "End Game — เราจะไปถึงไหน", 0.5, 0.2, 12, 0.6, size=28, bold=True, color=C_WHITE)
add_text(s, "Tags คือ Foundation — ไม่ใช่ Destination", 0.5, 0.82, 12, 0.45,
         size=16, color=C_ACCENT, align=PP_ALIGN.LEFT)

layers = [
    ("1", "Tags",                  "รู้ว่าลูกค้าเป็นใคร",           "Foundation ทั้งหมด",              C_ACCENT),
    ("2", "Personalized Trigger",   "เสนอ offer ที่ใช่ ให้คนที่ใช่ ในเวลาที่เหมาะ", "ลด CAC, เพิ่ม Conversion", RGBColor(0x4D,0xD0,0xFF)),
    ("3", "Personalized Products", "Card ที่ benefit ตรงแต่ละคน",   "ลด Churn, เพิ่ม Spend per Card",  RGBColor(0x7B,0xC8,0xFF)),
    ("4", "Merchant Intelligence", "ขายความแม่นยำให้ Merchant",      "Revenue Stream ใหม่จาก Data",     RGBColor(0xA8,0xD8,0xEA)),
    ("5", "Lifestyle Platform",    "ส่วนหนึ่งของชีวิตลูกค้า",        "Brand Moat คู่แข่ง copy ไม่ได้",  RGBColor(0xB0,0xC4,0xDE)),
]
for i, (num, name, goal, result, col) in enumerate(layers):
    y = 1.45 + i * 1.02
    # number circle
    add_rect(s, 0.5, y, 0.7, 0.78, col)
    add_text(s, num, 0.5, y + 0.1, 0.7, 0.58, size=22, bold=True,
             color=C_BG, align=PP_ALIGN.CENTER)
    # name
    add_rect(s, 1.3, y, 3.2, 0.78, C_CARD)
    add_text(s, name, 1.45, y + 0.12, 3.0, 0.55, size=14, bold=True, color=col)
    # goal
    add_rect(s, 4.6, y, 4.3, 0.78, RGBColor(0x12,0x24,0x35))
    add_text(s, goal, 4.75, y + 0.12, 4.0, 0.55, size=13, color=C_WHITE)
    # result
    add_rect(s, 9.0, y, 4.2, 0.78, RGBColor(0x0A,0x1E,0x2E))
    add_text(s, "→  " + result, 9.15, y + 0.12, 3.9, 0.55, size=13, color=C_LIGHT)

# Hyundai reference
add_rect(s, 0.5, 6.55, 12.3, 0.45, RGBColor(0x00, 0x3A, 0x50))
add_text(s, "Hyundai Card: Music Library · Travel Library · Cooking Library — driven by Tag data ไม่ใช่ CSR",
         0.7, 6.6, 12, 0.35, size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

footer(s); slide_number(s, 8)

# ─── Slide 9: Roadmap ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "แผนการทำงาน", 0.5, 0.2, 10, 0.6, size=28, bold=True, color=C_WHITE)

phases = [
    ("Phase 1", "ทันที", C_ACCENT,
     ["เริ่ม 10 Quick Win Tags ที่ชัดเจน",
      "วัดผล Campaign A/B Test เทียบ approach เดิม",
      "ตัวอย่าง: frequent_diner, high_spender, car_owner"]),
    ("Phase 2", "หลัง Pilot", RGBColor(0x7B, 0xC8, 0xFF),
     ["ขยายเป็น 117 Tags ครบทุกมิติ",
      "Calibrate ความแม่นยำจาก Phase 1",
      "ใช้ Tags เป็น Feature ใน ML Models"]),
    ("Phase 3", "Full Intelligence", RGBColor(0xB0, 0xC4, 0xDE),
     ["เพิ่ม 34 Tags ด้วยข้อมูล Enrichment",
      "Cuisine Preference, Travel Destination",
      "Merchant Partnership Targeting"]),
]
for i, (phase, timing, col, bullets) in enumerate(phases):
    x = 0.5 + i * 4.3
    add_rect(s, x, 1.1, 4.0, 5.5, C_CARD)
    add_rect(s, x, 1.1, 4.0, 0.65, col)
    add_text(s, phase, x + 0.15, 1.15, 2.5, 0.5, size=18, bold=True, color=C_BG)
    add_text(s, timing, x + 0.15, 1.45, 3.7, 0.3, size=11, color=C_BG)
    for j, b in enumerate(bullets):
        by = 2.0 + j * 1.4
        add_rect(s, x + 0.2, by, 3.6, 1.2, RGBColor(0x0D, 0x3A, 0x52))
        add_text(s, "→  " + b, x + 0.35, by + 0.1, 3.3, 1.0, size=13, color=C_WHITE)

# Arrow connectors between phases
for xi in [4.55, 8.85]:
    add_text(s, "▶", xi - 0.05, 3.6, 0.5, 0.5, size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

footer(s); slide_number(s, 9)

# ─── Slide 10: Ask ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "สิ่งที่ต้องการจากการตัดสินใจวันนี้", 0.5, 0.2, 12, 0.6,
         size=28, bold=True, color=C_WHITE)

asks = [
    ("01", "Approve Phase 1", "เริ่ม implement 10 Quick Win Tags แรก"),
    ("02", "กำหนด KPI ร่วมกัน", "วัดความสำเร็จก่อน scale — Conversion, Cost per Acquisition"),
    ("03", "ระบุ Campaign ที่จะ Test", "ให้มีผล Business ที่จับต้องได้จริงใน Phase 1"),
]
for i, (num, title, desc) in enumerate(asks):
    y = 1.2 + i * 1.8
    add_rect(s, 0.5, y, 12.3, 1.55, C_CARD)
    add_rect(s, 0.5, y, 1.0, 1.55, C_ACCENT)
    add_text(s, num, 0.5, y + 0.3, 1.0, 0.9, size=30, bold=True,
             color=C_BG, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.7, y + 0.1, 10.8, 0.55, size=20, bold=True, color=C_WHITE)
    add_text(s, desc, 1.7, y + 0.7, 10.8, 0.7, size=15, color=C_LIGHT)

footer(s); slide_number(s, 10)

# ─── Slide 11: Bottom Line ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
circ2 = s.shapes.add_shape(9, Inches(-1.5), Inches(4), Inches(6), Inches(6))
circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x52)
circ2.line.fill.background()

add_text(s, "Bottom Line", 0.8, 0.8, 11, 0.7, size=22, color=C_ACCENT, bold=True)
lines = [
    ("เรามีข้อมูลลูกค้าอยู่แล้ว", 30, C_WHITE, True),
    ("คำถามคือ  เราอ่านมันได้ดีแค่ไหน", 22, C_LIGHT, False),
    ("", 10, C_WHITE, False),
    ("Customer Tagging คือการเปลี่ยนจาก", 18, C_LIGHT, False),
    ("\"รู้ว่าลูกค้าใช้จ่ายเท่าไหร่\"", 22, C_WHITE, True),
    ("เป็น", 18, C_LIGHT, False),
    ("\"รู้ว่าลูกค้าเป็นใคร\"", 32, C_ACCENT, True),
]
y = 1.6
for text, size, color, bold in lines:
    h = size / 30.0 * 0.65
    add_text(s, text, 1.5, y, 10, h + 0.1, size=size, color=color,
             bold=bold, align=PP_ALIGN.CENTER)
    y += h + 0.15

footer(s); slide_number(s, 11)

# ─── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_CEO_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
