"""
insert_agenda.py
แทรก Agenda slide ที่ตำแหน่ง 2 (index 1) ใน CustomerTagging_FullDeck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy

BASE = "/Users/adisornj/Desktop/Thena/lab"
TARGET = f"{BASE}/CustomerTagging_FullDeck.pptx"

# ── Color Palette ──────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x1A, 0x2E, 0x44)
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x00, 0xE5, 0x96)

# ── Helpers ────────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=C_BG):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, C_BG)

def footer(slide, n, note="Data Science Team | เมษายน 2025"):
    rect(slide, 0, 7.15, 13.33, 0.35, RGBColor(0x0A, 0x14, 0x20))
    txt(slide, note, 0.5, 7.18, 11, 0.28, size=9, color=C_LIGHT)
    txt(slide, str(n), 12.7, 7.18, 0.5, 0.28, size=9,
        color=C_LIGHT, align=PP_ALIGN.RIGHT)

# ── Build Agenda slide on a temp presentation ──────────────────────────────────
tmp = Presentation()
tmp.slide_width  = Inches(13.33)
tmp.slide_height = Inches(7.5)
BLANK = tmp.slide_layouts[6]

s = tmp.slides.add_slide(BLANK)
bg(s)

# Dark header bar
rect(s, 0, 0, 13.33, 0.75, RGBColor(0x0A, 0x14, 0x20))
rect(s, 0, 0.75, 13.33, 0.05, C_ACCENT)
txt(s, "Agenda", 0.5, 0.1, 11, 0.6, size=24, bold=True, color=C_WHITE)

# Section label
txt(s, "สิ่งที่เราจะคุยกันวันนี้", 0.5, 0.9, 12, 0.4, size=13, color=C_LIGHT)

# Agenda items
agenda = [
    ("01", "Big Picture",          "5 Layers of Intelligence",                     C_ACCENT),
    ("02", "Layer 1: Customer Tags","Tag Library · สร้าง/แปลง/Validate",           C_GREEN),
    ("03", "Layer 2: Personalized Trigger", "Automation Flow · Use Cases",          C_YELLOW),
    ("04", "Roadmap & Next Steps", "แผนการพัฒนาต่อ",                               C_LIGHT),
    ("05", "Appendix",             "Tag Library เต็ม — 323 Tags",                  RGBColor(0xA0, 0xA0, 0xA0)),
]

# Layout: 5 cards stacked, left half wide card + right number+detail
card_top   = 1.4
card_h     = 0.95
card_gap   = 0.08
card_l     = 0.5
card_w     = 12.33

accent_colors = [C_ACCENT, C_GREEN, C_YELLOW, C_LIGHT, RGBColor(0xA0, 0xA0, 0xA0)]

for i, (num, title, sub, accent) in enumerate(agenda):
    top = card_top + i * (card_h + card_gap)

    # Card background
    rect(s, card_l, top, card_w, card_h, C_CARD)

    # Left accent bar
    rect(s, card_l, top, 0.07, card_h, accent)

    # Number
    txt(s, num, card_l + 0.15, top + 0.18, 0.6, 0.55,
        size=26, bold=True, color=accent)

    # Title
    txt(s, title, card_l + 0.85, top + 0.1, 5.5, 0.45,
        size=18, bold=True, color=C_WHITE)

    # Subtitle
    txt(s, sub, card_l + 0.85, top + 0.55, 8.0, 0.35,
        size=11, color=C_LIGHT)

footer(s, 2)

# ── Merge: insert agenda slide at index 1 in target pptx ──────────────────────
prs = Presentation(TARGET)

# Copy agenda slide's XML elements into a new blank slide on the target prs
blank = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank)   # appended at end

# Clear default shapes from the new slide
sp_tree = new_slide.shapes._spTree
for el in list(sp_tree):
    sp_tree.remove(el)

# Copy all shapes from tmp agenda slide
for shape in s.shapes:
    sp_tree.append(copy.deepcopy(shape.element))

# Move slide from last position → index 1 (position 2)
# python-pptx tracks slides via slide_layout + rId — move via _sldIdLst
xml_slides = prs.slides._sldIdLst
# The newly added slide is the last entry in _sldIdLst
entries = list(xml_slides)
last_entry = entries[-1]
xml_slides.remove(last_entry)
# Insert at position 1 (after slide 0 = Title)
xml_slides.insert(1, last_entry)

prs.save(TARGET)
print(f"✅  Agenda slide inserted at position 2")
print(f"   Total slides: {len(prs.slides)}")
