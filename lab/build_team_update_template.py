from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Color palette — white theme, high contrast
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy-black
ACCENT      = RGBColor(0x16, 0x21, 0x3E)   # dark navy (header bar, shapes)
ACCENT2     = RGBColor(0x0F, 0x3D, 0x8C)   # vivid blue (rule lines, bullets)
LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF8)   # slide background tint
MID_GRAY    = RGBColor(0xCC, 0xD3, 0xE0)   # separator / sub-bullet indent
MUTED       = RGBColor(0x55, 0x65, 0x80)   # sub-bullet text

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size, bold=False,
                color=NEAR_BLACK, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


# ── BUILD ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # truly blank


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, WHITE)

# Left navy panel
add_rect(s1, Inches(0), Inches(0), Inches(4.2), H, ACCENT)

# Accent blue vertical stripe
add_rect(s1, Inches(4.2), Inches(0), Inches(0.08), H, ACCENT2)

# Top-left logo placeholder (white box)
add_rect(s1, Inches(0.4), Inches(0.35), Inches(1.6), Inches(0.55), WHITE)
add_textbox(s1, Inches(0.45), Inches(0.38), Inches(1.5), Inches(0.5),
            "LOGO", 10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Main title (right panel)
add_textbox(s1, Inches(4.7), Inches(1.8), Inches(8.2), Inches(1.2),
            "Team Project Update", 42, bold=True, color=ACCENT)

# Subtitle
add_textbox(s1, Inches(4.7), Inches(3.1), Inches(8.2), Inches(0.6),
            "Q[X] · [Month] [Year]", 22, bold=False, color=ACCENT2)

# Thin rule under subtitle
add_rect(s1, Inches(4.7), Inches(3.85), Inches(5.5), Inches(0.04), ACCENT2)

# Presenter / date line
add_textbox(s1, Inches(4.7), Inches(4.1), Inches(8.2), Inches(0.5),
            "Presented by: [Name]  ·  [Date]", 13, color=MUTED)

# Left panel — chapter label
add_textbox(s1, Inches(0.3), Inches(3.2), Inches(3.5), Inches(0.5),
            "PROJECT UPDATE", 11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_rect(s1, Inches(0.3), Inches(3.6), Inches(2.4), Inches(0.035), ACCENT2)

# Left panel — team name placeholder
add_textbox(s1, Inches(0.3), Inches(3.75), Inches(3.5), Inches(0.8),
            "[Team Name]", 18, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDES 2 & 3 — Content (topic + bullets)
# ═══════════════════════════════════════════════════════════════════════════
MAIN_LABELS = [
    "Main Point 1",
    "Main Point 2",
    "Main Point 3",
    "Main Point 4",
]
SUB_LABELS = [
    "Sub-bullet [detail]",
    "Sub-bullet [detail]",
    "Sub-bullet [detail]",
    "Sub-bullet [detail]",
]
SLIDE_TITLES = ["Agenda / Topics — Part 1", "Agenda / Topics — Part 2"]

for idx, title_text in enumerate(SLIDE_TITLES, start=2):
    sc = prs.slides.add_slide(blank_layout)
    set_bg(sc, LIGHT_GRAY)

    # Header bar
    add_rect(sc, Inches(0), Inches(0), W, Inches(1.05), ACCENT)

    # Slide title
    add_textbox(sc, Inches(0.45), Inches(0.18), Inches(10), Inches(0.7),
                title_text, 26, bold=True, color=WHITE)

    # Slide number (top right)
    add_textbox(sc, Inches(12.5), Inches(0.22), Inches(0.7), Inches(0.5),
                str(idx), 14, bold=True, color=ACCENT2, align=PP_ALIGN.RIGHT)

    # Thin accent rule below header
    add_rect(sc, Inches(0), Inches(1.05), W, Inches(0.045), ACCENT2)

    # Content area — 2 columns × 2 rows  (4 blocks)
    col_x   = [Inches(0.45), Inches(6.9)]
    row_y   = [Inches(1.35), Inches(4.25)]
    box_w   = Inches(6.1)
    box_h   = Inches(2.6)

    block = 0
    for row in range(2):
        for col in range(2):
            lx = col_x[col]
            ty = row_y[row]

            # White card background
            card = add_rect(sc, lx, ty, box_w, box_h, WHITE)
            card.line.color.rgb = MID_GRAY
            card.line.width = Pt(0.75)

            # Left accent stripe on card
            add_rect(sc, lx, ty, Inches(0.07), box_h, ACCENT2)

            # Main bullet label
            add_textbox(sc, lx + Inches(0.18), ty + Inches(0.12),
                        box_w - Inches(0.3), Inches(0.42),
                        f"●  {MAIN_LABELS[block]}", 15, bold=True, color=ACCENT)

            # Sub-bullets
            for sb in range(4):
                sub_y = ty + Inches(0.62) + sb * Inches(0.46)
                add_textbox(sc, lx + Inches(0.38), sub_y,
                            box_w - Inches(0.55), Inches(0.42),
                            f"–  {SUB_LABELS[sb]}", 11, bold=False, color=MUTED)

            block += 1

    # Footer rule + page hint
    add_rect(sc, Inches(0), Inches(7.2), W, Inches(0.04), MID_GRAY)
    add_textbox(sc, Inches(0.45), Inches(7.25), Inches(6), Inches(0.22),
                "[Confidential — Internal Use Only]", 9, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Thank You
# ═══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, WHITE)

# Full-width top navy strip
add_rect(s4, Inches(0), Inches(0), W, Inches(0.55), ACCENT)

# Bottom navy strip
add_rect(s4, Inches(0), Inches(6.9), W, Inches(0.6), ACCENT)

# Centered accent circle (decorative)
add_rect(s4, Inches(5.5), Inches(1.5), Inches(2.33), Inches(2.33), LIGHT_GRAY)

# Big "Thank You"
add_textbox(s4, Inches(0), Inches(1.8), W, Inches(1.5),
            "Thank You", 60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Horizontal rule
add_rect(s4, Inches(4.5), Inches(3.5), Inches(4.33), Inches(0.06), ACCENT2)

# Sub-line
add_textbox(s4, Inches(0), Inches(3.7), W, Inches(0.6),
            "Questions & Discussion", 20, bold=False,
            color=ACCENT2, align=PP_ALIGN.CENTER)

# Contact / closing note
add_textbox(s4, Inches(0), Inches(4.5), W, Inches(0.5),
            "[Presenter Name]  ·  [email@company.com]  ·  [Team]", 13,
            color=MUTED, align=PP_ALIGN.CENTER)

# Bottom strip text
add_textbox(s4, Inches(0.4), Inches(6.93), Inches(6), Inches(0.5),
            "[Confidential — Internal Use Only]", 10, color=WHITE)
add_textbox(s4, Inches(9), Inches(6.93), Inches(4), Inches(0.5),
            "[Date]", 10, color=WHITE, align=PP_ALIGN.RIGHT)


# ── SAVE ───────────────────────────────────────────────────────────────────
out = "/Users/adisornj/Desktop/Thena/lab/TeamUpdate_Template.pptx"
prs.save(out)
print(f"Saved: {out}")
