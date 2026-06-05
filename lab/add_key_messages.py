"""
add_key_messages.py
เพิ่ม Key Message bar ด้านล่างแต่ละ slide ที่กำหนด
- covering rect y=6.58–7.15 (ทับ content ที่ล้นลงมา)
- KM bar y=6.6, h=0.44, size=24
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TARGET = "/Users/adisornj/Desktop/Thena/lab/CustomerTagging_FullDeck.pptx"

C_BG    = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC2, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
KM_BG   = RGBColor(0x04, 0x1A, 0x2C)

def add_rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def add_txt(slide, text, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(24)
    r.font.bold = False
    r.font.color.rgb = C_ACCENT
    r.font.name = "Helvetica Neue"
    return tb

def km_bar(slide, text):
    add_rect(slide, 0, 6.58, 13.33, 0.57, C_BG)     # cover content ที่ล้นลงมา
    add_rect(slide, 0, 6.60, 13.33, 0.44, KM_BG)    # KM background
    add_rect(slide, 0, 6.60, 0.07,  0.44, C_ACCENT) # left accent bar
    add_txt(slide, text, 0.15, 6.62, 13.0, 0.42)    # KM text

# ── Key messages (0-based slide index) ─────────────────────────────────────────
KM = {
    3:  "ธุรกิจที่ชนะระยะยาวไม่แข่งด้วย cashback — แต่แข่งด้วยความเข้าใจลูกค้า",
    4:  "ทุก transaction ที่รูด คือข้อมูลที่บอกว่าลูกค้าเป็นใคร",
    5:  "tag ที่ถูกต้อง → trigger offer ที่ใช่ → ให้คนที่ใช่ → ในเวลาที่เหมาะ",
    6:  "tag ที่ถูกต้อง → trigger offer ที่ใช่ → ให้คนที่ใช่ → ในเวลาที่เหมาะ",
    7:  "automate ได้จริง ต้องมี 5 ส่วนพร้อมกัน — ขาดอันเดียวระบบหยุด",
    8:  "เปลี่ยนจาก 'ลูกค้าอยู่กลุ่มไหน' เป็น 'ลูกค้าเป็นคนแบบไหน'",
    9:  "323 tags ครอบทุกมิติพฤติกรรม — 254 เริ่มได้เลย ไม่ต้องลงทุนเพิ่ม",
    10: "323 tags ครอบทุกมิติพฤติกรรม — 254 เริ่มได้เลย ไม่ต้องลงทุนเพิ่ม",
    11: "ทุก tag มีนิยามชัด วัดจากพฤติกรรมจริง — ไม่ใช่การเดา",
    12: "ทุก tag มีนิยามชัด วัดจากพฤติกรรมจริง — ไม่ใช่การเดา",
    13: "เริ่มจาก 10 tags — พิสูจน์ผล — แล้วค่อย scale",
    14: "เรามีข้อมูลอยู่แล้ว — คำถามคือเราอ่านมันได้ดีแค่ไหน",
    15: "score บอกว่า active แค่ไหน — dimension บอกว่า active แบบไหน",
    16: "score บอกว่า active แค่ไหน — dimension บอกว่า active แบบไหน",
    17: "campaign manager query ได้ตรง — ไม่ต้องรู้ logic ข้างใน",
    18: "tag ที่ดีต้องผ่าน 3 ด่าน — sanity, quality, business signal",
    19: "ตั้ง rule ครั้งเดียว — ระบบทำงานแทนทุกวัน",
}

prs = Presentation(TARGET)

for idx, text in KM.items():
    km_bar(prs.slides[idx], text)
    print(f"  slide {idx+1:2d}: {text[:50]}...")

prs.save(TARGET)
print(f"\n✅  Key messages added — {len(KM)} slides updated")
print(f"   Total slides: {len(prs.slides)}")
